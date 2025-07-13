from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime


def export_client_dashboard(fund_data, client_name=None, output_path="client_dashboard.pptx"):
    """
    Creates a branded FidSync client dashboard PPTX with selected fund data and rationale.

    Parameters:
        fund_data (list of dict): Each dict should have:
            - 'fund_name': str
            - 'key_metrics': str or list of str
            - 'rationale': str
        client_name (str): Optional name to include on the cover
        output_path (str): File path to save the PPTX
    """
    prs = Presentation()
    
    # === Cover Slide ===
    cover_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(cover_layout)

    # Logo Banner
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = "FidSync Beta – Client Dashboard"
    run.font.name = 'Times New Roman'
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = 1

    # Client Name
    if client_name:
        name_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        name_tf = name_box.text_frame
        p = name_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"For: {client_name}"
        run.font.name = 'Times New Roman'
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = RGBColor(80, 80, 80)
        p.alignment = 1

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.6))
    date_tf = date_box.text_frame
    p = date_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Prepared on {datetime.now().strftime('%B %d, %Y')}"
    run.font.name = 'Times New Roman'
    run.font.size = Pt(14)
    p.alignment = 1

    # === Fund Slides ===
    for fund in fund_data:
        fund_name = fund.get("fund_name", "Unnamed Fund")
        key_metrics = fund.get("key_metrics", "")
        rationale = fund.get("rationale", "")

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Fund Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        p = title_tf.paragraphs[0]
        run = p.add_run()
        run.text = fund_name
        run.font.name = 'Times New Roman'
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 32, 96)

        # Metrics
        metrics_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(4), Inches(4))
        metrics_tf = metrics_box.text_frame
        metrics_tf.word_wrap = True
        p = metrics_tf.paragraphs[0]
        run = p.add_run()
        if isinstance(key_metrics, list):
            run.text = "\n".join(key_metrics)
        else:
            run.text = key_metrics
        run.font.name = 'Times New Roman'
        run.font.size = Pt(16)

        # Rationale
        rationale_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(4))
        rationale_tf = rationale_box.text_frame
        rationale_tf.word_wrap = True
        p = rationale_tf.paragraphs[0]
        run = p.add_run()
        run.text = rationale
        run.font.name = 'Times New Roman'
        run.font.size = Pt(16)
        run.font.italic = True

    # === Closing Slide ===
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank you for using FidSync Beta"
    run.font.name = 'Times New Roman'
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = 1

    # Save
    prs.save(output_path)
