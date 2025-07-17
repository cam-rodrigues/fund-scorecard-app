# utils/pptx_exporter.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from io import BytesIO

DARK_BLUE = RGBColor(0, 32, 96)
GRAY = RGBColor(80, 80, 80)
FONT = 'Times New Roman'


def create_fidsync_template_slide(fund_name, body_lines=None):
    """Returns a BytesIO PowerPoint stream with a single writeup slide."""
    if body_lines is None:
        body_lines = []

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Header: FidSync Beta
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = "FidSync Beta"
    run.font.name = FONT
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p.alignment = 1

    # Timestamp
    timestamp_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(10), Inches(0.5))
    timestamp_tf = timestamp_box.text_frame
    p = timestamp_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}"
    run.font.name = FONT
    run.font.size = Pt(14)
    p.alignment = 1

    # Fund Name + Writeup
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    title_paragraph = content_tf.add_paragraph()
    title_paragraph.text = fund_name
    title_paragraph.font.name = FONT
    title_paragraph.font.size = Pt(20)
    title_paragraph.font.bold = True
    title_paragraph.space_after = Pt(24)

    for line in body_lines:
        p = content_tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(18)
        p.space_after = Pt(24)

    # Export to memory
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def export_client_dashboard(fund_data, client_name=None):
    """Creates a multi-slide presentation for a client."""
    prs = Presentation()

    # Cover Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    run = p.add_run()
    run.text = "FidSync Beta – Client Dashboard"
    run.font.name = FONT
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p.alignment = 1

    if client_name:
        name_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        name_tf = name_box.text_frame
        p = name_tf.paragraphs[0]
        run = p.add_run()
        run.text = f"For: {client_name}"
        run.font.name = FONT
        run.font.size = Pt(20)
        run.font.italic = True
        run.font.color.rgb = GRAY
        p.alignment = 1

    date_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.6))
    date_tf = date_box.text_frame
    p = date_tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Prepared on {datetime.now().strftime('%B %d, %Y')}"
    run.font.name = FONT
    run.font.size = Pt(14)
    p.alignment = 1

    # Fund Slides
    for fund in fund_data:
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        fund_name = fund.get("fund_name", "Unnamed Fund")
        key_metrics = fund.get("key_metrics", [])
        rationale = fund.get("rationale", "")

        # Fund Name
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        p = title_tf.paragraphs[0]
        run = p.add_run()
        run.text = fund_name
        run.font.name = FONT
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = DARK_BLUE

        # Key Metrics
        metrics_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(4), Inches(4))
        metrics_tf = metrics_box.text_frame
        p = metrics_tf.paragraphs[0]
        run = p.add_run()
        run.text = "\n".join(key_metrics) if isinstance(key_metrics, list) else str(key_metrics)
        run.font.name = FONT
        run.font.size = Pt(16)

        # Rationale
        rationale_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(4))
        rationale_tf = rationale_box.text_frame
        p = rationale_tf.paragraphs[0]
        run = p.add_run()
        run.text = rationale
        run.font.name = FONT
        run.font.size = Pt(16)
        run.font.italic = True

    # Closing Slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank you for using FidSync Beta"
    run.font.name = FONT
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    p.alignment = 1

    # Return as in-memory file
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
