import re
import streamlit as st
import pdfplumber
from calendar import month_name
import pandas as pd
from rapidfuzz import fuzz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from io import BytesIO
import yfinance as yf

#───safer sentence splitting that protects common abbreviations like U.S.──────────────────────────────────────────────────────────────────
_ABBREV_PROTECT = {
    "U.S.": "__US__",
    "U.S": "__US__",
    "U.K.": "__UK__",
    "U.K": "__UK__",
    "e.g.": "__EG__",
    "e.g": "__EG__",
    "i.e.": "__IE__",
    "i.e": "__IE__",
    "etc.": "__ETC__",
    "etc": "__ETC__",
}

def safe_split_sentences(text):
    if not text:
        return []
    protected = text
    for k, v in _ABBREV_PROTECT.items():
        protected = protected.replace(k, v)
    sentences = re.split(r'(?<=[\.!?])\s+', protected.strip())
    def restore(s):
        for k, v in _ABBREV_PROTECT.items():
            s = s.replace(v, k)
        return s
    return [restore(s).strip() for s in sentences if s.strip()]


#───Performance Table──────────────────────────────────────────────────────────────────

def extract_performance_table(pdf, performance_page, fund_names, end_page=None):
    end = end_page if end_page is not None else len(pdf.pages) + 1
    lines = []
    for pnum in range(performance_page - 1, end - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")
    perf_data = []
    for name in fund_names:
        item = {"Fund Scorecard Name": name}
        idx = next((i for i, ln in enumerate(lines) if name in ln), None)
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower())) for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                continue
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        clean += [None] * (8 - len(clean))
        item["QTD"] = clean[0]
        item["1Yr"] = clean[2]
        item["3Yr"] = clean[3]
        item["5Yr"] = clean[4]
        item["10Yr"] = clean[5]
        item["Net Expense Ratio"] = clean[-2]
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]
        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None
        perf_data.append(item)
    return perf_data

#───Utility──────────────────────────────────────────────────────────────────

def extract_report_date(text):
    dates = re.findall(r'(\d{1,2})/(\d{1,2})/(20\d{2})', text or "")
    for month, day, year in dates:
        m, d = int(month), int(day)
        quarter_map = {(3,31): "1st", (6,30): "2nd", (9,30): "3rd", (12,31): "4th"}
        if (m, d) in quarter_map:
            return f"{quarter_map[(m, d)]} QTR, {year}"
        return f"As of {month_name[m]} {d}, {year}"
    return None

#───Page 1──────────────────────────────────────────────────────────────────

def process_page1(text):
    report_date = extract_report_date(text)
    st.session_state['report_date'] = report_date if report_date else None
    m = re.search(r"Total Options:\s*(\d+)", text or "")
    st.session_state['total_options'] = int(m.group(1)) if m else None
    m = re.search(r"Prepared For:\s*\n(.*)", text or "")
    st.session_state['prepared_for'] = m.group(1).strip() if m else None
    m = re.search(r"Prepared By:\s*(.*)", text or "")
    pb = m.group(1).strip() if m else ""
    if not pb or "mpi stylus" in pb.lower():
        pb = "Procyon Partners, LLC"
    st.session_state['prepared_by'] = pb

#───Info Card──────────────────────────────────────────────────────────────────

def show_report_summary():
    report_date    = st.session_state.get('report_date', 'N/A')
    total_options  = st.session_state.get('total_options', 'N/A')
    prepared_for   = st.session_state.get('prepared_for', 'N/A')
    prepared_by    = st.session_state.get('prepared_by', 'N/A')
    st.markdown(f"""
        <div style="
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.5rem;
            box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
            padding: 1.2rem 2.3rem 1.2rem 2.3rem;
            margin-bottom: 2rem;
            font-size: 1.08rem;
            border: 1.2px solid #b5d0eb;">
            <div style="color:#244369;">
                <b>Report Date:</b> {report_date} <br>
                <b>Total Options:</b> {total_options} <br>
                <b>Prepared For:</b> {prepared_for} <br>
                <b>Prepared By:</b> {prepared_by}
            </div>
        </div>
    """, unsafe_allow_html=True)

#───Table of Contents Extraction──────────────────────────────────────────────────────────

def process_toc(text):
    perf = re.search(r"Fund Performance[^\d]*(\d{1,3})", text or "")
    cy   = re.search(r"Fund Performance: Calendar Year\s+(\d{1,3})", text or "")
    r3yr = re.search(r"Risk Analysis: MPT Statistics \(3Yr\)\s+(\d{1,3})", text or "")
    r5yr = re.search(r"Risk Analysis: MPT Statistics \(5Yr\)\s+(\d{1,3})", text or "")

    sc            = re.search(r"Fund Scorecard\s+(\d{1,3})", text or "")
    sc_prop       = re.search(r"Fund Scorecard:\s*Proposed Funds\s+(\d{1,3})", text or "")

    fs            = re.search(r"Fund Factsheets\s+(\d{1,3})", text or "")
    fs_prop       = re.search(r"Fund Factsheets:\s*Proposed Funds\s+(\d{1,3})", text or "")

    st.session_state['performance_page'] = int(perf.group(1)) if perf else None
    st.session_state['calendar_year_page'] = int(cy.group(1)) if cy else None
    st.session_state['r3yr_page'] = int(r3yr.group(1)) if r3yr else None
    st.session_state['r5yr_page'] = int(r5yr.group(1)) if r5yr else None
    st.session_state['scorecard_page'] = int(sc.group(1)) if sc else None
    st.session_state['scorecard_proposed_page'] = int(sc_prop.group(1)) if sc_prop else None
    st.session_state['factsheets_page'] = int(fs.group(1)) if fs else None
    st.session_state['factsheets_proposed_page'] = int(fs_prop.group(1)) if fs_prop else None

#───IPS Investment Screening──────────────────────────────────────────────────────────────────

def infer_fund_type_guess(ticker):
    try:
        if not ticker:
            return ""
        info = yf.Ticker(ticker).info
        name = (info.get("longName") or info.get("shortName") or "").lower()
        summary = (info.get("longBusinessSummary") or "").lower()
        if "index" in name or "index" in summary:
            return "Passive"
        if "track" in summary and "index" in summary:
            return "Passive"
        if "actively managed" in summary or "actively-managed" in summary:
            return "Active"
        if "outperform" in summary or "manager selects" in summary:
            return "Active"
        return ""
    except Exception:
        return ""

def extract_scorecard_blocks(pdf, scorecard_page):
    metric_labels = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)"
    ]
    pages, fund_blocks, fund_name, metrics = [], [], None, []
    for p in pdf.pages[scorecard_page-1:]:
        pages.append(p.extract_text() or "")
    lines = "\n".join(pages).splitlines()
    for line in lines:
        if not any(metric in line for metric in metric_labels) and line.strip():
            if fund_name and metrics:
                fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
            fund_name = re.sub(
                r"Fund (Meets Watchlist Criteria|has been placed on watchlist for not meeting .* out of 14 criteria)",
                "",
                line.strip()
            ).strip()
            metrics = []
        for metric in metric_labels:
            if metric in line:
                m = re.match(r"^(.*?)\s+(Pass|Review|Fail)\s*(.*)", line.strip())
                if m:
                    metric_name, status, info = m.groups()
                    metrics.append({"Metric": metric_name, "Status": status, "Info": info.strip()})
    if fund_name and metrics:
        fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
    return fund_blocks

def extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page=None):
    def normalize_name(name):
        cleaned = re.sub(
            r"has been placed on watchlist for not meeting .*? criteria",
            "",
            name,
            flags=re.IGNORECASE
        )
        cleaned = re.sub(r"[^A-Za-z0-9 ]+", "", cleaned)
        return cleaned.strip().lower()

    end_page = factsheets_page - 1 if factsheets_page else len(pdf.pages)
    all_lines = []
    for p in pdf.pages[performance_page - 1:end_page]:
        txt = p.extract_text() or ""
        all_lines.extend([ln.strip() for ln in txt.splitlines() if ln.strip()])

    candidate_pairs = []
    ticker_rx = re.compile(r"\b([A-Z]{2,5})\b")
    for ln in all_lines:
        matches = ticker_rx.findall(ln)
        if not matches:
            continue
        for ticker in set(matches):
            parts = ln.rsplit(ticker, 1)
            if len(parts) >= 1:
                raw_name = parts[0].strip()
                if not raw_name:
                    continue
                norm_raw = normalize_name(raw_name)
                if not norm_raw:
                    continue
                candidate_pairs.append((norm_raw, ticker, raw_name))

    assigned = {}
    used_tickers = set()
    norm_expected_list = [(name, normalize_name(name)) for name in fund_names]

    for fund_name, norm_expected in norm_expected_list:
        best = (None, None, 0)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            score = fuzz.token_sort_ratio(norm_expected, norm_raw)
            if score > best[2]:
                best = (ticker, raw_name, score)
        if best[2] >= 70:
            assigned[fund_name] = best[0]
            used_tickers.add(best[0])
        else:
            assigned[fund_name] = ""

    for name in fund_names:
        if assigned.get(name):
            continue
        norm_expected = normalize_name(name)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            if norm_expected in norm_raw or norm_raw in norm_expected:
                assigned[name] = ticker
                used_tickers.add(ticker)
                break
        if not assigned.get(name):
            for ln in all_lines:
                if name.lower() in ln.lower():
                    m = re.search(r"\b([A-Z]{2,5})\b", ln)
                    if m:
                        assigned[name] = m.group(1)
                        break

    return {k: (v if v else "") for k, v in assigned.items()}

def scorecard_to_ips(fund_blocks, fund_types, tickers):
    metrics_order = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)",
    ]
    active_map  = [0,1,3,6,10,2,4,7,11,5,None]
    passive_map = [0,8,3,6,12,9,4,7,13,5,None]
    ips_labels = [f"IPS Investment Criteria {i+1}" for i in range(11)]
    ips_results, raw_results = [], []
    for fund in fund_blocks:
        fund_name = fund["Fund Name"]
        fund_type = fund_types.get(fund_name, "Passive" if "index" in fund_name.lower() else "Active")
        metrics = fund["Metrics"]
        scorecard_status = [next((m["Status"] for m in metrics if m["Metric"] == label), None) for label in metrics_order]
        idx_map = passive_map if fund_type == "Passive" else active_map
        ips_status = [scorecard_status[m_idx] if m_idx is not None else "Pass" for m_idx in idx_map]
        review_fail = sum(1 for status in ips_status if status in ["Review","Fail"])
        watch_status = "FW" if review_fail >= 6 else "IW" if review_fail >= 5 else "NW"
        def iconify(status): return "✔" if status == "Pass" else "✗" if status in ("Review", "Fail") else ""
        row = {
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: iconify(ips_status[i]) for i in range(11)},
            "IPS Watch Status": watch_status,
        }
        ips_results.append(row)
        raw_results.append({
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: ips_status[i] for i in range(11)},
            "IPS Watch Status": watch_status,
        })
    return pd.DataFrame(ips_results), pd.DataFrame(raw_results)

def watch_status_color(val):
    if val == "FW":
        return "background-color:#f8d7da; color:#c30000; font-weight:600;"
    if val == "IW":
        return "background-color:#fff3cd; color:#B87333; font-weight:600;"
    if val == "NW":
        return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
    return ""

def step3_5_6_scorecard_and_ips(pdf, scorecard_page, performance_page, factsheets_page, total_options):
    fund_blocks = extract_scorecard_blocks(pdf, scorecard_page)
    fund_names = [fund["Fund Name"] for fund in fund_blocks]
    if not fund_blocks:
        st.error("Could not extract fund scorecard blocks. Check the PDF and page number.")
        return

    tickers = extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page)

    inferred_guesses = []
    for name in fund_names:
        guess = ""
        if tickers.get(name):
            guess = infer_fund_type_guess(tickers.get(name, "")) or ""
        inferred_guesses.append("Passive" if guess.lower() == "passive" else ("Passive" if "index" in name.lower() else "Active"))

    df_types_base = pd.DataFrame({
        "Fund Name":       fund_names,
        "Ticker":          [tickers.get(n, "") for n in fund_names],
        "Inferred Type":   inferred_guesses,
        "Fund Type":       inferred_guesses,
    })

    if "show_edit_fund_type" not in st.session_state:
        st.session_state["show_edit_fund_type"] = False
    toggle_label = "Hide Fund Type Editor" if st.session_state["show_edit_fund_type"] else "Edit Fund Type"
    if st.button(toggle_label, key="toggle_edit_fund_type"):
        st.session_state["show_edit_fund_type"] = not st.session_state["show_edit_fund_type"]

    if st.session_state["show_edit_fund_type"]:
        st.markdown("### Fund Type Overrides")
        st.caption("Edit any fund type here; once you modify at least one, your edits take precedence over inferred types.")
        edited_types = st.data_editor(
            df_types_base,
            column_config={
                "Fund Type": st.column_config.SelectboxColumn("Fund Type", options=["Active", "Passive"]),
            },
            disabled=["Fund Name", "Ticker", "Inferred Type"],
            hide_index=True,
            key="data_editor_fundtype_ips",
            use_container_width=True,
        )
        manual_override = any(
            row["Fund Type"] != row["Inferred Type"]
            for _, row in edited_types.iterrows()
        )
        if manual_override:
            fund_types = {row["Fund Name"]: row["Fund Type"] for _, row in edited_types.iterrows()}
        else:
            fund_types = {row["Fund Name"]: row["Inferred Type"] for _, row in edited_types.iterrows()}
    else:
        fund_types = {name: inferred_guesses[i] for i, name in enumerate(fund_names)}

    df_icon, df_raw = scorecard_to_ips(fund_blocks, fund_types, tickers)

    st.subheader("IPS Screening Results")
    st.markdown(
        '<div style="display:flex; gap:1rem; margin-bottom:0.5rem;">'
        '<div style="padding:4px 10px; background:#d6f5df; border-radius:4px; font-weight:600;">NW: No Watch</div>'
        '<div style="padding:4px 10px; background:#fff3cd; border-radius:4px; font-weight:600;">IW: Informal Watch</div>'
        '<div style="padding:4px 10px; background:#f8d7da; border-radius:4px; font-weight:600;">FW: Formal Watch</div>'
        '</div>', unsafe_allow_html=True
    )

    if df_icon.empty:
        st.info("No IPS screening data available.")
    else:
        display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
        display_df = df_icon.rename(columns=display_columns)

        def iconify(s):
            if s == "Pass":
                return "✔"
            if s in ("Review", "Fail"):
                return "✗"
            return ""

        compact_df = display_df.copy()
        for i in range(1, 12):
            orig = f"IPS Investment Criteria {i}"
            if orig in compact_df.columns:
                compact_df[str(i)] = compact_df[orig].apply(iconify)
                compact_df.drop(columns=[orig], inplace=True)
        cols_order = ["Fund Name", "Ticker", "Fund Type"] + [str(i) for i in range(1, 12)] + ["IPS Watch Status"]
        compact_df = compact_df[[c for c in cols_order if c in compact_df.columns]]

        def watch_style(val):
            if val == "NW":
                return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
            if val == "IW":
                return "background-color:#fff3cd; color:#B87333; font-weight:600;"
            if val == "FW":
                return "background-color:#f8d7da; color:#c30000; font-weight:600;"
            return ""

        styled = compact_df.style.applymap(watch_style, subset=["IPS Watch Status"])
        st.dataframe(styled, use_container_width=True, hide_index=True)

    st.session_state["fund_blocks"] = fund_blocks
    st.session_state["fund_types"] = fund_types
    st.session_state["fund_tickers"] = tickers
    st.session_state["ips_icon_table"] = df_icon
    st.session_state["ips_raw_table"] = df_raw

    perf_data = extract_performance_table(pdf, performance_page, fund_names, factsheets_page)
    for itm in perf_data:
        itm["Ticker"] = tickers.get(itm["Fund Scorecard Name"], "")
    st.session_state["fund_performance_data"] = perf_data
    st.session_state["tickers"] = tickers


#───Proposed Funds Extraction──────────────────────────────────────────────────────────────────

def extract_proposed_scorecard_blocks(pdf, *, fuzzy_threshold=78, min_token_overlap=2, max_pages_per_section=4):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    # --- helpers -------------------------------------------------------------
    def norm_text(s: str) -> str:
        return " ".join((s or "").strip().upper().split())

    def page_text(i: int) -> str:
        return pdf.pages[i - 1].extract_text() or ""

    def tokens(s: str):
        # alphanumeric tokens length >=3
        return {t for t in re.findall(r"[A-Za-z0-9]+", (s or "")) if len(t) >= 3}

    RE_PROPOSED_HDR = re.compile(r"FUND\s*SCORECARD\s*:\s*PROPOSED\s*FUNDS", re.I)

    # Use TOC anchors as hard bounds if available
    anchor_keys = (
        "factsheets_proposed_page", "factsheets_page", "calendar_year_page",
        "r3yr_page", "r5yr_page", "performance_page", "scorecard_page"
    )
    anchors = sorted(
        p for k in anchor_keys
        for p in [st.session_state.get(k)]
        if isinstance(p, int)
    )

    # --- 1) find ALL Proposed header pages ----------------------------------
    starts = []
    for p in range(1, len(pdf.pages) + 1):
        t = page_text(p)
        if RE_PROPOSED_HDR.search(t) or "FUND SCORECARD: PROPOSED FUNDS" in norm_text(t):
            starts.append(p)

    # include TOC start if OCR missed header
    toc_start = st.session_state.get("scorecard_proposed_page")
    if isinstance(toc_start, int) and toc_start not in starts:
        starts.append(toc_start)
    starts = sorted(set(starts))

    if not starts:
        st.session_state["proposed_funds_confirmed_df"] = pd.DataFrame()
        return pd.DataFrame()

    # --- 2) collect ONLY lines from pages that STILL show the header ----------
    sections = []
    for s in starts:
        next_anchor = next((a for a in anchors if a > s), len(pdf.pages) + 1)
        next_header = next((h for h in starts if h > s), len(pdf.pages) + 1)
        hard_end = min(next_anchor, next_header, len(pdf.pages) + 1)

        lines, pages_used = [], 0
        for p in range(s, hard_end):
            txt = page_text(p)
            # stop if header disappears (keeps scope tight)
            if not (RE_PROPOSED_HDR.search(txt) or "FUND SCORECARD: PROPOSED FUNDS" in norm_text(txt)):
                break
            # keep non-empty lines
            lines.extend([ln.strip() for ln in (txt.splitlines() or []) if ln.strip()])
            pages_used += 1
            if max_pages_per_section and pages_used >= max_pages_per_section:
                break

        if lines:
            sections.append({"start": s, "lines": lines})

    if not sections:
        st.session_state["proposed_funds_confirmed_df"] = pd.DataFrame()
        return pd.DataFrame()

    # --- 3) candidate funds (with tickers to backfill later) -----------------
    perf_data = st.session_state.get("fund_performance_data", []) or []
    if not perf_data:
        st.session_state["proposed_funds_confirmed_df"] = pd.DataFrame()
        return pd.DataFrame()

    # map name -> ticker for backfill
    name_to_ticker = { (it.get("Fund Scorecard Name") or "").strip(): (it.get("Ticker") or "").strip().upper()
                       for it in perf_data }

    # --- 4) name-only matching within proposed sections ----------------------
    results = []
    for it in perf_data:
        fund_name = (it.get("Fund Scorecard Name") or "").strip()
        if not fund_name:
            continue

        name_tok = tokens(fund_name)
        best_score, best_page, best_line = 0, None, ""

        for sec in sections:
            for ln in sec["lines"]:
                # quick token overlap guard to reduce false hits
                if len(name_tok.intersection(tokens(ln))) < min_token_overlap:
                    continue
                score = fuzz.token_set_ratio(fund_name.lower(), ln.lower())
                if score > best_score:
                    best_score, best_page, best_line = score, sec["start"], ln

        if best_score >= fuzzy_threshold:
            results.append({
                "Fund Scorecard Name": fund_name,
                "Ticker": name_to_ticker.get(fund_name, ""),  # backfill ticker here
                "Proposed Section Start Page": best_page,
                "Match Score": best_score,
                "Matched Line": best_line
            })

    df = pd.DataFrame(results).drop_duplicates(subset=["Fund Scorecard Name"])
    st.session_state["proposed_funds_confirmed_df"] = df
    return df

#───Side-by-side Info Card Helpers──────────────────────────────────────────────────────

def _shared_cards_css():
    return """
    <style>
      /* Card shell */
      .fid-card {
        background: linear-gradient(120deg, #e6f0fb 84%, #cfe4f8 100%);
        color: #23395d;
        border-radius: 1.25rem;
        box-shadow: 0 2px 14px rgba(44,85,130,0.08), 0 1px 4px rgba(36,67,105,0.07);
        border: 1.2px solid #b5d0eb;
        padding: 1.1rem 1.6rem;
        margin: 0 0 1.1rem 0;
      }
      .fid-card h4 {
        margin: 0 0 .45rem 0;
        font-weight: 700;
        font-size: 1.08rem;
        color: #223d63;
        letter-spacing: -.2px;
      }
      .fid-card .sub {
        font-size: .98rem;
        color: #2b4770;
        margin: 0 0 .6rem 0;
      }

      /* Tables */
      table.fid-table {
        width: 100%;
        border-collapse: collapse;
        font-family: 'Inter', 'Segoe UI', Arial, sans-serif;
        font-size: .98rem;
      }
      table.fid-table th, table.fid-table td {
        border: none;
        padding: .42rem .9rem;
        text-align: left;
      }
      table.fid-table th {
        background: #244369;
        color: #fff;
        font-weight: 700;
        letter-spacing: .01em;
      }
      table.fid-table tr:nth-child(even) { background: #e9f2fc; }
      table.fid-table tr:nth-child(odd)  { background: #f8fafc; }

      /* Column alignment helpers */
      .center { text-align: center; }
      .ticker  { font-family: ui-monospace, SFMono-Regular, Menlo, Consolas, "Liberation Mono", monospace; letter-spacing: .2px; }

      /* Watch badges (in-table values or summary blocks) */
      .badge-nw { background:#d6f5df; color:#217a3e; border-radius:.5rem; padding:.15rem .5rem; font-weight:600; }
      .badge-iw { background:#fff3cd; color:#B87333; border-radius:.5rem; padding:.15rem .5rem; font-weight:600; }
      .badge-fw { background:#f8d7da; color:#c30000; border-radius:.5rem; padding:.15rem .5rem; font-weight:600; }

      /* Narrow tweaks for Proposed table: center ticker column */
      table.proposed-fund-table td:nth-child(2),
      table.proposed-fund-table th:nth-child(2) { text-align:center; }
      table.proposed-fund-table td:nth-child(2) { font-family: ui-monospace, monospace; letter-spacing:.2px; }
    </style>
    """

def get_ips_fail_card_html():
    df = st.session_state.get("ips_icon_table")
    if df is None or df.empty:
        return "", ""
    fail_df = df[df["IPS Watch Status"].isin(["FW", "IW"])][["Fund Name", "IPS Watch Status"]]
    if fail_df.empty:
        return "", ""

    # Map status to a colored badge label
    def badge(s):
        if s == "FW": return '<span class="badge-fw">FW</span>'
        if s == "IW": return '<span class="badge-iw">IW</span>'
        return s

    display = fail_df.rename(columns={"Fund Name": "Fund", "IPS Watch Status": "Watch Status"}).copy()
    display["Watch Status"] = display["Watch Status"].map(badge)

    table_html = display.to_html(index=False, border=0, justify="center",
                                 classes="fid-table ips-fail-table",
                                 escape=False)

    card_html = f"""
      <div class="fid-card">
        <h4>Funds on Watch</h4>
        <div class="sub">The following funds failed five or more IPS criteria and are currently on watch.</div>
        {table_html}
      </div>
    """
    return card_html, _shared_cards_css()

def get_proposed_fund_card_html(*, only_with_tickers=True, min_score=74):
    df = st.session_state.get("proposed_funds_confirmed_df")

    if df is None or df.empty:
        card_html = """
          <div class="fid-card">
            <h4>Confirmed Proposed Funds</h4>
            <div class="sub">No confirmed proposed funds were found on the Proposed Funds scorecard pages.</div>
          </div>
        """
        return card_html, _shared_cards_css()

    df = df.copy()

    # 1) Keep only rows we know came from Proposed sections
    if "Proposed Section Start Page" in df.columns:
        df = df[df["Proposed Section Start Page"].notna()]

    # 2) Enforce a minimum fuzzy match score, if available
    if "Match Score" in df.columns and min_score is not None:
        df = df[df["Match Score"] >= min_score]

    # 3) Require a ticker (optional but helps avoid noisy name-only matches)
    if only_with_tickers and "Ticker" in df.columns:
        df = df[df["Ticker"].astype(str).str.len() > 0]

    # 4) Final columns, sort, and de-dupe
    cols = [c for c in ["Fund Scorecard Name", "Ticker"] if c in df.columns]
    if not cols:
        # Fallback if columns got renamed upstream
        cols = df.columns[:2].tolist()
    display_df = (
        df[cols]
        .rename(columns={"Fund Scorecard Name": "Fund"})
        .drop_duplicates()
        .sort_values(by=["Fund", "Ticker"], kind="stable")
    )

    if display_df.empty:
        card_html = """
          <div class="fid-card">
            <h4>Confirmed Proposed Funds</h4>
            <div class="sub">No confirmed proposed funds were found on the Proposed Funds scorecard pages.</div>
          </div>
        """
        return card_html, _shared_cards_css()

    table_html = display_df.to_html(
        index=False, border=0, justify="center",
        classes="fid-table proposed-fund-table",
        escape=True
    )

    card_html = f"""
      <div class="fid-card">
        <h4>Confirmed Proposed Funds</h4>
        <div class="sub">The following funds were identified on the Proposed Funds scorecard pages.</div>
        {table_html}
      </div>
    """
    return card_html, _shared_cards_css()


def get_watch_summary_card_html():
    df = st.session_state.get("ips_icon_table")
    if df is None or df.empty:
        return "", ""
    counts = df["IPS Watch Status"].value_counts().to_dict()
    summary = {
        "No Watch": counts.get("NW", 0),
        "Informal Watch": counts.get("IW", 0),
        "Formal Watch": counts.get("FW", 0),
    }

    card_html = f"""
      <div class="fid-card" style="max-width:540px;">
        <h4>Watch Summary</h4>
        <div style="display:flex; gap:1rem; align-items:stretch; margin-top:.2rem;">
          <div class="badge-nw" style="flex:1; text-align:center;">
            No Watch<br><span style="font-size:1.35rem; font-weight:800;">{summary["No Watch"]}</span>
          </div>
          <div class="badge-iw" style="flex:1; text-align:center;">
            Informal Watch<br><span style="font-size:1.35rem; font-weight:800;">{summary["Informal Watch"]}</span>
          </div>
          <div class="badge-fw" style="flex:1; text-align:center;">
            Formal Watch<br><span style="font-size:1.35rem; font-weight:800;">{summary["Formal Watch"]}</span>
          </div>
        </div>
      </div>
    """
    return card_html, _shared_cards_css()


#───Step 6:Factsheets Pages──────────────────────────────────────────────────────────────────

def step6_process_factsheets(pdf, fund_names, suppress_output=True):
    # If you ever want UI, set suppress_output=False when calling

    factsheet_start = st.session_state.get("factsheets_page")
    total_declared = st.session_state.get("total_options")
    performance_data = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in st.session_state.get("tickers", {}).items()
    ]

    if not factsheet_start:
        if not suppress_output:
            st.error("❌ 'Fund Factsheets' page number not found in TOC.")
        return

    matched_factsheets = []
    for i in range(factsheet_start - 1, len(pdf.pages)):
        page = pdf.pages[i]
        words = page.extract_words(use_text_flow=True)
        header_words = [w['text'] for w in words if w['top'] < 100]
        first_line = " ".join(header_words).strip()

        if not first_line or "Benchmark:" not in first_line or "Expense Ratio:" not in first_line:
            continue

        ticker_match = re.search(r"\b([A-Z]{5})\b", first_line)
        ticker = ticker_match.group(1) if ticker_match else ""
        fund_name_raw = first_line.split(ticker)[0].strip() if ticker else first_line

        best_score = 0
        matched_name = matched_ticker = ""
        for item in performance_data:
            ref = f"{item['Fund Scorecard Name']} {item['Ticker']}".strip()
            score = fuzz.token_sort_ratio(f"{fund_name_raw} {ticker}".lower(), ref.lower())
            if score > best_score:
                best_score, matched_name, matched_ticker = score, item['Fund Scorecard Name'], item['Ticker']

        def extract_field(label, text, stop=None):
            try:
                start = text.index(label) + len(label)
                rest = text[start:]
                if stop and stop in rest:
                    return rest[:rest.index(stop)].strip()
                return rest.split()[0]
            except Exception:
                return ""

        benchmark = extract_field("Benchmark:", first_line, "Category:")
        category  = extract_field("Category:", first_line, "Net Assets:")
        net_assets= extract_field("Net Assets:", first_line, "Manager Name:")
        manager   = extract_field("Manager Name:", first_line, "Avg. Market Cap:")
        avg_cap   = extract_field("Avg. Market Cap:", first_line, "Expense Ratio:")
        expense   = extract_field("Expense Ratio:", first_line)

        matched_factsheets.append({
            "Page #": i + 1,
            "Parsed Fund Name": fund_name_raw,
            "Parsed Ticker": ticker,
            "Matched Fund Name": matched_name,
            "Matched Ticker": matched_ticker,
            "Benchmark": benchmark,
            "Category": category,
            "Net Assets": net_assets,
            "Manager Name": manager,
            "Avg. Market Cap": avg_cap,
            "Expense Ratio": expense,
            "Match Score": best_score,
            "Matched": "✅" if best_score > 20 else "❌"
        })

    df_facts = pd.DataFrame(matched_factsheets)
    st.session_state['fund_factsheets_data'] = matched_factsheets

    # Hide UI output unless suppress_output is False
    if not suppress_output:
        display_df = df_facts[[
            "Matched Fund Name", "Matched Ticker", "Benchmark", "Category",
            "Net Assets", "Manager Name", "Avg. Market Cap", "Expense Ratio", "Matched"
        ]].rename(columns={"Matched Fund Name": "Fund Name", "Matched Ticker": "Ticker"})

        st.dataframe(display_df, use_container_width=True)

        matched_count = sum(1 for r in matched_factsheets if r["Matched"] == "✅")
        st.write(f"Matched {matched_count} of {len(matched_factsheets)} factsheet pages.")
        if matched_count == total_declared:
            st.success(f"All {matched_count} funds matched the declared Total Options from Page 1.")
        else:
            st.error(f"Mismatch: Page 1 declared {total_declared}, but only matched {matched_count}.")

#───Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense Ratio & Bench QTD──────────────────────────────────────────────────────────────────

def step7_extract_returns(pdf):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    # 1) Where to scan
    perf_page = st.session_state.get("performance_page")
    end_page  = st.session_state.get("calendar_year_page") or (len(pdf.pages) + 1)
    perf_data = st.session_state.get("fund_performance_data", [])
    if perf_page is None or not perf_data:
        st.error("❌ Run Step 5 first to populate performance data.")
        return

    # 2) Prep output slots
    fields = [
        "QTD", "1Yr", "3Yr", "5Yr", "10Yr", "Net Expense Ratio",
        "Bench QTD", "Bench 1Yr", "Bench 3Yr", "Bench 5Yr", "Bench 10Yr"
    ]
    for itm in perf_data:
        for f in fields:
            itm.setdefault(f, None)

    # 3) Gather every nonblank line in the Performance section
    lines = []
    for pnum in range(perf_page - 1, end_page - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 4) Regex to pull decimal tokens (with optional % and parentheses)
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    matched = 0
    for item in perf_data:
        name = item["Fund Scorecard Name"]
        tk   = item["Ticker"].upper().strip()

        # a) Exact-ticker match
        idx = next(
            (i for i, ln in enumerate(lines)
             if re.search(rf"\b{re.escape(tk)}\b", ln)),
            None
        )
        # b) Fuzzy-name fallback
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                st.warning(f"⚠️ {name} ({tk}): no match found.")
                continue

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        if len(clean) < 8:
            clean += [None] * (8 - len(clean))

        # d) Map fund returns & net expense
        item["QTD"]               = clean[0]
        item["1Yr"]               = clean[2]
        item["3Yr"]               = clean[3]
        item["5Yr"]               = clean[4]
        item["10Yr"]              = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 1Yr, 3Yr, 5Yr, and 10Yr from the very next line(s)
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 5 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]

        item["Bench QTD"]  = bench_clean[0] if len(bench_clean) > 0 else None
        item["Bench 1Yr"] = bench_clean[1] if len(bench_clean) > 1 else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None
        item["Bench 10Yr"] = bench_clean[5] if len(bench_clean) > 5 else None

        matched += 1

    # 5) Save & display
    st.session_state["fund_performance_data"] = perf_data
    df = pd.DataFrame(perf_data)
    # Hide the per-fund warnings and overall success message, but still alert if something is off
    
    expected_count = len(perf_data)
    if matched < expected_count:
        st.error(f"❌ Only matched {matched} of {expected_count} funds with return data. Check your PDF or extraction logic.")
    # (Do NOT display per-fund warnings or success)

    st.dataframe(
        df[["Fund Scorecard Name", "Ticker"] + fields],
        use_container_width=True
    )


#───Step 8 Calendar Year Returns (funds + benchmarks──────────────────────────────────────────────────────────────────

def step8_calendar_returns(pdf):
    import re, streamlit as st, pandas as pd
    from rapidfuzz import fuzz

    # 1) Section bounds
    cy_page = st.session_state.get("calendar_year_page")
    if cy_page is None:
        st.error("❌ 'Fund Performance: Calendar Year' not found in TOC.")
        return

    next_page = st.session_state.get("r3yr_page")  # may be None
    # Fallback to end-of-PDF when missing, and ensure we scan at least one page
    end_page = next_page if isinstance(next_page, int) else (len(pdf.pages) + 1)
    end_page = max(end_page, cy_page + 1)
    end_page = min(end_page, len(pdf.pages) + 1)

    # 2) Pull lines in section
    all_lines = []
    for p in pdf.pages[cy_page - 1 : end_page - 1]:
        all_lines.extend((p.extract_text() or "").splitlines())

    # 3) Identify header & years
    header = next((ln for ln in all_lines if "Ticker" in ln and re.search(r"\b20\d{2}\b", ln)), None)
    if not header:
        st.error("❌ Couldn’t find header row with 'Ticker' + year.")
        return
    years = re.findall(r"\b20\d{2}\b", header)
    num_rx = re.compile(r"-?\d+\.\d+%?")

    # — A) Funds —
    fund_map = st.session_state.get("tickers", {}) or {}
    fund_records = []
    for name, tk in fund_map.items():
        ticker = (tk or "").upper()
        # robust ticker search (word-boundary)
        idx = next(
            (i for i, ln in enumerate(all_lines)
             if re.search(rf"\b{re.escape(ticker)}\b", ln)),
            None
        )
        # numbers typically appear on the line above the ticker row
        raw = num_rx.findall(all_lines[idx - 1]) if idx not in (None, 0) else []
        vals = raw[:len(years)] + [None] * (len(years) - len(raw))

        if idx is None:
            st.warning(f"⚠️ Calendar-year table: no ticker row found for {name} ({ticker}).")

        rec = {"Name": name, "Ticker": ticker}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        fund_records.append(rec)

    df_fund = pd.DataFrame(fund_records)
    if not df_fund.empty:
        st.markdown("**Fund Calendar-Year Returns**")
        st.dataframe(df_fund[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["step8_returns"] = fund_records

    # — B) Benchmarks per fund —
    facts = st.session_state.get("fund_factsheets_data", []) or []
    bench_records = []
    for f in facts:
        bench_name = (f.get("Benchmark") or "").strip()
        fund_tkr   = (f.get("Matched Ticker") or "").upper()
        if not bench_name:
            continue

        # exact contains
        idx = next((i for i, ln in enumerate(all_lines) if bench_name in ln), None)
        if idx is None:
            # fuzzy fallback (loose threshold)
            best = max(
                ((i, fuzz.token_set_ratio(bench_name.lower(), ln.lower()))
                 for i, ln in enumerate(all_lines)),
                key=lambda x: x[1],
                default=(None, 0)
            )
            idx = best[0] if best[1] >= 70 else None

        if idx is None:
            st.warning(f"⚠️ Calendar-year benchmark row not found for: {bench_name} ({fund_tkr}).")
            continue

        raw = num_rx.findall(all_lines[idx])
        vals = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec = {"Name": bench_name, "Ticker": fund_tkr}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        bench_records.append(rec)

    df_bench = pd.DataFrame(bench_records)
    if not df_bench.empty:
        st.markdown("**Benchmark Calendar-Year Returns**")
        st.dataframe(df_bench[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["benchmark_calendar_year_returns"] = bench_records
    else:
        st.warning("No benchmark returns extracted.")


#───Step 9: 3‑Yr Risk Analysis – Match & Extract MPT Stats (hidden matching)──────────────────────────────────────────────────────────────────

def step9_risk_analysis_3yr(pdf):
    import re, streamlit as st, pandas as pd
    from rapidfuzz import fuzz

    # 1) Get your fund→ticker map
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (3Yr)” page
    start_page = st.session_state.get("r3yr_page")
    if not start_page:
        st.error("❌ ‘Risk Analysis: MPT Statistics (3Yr)’ page not found; run Step 2 first.")
        return

    # 3) Scan forward until you’ve seen each ticker (no display)
    locs = {}
    for pnum in range(start_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for fname, tk in fund_map.items():
                if fname in locs: 
                    continue
                if tk.upper() in tokens:
                    locs[fname] = {"page": pnum, "line": li}
        if len(locs) == len(fund_map):
            break

    # 4) Extract the first four numeric MPT stats from that same line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, info in locs.items():
        page = pdf.pages[info["page"]-1]
        lines = (page.extract_text() or "").splitlines()
        line = lines[info["line"]]
        nums = num_rx.findall(line)
        nums += [None] * (4 - len(nums))
        alpha, beta, up, down = nums[:4]
        results.append({
            "Fund Name":               name,
            "Ticker":                  fund_map[name].upper(),
            "3 Year Alpha":            alpha,
            "3 Year Beta":             beta,
            "3 Year Upside Capture":   up,
            "3 Year Downside Capture": down
        })

    # 5) Display final table only
    st.session_state["step9_mpt_stats"] = results

#───Step 10: Risk Analysis (5Yr) – Match & Extract MPT Statistics──────────────────────────────────────────────────────────────────

def step10_risk_analysis_5yr(pdf):
    import re, streamlit as st, pandas as pd

    # 1) Your fund→ticker map from Step 5
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (5Yr)” section
    section_page = next(
        (i for i, pg in enumerate(pdf.pages, start=1)
         if "Risk Analysis: MPT Statistics (5Yr)" in (pg.extract_text() or "")),
        None
    )
    if section_page is None:
        st.error("❌ Could not find ‘Risk Analysis: MPT Statistics (5Yr)’ section.")
        return

    # 3) Under‑the‑hood: scan pages until each ticker is located
    locs = {}
    total = len(fund_map)
    for pnum in range(section_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for name, tk in fund_map.items():
                if name in locs:
                    continue
                if tk.upper() in tokens:
                    locs[name] = {"page": pnum, "line": li}
        if len(locs) == total:
            break

    # 4) Wrap‑aware extraction of the first four floats after each ticker line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, tk in fund_map.items():
        info = locs.get(name)
        vals = [None] * 4
        if info:
            page = pdf.pages[info["page"] - 1]
            text_lines = (page.extract_text() or "").splitlines()
            idx = info["line"]
            nums = []
            # look on the line of the ticker and up to the next 2 lines
            for j in range(idx, min(idx + 3, len(text_lines))):
                nums += num_rx.findall(text_lines[j])
                if len(nums) >= 4:
                    break
            nums += [None] * (4 - len(nums))
            vals = nums[:4]
        else:
            st.warning(f"⚠️ {name} ({tk.upper()}): not found after page {section_page}.")

        alpha5, beta5, up5, down5 = vals
        results.append({
            "Fund Name":               name,
            "Ticker":                  tk.upper(),
            "5 Year Alpha":            alpha5,
            "5 Year Beta":             beta5,
            "5 Year Upside Capture":   up5,
            "5 Year Downside Capture": down5,
        })

    # 5) Save & display only the consolidated table
    st.session_state["step10_mpt_stats"] = results

#───Step 11: Combined MPT Statistics Summary──────────────────────────────────────────────────────────────────

def step11_create_summary(pdf=None):
    import pandas as pd
    import streamlit as st

    # 1) Load your 3‑Yr and 5‑Yr stats from session state
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    if not mpt3 or not mpt5:
        st.error("❌ Missing MPT stats. Run Steps 9 & 10 first.")
        return

    # 2) Build DataFrames
    df3 = pd.DataFrame(mpt3)  # contains "3 Year Alpha", "3 Year Beta", etc.
    df5 = pd.DataFrame(mpt5)  # contains "5 Year Alpha", "5 Year Beta", etc.

    # 3) Merge on Fund Name & Ticker
    df = pd.merge(
        df3,
        df5,
        on=["Fund Name", "Ticker"],
        how="outer",
        suffixes=("_3yr", "_5yr")
    )

    # 4) Build the Investment Manager column
    df.insert(0, "Investment Manager", df["Fund Name"] + " (" + df["Ticker"] + ")")

    # 5) Select & order the columns
    df = df[[
        "Investment Manager",
        "3 Year Alpha",
        "5 Year Alpha",
        "3 Year Beta",
        "5 Year Beta",
        "3 Year Upside Capture",
        "3 Year Downside Capture",
        "5 Year Upside Capture",
        "5 Year Downside Capture"
    ]]

    # 6) Display
    st.session_state["step11_summary"] = df.to_dict("records")
    st.dataframe(df)

#───Step 12: Extract “FUND FACTS” & Its Table Details in One Go──────────────────────────────────────────────────────────────────

def step12_process_fund_facts(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # the exact labels and the order you want them in the table
    labels = [
        "Manager Tenure Yrs.",
        "Expense Ratio",
        "Expense Ratio Rank",
        "Total Number of Holdings",
        "Turnover Ratio"
    ]

    records = []
    # scan each factsheet page
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = pdf.pages[pnum-1].extract_text().splitlines()

        for idx, line in enumerate(lines):
            if line.lstrip().upper().startswith("FUND FACTS"):
                # grab the next 8 lines (should contain your 5 labels)
                snippet = lines[idx+1 : idx+1+8]
                rec = {"Fund Name": fund_name, "Ticker": ticker}
                for lab in labels:
                    val = None
                    for ln in snippet:
                        norm = " ".join(ln.strip().split())
                        if norm.startswith(lab):
                            rest = norm[len(lab):].strip(" :\t")
                            m = re.match(r"(-?\d+\.\d+)", rest)
                            val = m.group(1) if m else (rest.split()[0] if rest else None)
                            break
                    rec[lab] = val
                records.append(rec)
                break  # move on to the next page once Fund Facts is processed

    if not records:
        st.warning("No Fund Facts tables found.")
        return

    # save & show
    st.session_state["step12_fund_facts_table"] = records

#───Step 13: Extract Risk‑Adjusted Returns Metrics──────────────────────────────────────────────────────────────────

def step13_process_risk_adjusted_returns(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # which metrics to pull
    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    num_rx  = re.compile(r"-?\d+\.\d+")

    records = []
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()

        # find the heading
        for idx, line in enumerate(lines):
            norm = " ".join(line.strip().split()).upper()
            if norm.startswith("RISK-ADJUSTED RETURNS"):
                snippet = lines[idx+1 : idx+1+6]  # grab next few lines
                rec = {"Fund Name": fund_name, "Ticker": ticker}

                for metric in metrics:
                    # find the snippet line for this metric
                    text_line = next(
                        ( " ".join(ln.strip().split())
                          for ln in snippet
                          if ln.strip().upper().startswith(metric.upper()) ),
                        None
                    ) or ""
                    # extract up to 4 numbers
                    nums = num_rx.findall(text_line)
                    nums += [None] * (4 - len(nums))

                    # assign into rec
                    rec[f"{metric} 1Yr"]  = nums[0]
                    rec[f"{metric} 3Yr"]  = nums[1]
                    rec[f"{metric} 5Yr"]  = nums[2]
                    rec[f"{metric} 10Yr"] = nums[3]

                records.append(rec)
                break  # done with this page

    if not records:
        st.warning("No 'RISK‑ADJUSTED RETURNS' tables found.")
        return

    # save & show
    st.session_state["step13_risk_adjusted_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#───Step 14: Peer Risk-Adjusted Return Rank──────────────────────────────────────────────────────────────────

def step14_extract_peer_risk_adjusted_return_rank(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.write("Peer Rank")

    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    page_map = {
        f["Page #"]:(f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    records = []

    for pnum, (fund, ticker) in page_map.items():
        page = pdf.pages[pnum-1]
        text = page.extract_text() or ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        # 1) locate the Risk-Adjusted Returns header
        try:
            risk_idx = next(i for i, ln in enumerate(lines)
                            if "RISK-ADJUSTED RETURNS" in ln.upper())
        except StopIteration:
            st.warning(f"⚠️ {fund} ({ticker}): Risk-Adjusted header not found.")
            continue

        # 2) find all “1 Yr 3 Yrs 5 Yrs 10 Yrs” lines after that
        header_idxs = [i for i, ln in enumerate(lines)
                       if re.match(r"1\s*Yr", ln)]
        peer_header_idxs = [i for i in header_idxs if i > risk_idx]

        if not peer_header_idxs:
            st.warning(f"⚠️ {fund} ({ticker}): peer header not found.")
            continue

        # take the *second* header occurrence (first is Risk-Adjusted, next is Peer)
        peer_hdr = peer_header_idxs[0] if len(peer_header_idxs)==1 else peer_header_idxs[1]

        rec = {"Fund Name": fund, "Ticker": ticker}

        # 3) read the three lines immediately below that header
        for offset, metric in enumerate(metrics, start=1):
            if peer_hdr + offset < len(lines):
                parts = lines[peer_hdr + offset].split()
                # parts[0:2] = metric name words, parts[2:6] = the four integer ranks
                vals = parts[2:6] if len(parts) >= 6 else []
            else:
                vals = []

            # fill into record (pad with None if missing)
            for idx, period in enumerate(["1Yr","3Yr","5Yr","10Yr"]):
                rec[f"{metric} {period}"] = vals[idx] if idx < len(vals) else None

            if len(vals) < 4:
                st.warning(f"⚠️ {fund} ({ticker}): only {len(vals)} peer values found for '{metric}'.")

        records.append(rec)

    if not records:
        st.warning("❌ No Peer Risk-Adjusted Return Rank data extracted.")
        return

    df = pd.DataFrame(records)
    st.session_state["step14_peer_rank_table"] = records
    st.dataframe(df, use_container_width=True)

#───Step 14.5: IPS Fail Table──────────────────────────────────────────────────────────────────

def step14_5_ips_fail_table():
    import streamlit as st
    import pandas as pd

    df = st.session_state.get("ips_icon_table")
    if df is None or df.empty:
        return

    fail_df = df[df["IPS Watch Status"].isin(["FW", "IW"])][["Fund Name", "IPS Watch Status"]]
    if fail_df.empty:
        return

    table_html = fail_df.rename(columns={
        "Fund Name": "Fund",
        "IPS Watch Status": "Watch Status"
    }).to_html(index=False, border=0, justify="center", classes="ips-fail-table")

    st.markdown(f"""
    <div style='
        background: linear-gradient(120deg, #e6f0fb 85%, #c8e0f6 100%);
        color: #23395d;
        border-radius: 1.3rem;
        box-shadow: 0 2px 14px rgba(44,85,130,0.08), 0 1px 4px rgba(36,67,105,0.07);
        padding: 1.6rem 2.0rem 1.6rem 2.0rem;
        max-width: 650px;
        margin: 1.4rem auto 1.2rem auto;
        border: 1.5px solid #b5d0eb;'>
        <div style='font-weight:700; color:#23395d; font-size:1.15rem; margin-bottom:0.5rem; letter-spacing:-0.5px;'>
            Funds on Watch
        </div>
        <div style='font-size:1rem; margin-bottom:1rem; color:#23395d;'>
            The following funds failed five or more IPS criteria and are currently on watch.
        </div>
        {table_html}
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <style>
    .ips-fail-table {
        width: 100%;
        border-collapse: collapse;
        margin-top: 0.7em;
        font-family: 'Inter', 'Segoe UI', Arial, sans-serif;
    }
    .ips-fail-table th, .ips-fail-table td {
        border: none;
        padding: 0.48em 1.1em;
        text-align: left;
        font-size: 1.07em;
    }
    .ips-fail-table th {
        background: #244369;
        color: #fff;
        font-weight: 700;
        letter-spacing: 0.01em;
    }
    .ips-fail-table td {
        color: #244369;
    }
    .ips-fail-table tr:nth-child(even) {background: #e6f0fb;}
    .ips-fail-table tr:nth-child(odd)  {background: #f8fafc;}
    </style>
    """, unsafe_allow_html=True)

#───Step 14.7: Proposal──────────────────────────────────────────────────────────────────

#───Step 15: Single Fund──────────────────────────────────────────────────────────────────

def step15_display_selected_fund():
    import pandas as pd
    import streamlit as st
    import re

    facts = st.session_state.get("fund_factsheets_data", [])
    if not facts:
        st.info("Run Steps 1–14 to populate data before viewing fund details.")
        return

    fund_names = [f["Matched Fund Name"] for f in facts]
    selected_fund = st.selectbox("Select a fund to view details:", fund_names)
    st.session_state.selected_fund = selected_fund

    # --- NEW: pull confirmed proposed funds once, independent of selection ---
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposed_fund_names = (
        confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        if not confirmed_proposed_df.empty else []
    )

    st.write(f"Details for: {selected_fund}")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    factsheet_rec = next((row for row in factsheets if row["Matched Fund Name"] == selected_fund), None)
    
    fund_facts_table = st.session_state.get("step12_fund_facts_table", [])
    
    # Filter out metadata rows if present
    fund_facts_table = [row for row in fund_facts_table if row.get("Fund Name") and row.get("Fund Name").lower() != "metadata"]
    
    # Robust matching
    facts_rec = next((row for row in fund_facts_table if row.get("Fund Name") == selected_fund), None)
    if not facts_rec and factsheet_rec:
        factsheet_ticker = factsheet_rec.get("Matched Ticker")
        facts_rec = next((row for row in fund_facts_table if row.get("Ticker") == factsheet_ticker), None)
    if not facts_rec:
        facts_rec = next(
            (row for row in fund_facts_table if selected_fund.lower() in row.get("Fund Name", "").lower()),
            None
        )
    
    left_box = (
        f"""<div style='
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.2rem;
            box-shadow: 0 2px 12px rgba(44,85,130,0.09), 0 1px 4px rgba(36,67,105,0.07);
            padding: 1rem 1.2rem;
            min-width: 220px;
            max-width: 260px;
            margin: 0.3rem 1.2rem 0.3rem 0;
            border: 1.2px solid #b5d0eb;
            font-size: 1rem;
            display: inline-block;
            vertical-align: top;'>
            <div><b>Category:</b> {factsheet_rec.get("Category", "—")}</div>
            <div><b>Benchmark:</b> {factsheet_rec.get("Benchmark", "—")}</div>
            <div><b>Net Assets:</b> {factsheet_rec.get("Net Assets", "—")}</div>
            <div><b>Manager:</b> {factsheet_rec.get("Manager Name", "—")}</div>
            <div><b>Avg. Market Cap:</b> {factsheet_rec.get("Avg. Market Cap", "—")}</div>
        </div>"""
        if factsheet_rec else "<div style='display:inline-block; min-width:220px; color:#666;'>No factsheet info found.</div>"
    )
    
    right_box = (
        f"""<div style='
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.2rem;
            box-shadow: 0 2px 12px rgba(44,85,130,0.09), 0 1px 4px rgba(36,67,105,0.07);
            padding: 1rem 1.2rem;
            min-width: 220px;
            max-width: 260px;
            margin: 0.3rem 0 0.3rem 0;
            border: 1.2px solid #b5d0eb;
            font-size: 1rem;
            display: inline-block;
            vertical-align: top;'>
            <div><b>Manager Tenure:</b> {facts_rec.get("Manager Tenure Yrs.", "—")}</div>
            <div><b>Expense Ratio:</b> {facts_rec.get("Expense Ratio", "—")}</div>
            <div><b>Expense Ratio Rank:</b> {facts_rec.get("Expense Ratio Rank", "—")}</div>
            <div><b>Total Number of Holdings:</b> {facts_rec.get("Total Number of Holdings", "—")}</div>
            <div><b>Turnover Ratio:</b> {facts_rec.get("Turnover Ratio", "—")}</div>
        </div>"""
        if facts_rec else "<div style='display:inline-block; min-width:220px; color:#666;'>No Fund Facts available.</div>"
    )
    
    st.markdown(
        f"""
        <div style='
            width:100%;
            display:flex;
            flex-wrap:wrap;
            justify-content:center;
            align-items:flex-start;
            gap:24px;
            margin: 0.6rem 0 2rem 0;
        '>
            {left_box}{right_box}
        </div>
        """,
        unsafe_allow_html=True
    )


    # --- IPS Results ---
    ips_icon_table = st.session_state.get("ips_icon_table")
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)

    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        if not row.empty:
            row_dict = row.iloc[0].to_dict()
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            row_df = pd.DataFrame([{
                "Category": fs_rec.get("Category", "") if fs_rec else "",
                "Time Period": st.session_state.get("report_date", ""),
                "Plan Assets": "$",  # Or replace with actual variable if you store this elsewhere!
                **{display_columns.get(k, k): v for k, v in row_dict.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row_dict.get("IPS Watch Status", "")
            }])

            def color_bool(v):
                return "background-color: green" if v == "✔" else ("background-color: red" if v == "✗" else "")
            
            def style_status(v):
                if v == "NW": return "background-color: green; color: white; font-weight: 600;"
                if v == "IW": return "background-color: orange; color: white; font-weight: 600;"
                if v == "FW": return "background-color: red; color: white; font-weight: 600;"
                return ""
            
            styled = row_df.style.applymap(color_bool, subset=[str(i) for i in range(1, 12)]) \
                                 .applymap(style_status, subset=["IPS Status"])
            
            st.dataframe(styled, use_container_width=True)

        else:
            st.warning("No IPS screening result found for selected fund.")
    else:
        st.warning("IPS screening table not found. Run earlier steps first.")

    # --- Expense & Return: Table 1 ---
    st.markdown("**Net Expense Ratio**")
    
    def format_expense(val):
        if val is None or val == "":
            return ""
        s = str(val)
        if s.endswith("%"):
            return s
        return f"{s}%"
    
    def build_expense_row(fund_name):
        perf_data = st.session_state.get("fund_performance_data", [])
        item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        net_exp = format_expense(item.get("Net Expense Ratio", ""))
        return {
            "Investment Manager": inv_mgr,
            "Net Expense Ratio": net_exp
        }
    
    # Selected fund row
    row_selected = build_expense_row(selected_fund)
    
    # Proposed fund(s) — persistent, independent of selection
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_expense_row(pf))
    
    # Assemble final table: selected fund first, then proposed(s)
    all_rows = [row_selected] + proposed_rows
    df_ear_table1 = pd.DataFrame(all_rows)
    
    # Save & display
    st.session_state["ear_table1_data"] = df_ear_table1
    st.dataframe(df_ear_table1, use_container_width=True)


    # --- Expense & Return: Table 2: Returns ---
    st.markdown("**Returns**")
    date_label = st.session_state.get("report_date", "QTD")
    
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s == "" else f"{s}%"
    
    def build_return_row(fund_name):
        perf_data = st.session_state.get("fund_performance_data", [])
        item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
    
        qtd   = append_pct(item.get("QTD", ""))
        one   = append_pct(item.get("1Yr", ""))
        three = append_pct(item.get("3Yr", ""))
        five  = append_pct(item.get("5Yr", ""))
        ten   = append_pct(item.get("10Yr", ""))
    
        return {
            "Investment Manager": inv_mgr,
            date_label:           qtd,
            "1 Year":             one,
            "3 Year":             three,
            "5 Year":             five,
            "10 Year":            ten
        }
    
    # Selected fund row
    row_selected = build_return_row(selected_fund)
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_return_row(pf))
    
    # Benchmark row (same as before)
    fs_rec = next((f for f in st.session_state.get("fund_factsheets_data", []) if f["Matched Fund Name"] == selected_fund), {})
    bench_name = fs_rec.get("Benchmark", "") if fs_rec else ""
    bench_ticker = fs_rec.get("Matched Ticker", "") if fs_rec else ""
    bench_inv_mgr = f"{bench_name} ({bench_ticker})" if bench_name and bench_ticker else (bench_name or "Benchmark")
    
    # Need the selected fund's perf_item to get benchmark returns
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    bench_qtd   = append_pct(perf_item.get("Bench QTD", ""))
    bench_one   = append_pct(perf_item.get("Bench 1Yr", ""))
    bench_3yr   = append_pct(perf_item.get("Bench 3Yr", ""))
    bench_5yr   = append_pct(perf_item.get("Bench 5Yr", ""))
    bench_ten   = append_pct(perf_item.get("Bench 10Yr", ""))
    
    row_benchmark = {
        "Investment Manager": bench_inv_mgr,
        date_label:           bench_qtd,
        "1 Year":             bench_one,
        "3 Year":             bench_3yr,
        "5 Year":             bench_5yr,
        "10 Year":            bench_ten
    }
    
    # Assemble: selected, proposed(s), then benchmark
    all_rows = [row_selected] + proposed_rows + [row_benchmark]
    df_ear_table2 = pd.DataFrame(all_rows)
    
    # Save & display
    st.session_state["ear_table2_data"] = df_ear_table2
    st.dataframe(df_ear_table2, use_container_width=True)

    # --- Expense & Return: Table 3: Calendar Returns ---
    st.markdown("**Calendar Returns**")
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    if not fund_cy or not bench_cy:
        st.error("❌ No calendar year returns data found. Ensure Step 8 has been run correctly.")
        return
    
    def build_cy_row(name, is_selected=True):
        # find fund record
        fund_rec = next((r for r in fund_cy if r.get("Name") == name), None)
        if not fund_rec:
            return None
        ticker = fund_rec.get("Ticker", "")
        row = {"Investment Manager": f"{name} ({ticker})" if ticker else name}
        year_cols = [col for col in fund_rec.keys() if re.match(r"20\d{2}", col)]
        for year in year_cols:
            row[year] = fund_rec.get(year, "")
        return row, year_cols
    
    # Selected
    selected_row_tuple = build_cy_row(selected_fund)
    if not selected_row_tuple:
        st.error(f"❌ Could not find data for selected fund: {selected_fund}")
        return
    row_selected, year_cols = selected_row_tuple
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            tup = build_cy_row(pf)
            if tup:
                row_pf, _ = tup
                proposed_rows.append(row_pf)
    
    # Benchmark for selected fund
    bench_rec = None
    fund_rec = next((r for r in fund_cy if r.get("Name") == selected_fund), None)
    if fund_rec:
        bench_rec = next(
            (r for r in bench_cy if r.get("Name") == selected_fund or r.get("Ticker") == fund_rec.get("Ticker")),
            None
        )
    if not bench_rec:
        st.error(f"❌ Could not find benchmark data for selected fund: {selected_fund}")
        return
    bench_ticker = bench_rec.get("Ticker", "")
    bench_name_display = bench_rec.get("Name", "Benchmark")
    row_benchmark = {"Investment Manager": f"{bench_name_display} ({bench_ticker})" if bench_ticker else bench_name_display}
    for year in year_cols:
        row_benchmark[year] = bench_rec.get(year, "")
    
    # Assemble in order: selected, proposed(s), benchmark
    all_rows = [row_selected] + proposed_rows + [row_benchmark]
    df_ear_table3 = pd.DataFrame(all_rows, columns=["Investment Manager"] + year_cols)
    
    st.session_state["ear_table3_data"] = df_ear_table3
    st.dataframe(df_ear_table3, use_container_width=True)



    # --- Risk Adjusted Returns Table 1 ---
    st.markdown("**MPT Statistics Summary**")
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_risk_analysis_5yr", []) or st.session_state.get("step10_mpt_stats", [])
    
    def build_mpt_row(fund_name, stats3_list, stats5_list):
        stats3 = next((r for r in stats3_list if r["Fund Name"] == fund_name), {})
        stats5 = next((r for r in stats5_list if r["Fund Name"] == fund_name), {})
        ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        return {
            "Investment Manager":        inv_mgr,
            "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
            "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
            "3 Year Beta":               stats3.get("3 Year Beta", ""),
            "5 Year Beta":               stats5.get("5 Year Beta", ""),
            "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
            "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
            "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
            "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
        }
    
    # Selected fund row
    row_selected = build_mpt_row(selected_fund, mpt3, mpt5)
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_mpt_row(pf, mpt3, mpt5))
    
    # Assemble: selected then proposed(s)
    all_rows = [row_selected] + proposed_rows
    df_raj_table1 = pd.DataFrame(all_rows)
    
    st.session_state["raj_table1_data"] = df_raj_table1
    st.dataframe(df_raj_table1, use_container_width=True)

    # --- Risk Adjusted Returns Table 2 ---
    st.markdown("**Risk-Adjusted Returns / Peer Ranking %**")
    risk_table = st.session_state.get("step13_risk_adjusted_table", [])
    peer_table = st.session_state.get("step14_peer_rank_table", [])
    
    def build_ratio_row(fund_name):
        risk_rec = next((r for r in risk_table if r.get("Fund Name") == fund_name), {})
        peer_rec = next((r for r in peer_table if r.get("Fund Name") == fund_name), {})
        ticker = risk_rec.get("Ticker") or peer_rec.get("Ticker", "")
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        def frac(metric, period):
            r = risk_rec.get(f"{metric} {period}", "")
            p = peer_rec.get(f"{metric} {period}", "")
            return f"{r} / {p}"
        return {
            "Investment Manager": inv_mgr,
            "3 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "3Yr"),
            "5 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "5Yr"),
            "3 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "3Yr"),
            "5 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "5Yr"),
            "3 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "3Yr"),
            "5 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "5Yr"),
        }
    
    # Selected fund
    rows = [build_ratio_row(selected_fund)]
    
    # Proposed fund(s)
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_ratio_row(pf))
    
    df_raj_table2 = pd.DataFrame(rows)

    # Save for Step 17 to use
    st.session_state["raj_table2_data"] = df_raj_table2
    st.dataframe(df_raj_table2, use_container_width=True)

    # --- Qualitative Factors Table 1 ---
    st.markdown("**Manager Tenure**")
    blocks = st.session_state.get("fund_blocks", [])
    
    def build_tenure_row(fund_name):
        block = next((b for b in blocks if b["Fund Name"] == fund_name), {})
        raw_tenure = next((m["Info"] for m in block.get("Metrics", []) if m["Metric"] == "Manager Tenure"), "")
        m = re.search(r"(\d+(\.\d+)?)", raw_tenure)
        tenure = f"{m.group(1)} years" if m else raw_tenure
        # attempt to get ticker from performance_data for formatting
        perf_data = st.session_state.get("fund_performance_data", [])
        perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (perf_item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        return {
            "Investment Manager": inv_mgr,
            "Manager Tenure":     tenure
        }
    
    rows = [build_tenure_row(selected_fund)]
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_tenure_row(pf))
    
    df_qualfact_table1 = pd.DataFrame(rows)

    # Save for Step 17 to use
    st.session_state["qualfact_table1_data"] = df_qualfact_table1
    st.dataframe(df_qualfact_table1, use_container_width=True)

    
    # --- Qualitative Factors Table 2 ---
    st.markdown("**Assets**")
    facts = st.session_state.get("fund_factsheets_data", [])
    
    def build_asset_row(fund_name):
        # factsheet for the fund
        fs_rec = next((f for f in facts if f["Matched Fund Name"] == fund_name), None)
        # performance for ticker formatting
        perf_data = st.session_state.get("fund_performance_data", [])
        perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (perf_item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        assets = fs_rec.get("Net Assets", "") if fs_rec else ""
        avg_cap = fs_rec.get("Avg. Market Cap", "") if fs_rec else ""
        return {
            "Investment Manager":            inv_mgr,
            "Assets Under Management":       assets,
            "Average Market Capitalization": avg_cap
        }
    
    rows = [build_asset_row(selected_fund)]
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_asset_row(pf))


    df_qualfact_table2 = pd.DataFrame(rows)
    
    # Save for Step 17 to use
    st.session_state["qualfact_table2_data"] = df_qualfact_table2
    st.dataframe(df_qualfact_table2, use_container_width=True)

#───Bullet Points──────────────────────────────────────────────────────────────────
def step16_3_selected_overview_lookup(pdf, context_lines=3, min_score=50):
    import streamlit as st
    from rapidfuzz import fuzz
    import re

    selected_fund = st.session_state.get("selected_fund")
    if not selected_fund:
        return {}

    lookup = {
        "Fund": selected_fund,
        "Best Page": None,
        "Match Score": 0,
        "Found Investment Overview": False,
        "Overview Paragraph": "",
        "Overview Context": "",
    }

    # Normalize helper
    def normalize(s):
        return re.sub(r"[^A-Za-z0-9 ]+", "", (s or "").lower()).strip()

    norm_target = re.sub(r"\s*\(.*\)$", "", selected_fund).strip()
    norm_target_clean = normalize(norm_target)

    # 1) If factsheet data already has a match for this fund, use its page
    factsheets = st.session_state.get("fund_factsheets_data", [])
    best_page = None
    best_score = 0
    if factsheets:
        # Try to find the row where Matched Fund Name equals selected_fund (allow fuzzy)
        for f in factsheets:
            name = f.get("Matched Fund Name", "")
            page_no = f.get("Page #")
            if not name or not page_no:
                continue
            score = fuzz.token_sort_ratio(normalize(name), normalize(selected_fund))
            if score > best_score:
                best_score = score
                best_page = page_no

    # 2) If we didn't get a confident page from existing factsheets, scan pages starting at factsheets_page
    if best_page is None or best_score < 80:  # allow fallback if existing was weak
        factsheets_start = st.session_state.get("factsheets_page") or 1
        for i in range(factsheets_start - 1, len(pdf.pages)):
            page = pdf.pages[i]
            text = (page.extract_text() or "").lower()
            name_score = fuzz.token_sort_ratio(norm_target_clean, normalize(text))
            # also consider ticker if available
            ticker = ""
            perf_data = st.session_state.get("fund_performance_data", [])
            item = next((x for x in perf_data if x.get("Fund Scorecard Name") == selected_fund), {})
            ticker = (item.get("Ticker") or "").upper().strip()
            ticker_score = 0
            if ticker:
                ticker_score = 100 if re.search(rf"\b{re.escape(ticker.lower())}\b", text) else 0
            combined_score = max(name_score, ticker_score)
            if combined_score > best_score:
                best_score = combined_score
                best_page = i + 1

    lookup["Best Page"] = best_page
    lookup["Match Score"] = best_score

    if best_page is None:
        st.warning(f"Could not locate any candidate page for '{selected_fund}'.")
        st.session_state["step16_3_selected_overview_lookup"] = lookup
        return lookup

    page_obj = pdf.pages[best_page - 1]
    raw_lines = (page_obj.extract_text() or "").splitlines()

    # Find heading
    heading_idx = None
    for idx, ln in enumerate(raw_lines):
        if re.search(r"INVESTMENT\s+OVERVIEW", ln, re.IGNORECASE):
            heading_idx = idx
            break

    if heading_idx is None:
        st.warning(f"'INVESTMENT OVERVIEW' heading not found on page {best_page} for selected fund '{selected_fund}'. Best score {best_score}.")
        st.session_state["step16_3_selected_overview_lookup"] = lookup
        return lookup

    # Context snippet
    start_ctx = max(0, heading_idx - context_lines)
    end_ctx = min(len(raw_lines), heading_idx + context_lines + 1)
    context_snippet = "\n".join(raw_lines[start_ctx:end_ctx])
    lookup["Overview Context"] = context_snippet

    # Collect paragraph beneath heading
    def is_new_section_heading(line):
        stripped = line.strip()
        if not stripped:
            return False
        word_count = len(stripped.split())
        has_digit = bool(re.search(r"\d", stripped))
        if (stripped.upper() == stripped or (stripped.istitle() and not stripped.endswith("."))) and word_count <= 7 and not has_digit:
            return True
        return False

    collected = []
    for ln in raw_lines[heading_idx + 1:]:
        if not ln.strip():
            if collected:
                break
            continue
        if is_new_section_heading(ln):
            break
        collected.append(ln.strip())
        if len(collected) >= 80:
            break

    full_text = " ".join(collected)
    sentences = safe_split_sentences(full_text)
    overview_paragraph = " ".join(sentences[:3]).strip() if sentences else full_text.strip()

    if overview_paragraph:
        lookup["Found Investment Overview"] = True
        lookup["Overview Paragraph"] = overview_paragraph

    st.session_state["step16_3_selected_overview_lookup"] = lookup
    return lookup

    
#───Bullet Points──────────────────────────────────────────────────────────────────
def markdown_bold_to_html(text: str) -> str:
    # escape to avoid injection, then convert **bold** to <strong>...</strong>
    escaped = html.escape(text)
    return re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', escaped)

def step16_bullet_points(pdf=None):
    import streamlit as st

    selected_fund = st.session_state.get("selected_fund")
    if not selected_fund:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    # Ensure overview lookup has been performed so the overview bullet can be inserted
    if pdf is not None:
        existing = st.session_state.get("step16_3_selected_overview_lookup", {})
        if existing.get("Fund") != selected_fund:
            step16_3_selected_overview_lookup(pdf, context_lines=3, min_score=50)

    perf_data = st.session_state.get("fund_performance_data", [])
    item = next((x for x in perf_data if x.get("Fund Scorecard Name") == selected_fund), None)
    if not item:
        st.error(f"❌ Performance data for '{selected_fund}' not found.")
        return

    bullets = []

    # Bullet 1: Performance vs. Benchmark, using template
    template = st.session_state.get("bullet_point_templates", [""])[0]
    b1 = template
    for fld, val in item.items():
        b1 = b1.replace(f"[{fld}]", str(val))
    bullets.append(b1)

    # Get IPS status
    ips_icon_table = st.session_state.get("ips_icon_table")
    ips_status = None
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        ips_status = row.iloc[0]["IPS Watch Status"] if not row.empty else None

    # Bullet 2: Watch status and return comparison
    if ips_status == "NW":
        b2 = "- This fund is **not on watch**."
        bullets.append(b2)
    else:
        status_label = (
            "Formal Watch" if ips_status == "FW" else
            "Informal Watch" if ips_status == "IW" else
            ips_status or "on watch"
        )

        def to_float(x):
            try:
                return float(x)
            except:
                return 0.0

        three   = to_float(item.get("3Yr"))
        bench3  = to_float(item.get("Bench 3Yr"))
        five    = to_float(item.get("5Yr"))
        bench5  = to_float(item.get("Bench 5Yr"))
        bps3 = round((three - bench3) * 100, 1)
        bps5 = round((five  - bench5) * 100, 1)

        peer = st.session_state.get("step14_peer_rank_table", [])
        raw3 = next((r.get("Sharpe Ratio 3Yr") or r.get("Sharpe Ratio Rank 3Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        raw5 = next((r.get("Sharpe Ratio 5Yr") or r.get("Sharpe Ratio Rank 5Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        try:
            pos3 = "top" if raw3 and int(raw3) <= 50 else "bottom"
        except:
            pos3 = "bottom"
        try:
            pos5 = "top" if raw5 and int(raw5) <= 50 else "bottom"
        except:
            pos5 = "bottom"

        b2 = (
            f"The fund is now on **{status_label}**. Its 3-year return trails the benchmark by "
            f"{bps3} bps ({three:.2f}% vs. {bench3:.2f}%) and its 5-year return trails by "
            f"{bps5} bps ({five:.2f}% vs. {bench5:.2f}%). Its 3-Yr Sharpe ranks in the {pos3} half of peers "
            f"and the {pos5} half of its 5-Yr Sharpe ranks."
        )
        bullets.append(b2)

    # --- Investment Overview bullet ---
    overview_info = st.session_state.get("step16_3_selected_overview_lookup", {}) or {}
    overview_paragraph = overview_info.get("Overview Paragraph", "")
    if overview_paragraph:
        # safe sentence splitting with fallback
        def safe_split_sentences_local(text):
            abbrev_map = {
                "U.S.": "__US__", "U.S": "__US__", "U.K.": "__UK__", "U.K": "__UK__",
                "e.g.": "__EG__", "e.g": "__EG__", "i.e.": "__IE__", "i.e": "__IE__",
                "etc.": "__ETC__", "etc": "__ETC__",
            }
            protected = text
            for k, v in abbrev_map.items():
                protected = protected.replace(k, v)
            sentences = re.split(r'(?<=[\.!?])\s+', protected.strip())
            def restore(s):
                for k, v in abbrev_map.items():
                    s = s.replace(v, k)
                return s
            return [restore(s).strip() for s in sentences if s.strip()]

        sentences = safe_split_sentences_local(overview_paragraph)
        overview_bullet = " ".join(sentences[:3]) if sentences else overview_paragraph
        b_overview = overview_bullet  # no prefix
        bullets.append(b_overview)

    # Bullet 3: Action for Formal Watch only
    if ips_status == "FW":
        confirmed = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
        proposals = []
        if not confirmed.empty:
            seen = set()
            for name, ticker in zip(confirmed["Fund Scorecard Name"], confirmed["Ticker"]):
                display = f"{name} ({ticker})" if ticker else name
                if display not in seen:
                    seen.add(display)
                    proposals.append(display)
        replacement = ", ".join(proposals) if proposals else "a proposed fund"
        b3 = f"**Action:** Consider replacing this fund with {replacement}."
        bullets.append(b3)

    # Persist updated bullets
    st.session_state["bullet_points"] = bullets


#───Bullet Points 2──────────────────────────────────────────────────────────────────
from rapidfuzz import fuzz
import re
import streamlit as st
import pandas as pd

def step16_5_locate_proposed_factsheets_with_overview(pdf, context_lines=3, min_score=60):
    """
    For each confirmed proposed fund, locate its best matching factsheet page and extract the
    first few sentences under the 'INVESTMENT OVERVIEW' heading. Returns per-fund metadata,
    confidence scores, context snippet, and extracted overview paragraph.
    """
    def normalize(text: str) -> str:
        return re.sub(r"[^A-Za-z0-9 ]+", "", (text or "").lower()).strip()

    def split_into_sentences(text: str) -> list[str]:
        abbrev_map = {
            "U.S.": "__US__", "U.S": "__US__", "U.K.": "__UK__", "U.K": "__UK__",
            "e.g.": "__EG__", "e.g": "__EG__", "i.e.": "__IE__", "i.e": "__IE__",
            "etc.": "__ETC__", "etc": "__ETC__",
        }
        protected = text
        for k, v in abbrev_map.items():
            protected = protected.replace(k, v)
        sentences = re.split(r'(?<=[\.!?])\s+', protected.strip())
        def restore(s):
            for k, v in abbrev_map.items():
                s = s.replace(v, k)
            return s
        return [restore(s).strip() for s in sentences if s.strip()]

    def is_new_section_heading(line: str) -> bool:
        stripped = line.strip()
        if not stripped:
            return False
        word_count = len(stripped.split())
        if word_count > 7:
            return False
        has_digit = bool(re.search(r"\d", stripped))
        if (stripped.upper() == stripped or (stripped.istitle() and not stripped.endswith("."))) and not has_digit:
            return True
        return False

    confirmed = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    factsheets_start = st.session_state.get("factsheets_page") or 1
    results: dict[str, dict] = {}

    # Preload pages once
    pages_cache = []
    for i in range(factsheets_start - 1, len(pdf.pages)):
        page_obj = pdf.pages[i]
        raw_text = page_obj.extract_text() or ""
        pages_cache.append((i + 1, raw_text, page_obj))  # 1-based page number

    for _, row in confirmed.iterrows():
        raw_name = row.get("Fund Scorecard Name", "")
        ticker = (row.get("Ticker") or "").upper().strip()
        fund_key = raw_name.strip().rstrip(".")
        norm_expected = normalize(fund_key)

        best_candidate = {"page": None, "score": 0, "match_type": None}
        # 1) Strong match via exact ticker presence
        if ticker:
            for page_num, text, page_obj in pages_cache:
                if re.search(rf"\b{re.escape(ticker.lower())}\b", (text or "").lower()):
                    best_candidate = {"page": page_num, "score": 100, "match_type": "ticker"}
                    break

        # 2) Fallback fuzzy name match
        if best_candidate["page"] is None:
            for page_num, text, page_obj in pages_cache:
                score = fuzz.token_sort_ratio(norm_expected, normalize(text))
                if score > best_candidate["score"]:
                    best_candidate = {"page": page_num, "score": score, "match_type": "name"}

        fund_result = {
            "Fund": fund_key,
            "Ticker": ticker,
            "Best Page": best_candidate["page"],
            "Match Score": best_candidate["score"],
            "Match Type": best_candidate["match_type"],
            "Found Investment Overview": False,
            "Overview Context": "",
            "Overview Paragraph": "",
            "Overview Bold Detected": False,
        }

        if best_candidate["page"] is None or best_candidate["score"] < min_score:
            st.warning(
                f"Could not confidently locate factsheet for proposed fund '{fund_key}' "
                f"({ticker}), best score {best_candidate['score']}."
            )
            results[fund_key] = fund_result
            continue

        page_obj = pdf.pages[best_candidate["page"] - 1]
        raw_lines = (page_obj.extract_text() or "").splitlines()
        words = page_obj.extract_words(use_text_flow=True, extra_attrs=["fontname"])

        # Locate heading "INVESTMENT OVERVIEW"
        target_re = re.compile(r"INVESTMENT\s+OVERVIEW", re.IGNORECASE)
        heading_idx = None
        bold_detected = False

        # Try via word-pair heuristic to detect bold heading
        for i in range(len(words) - 1):
            pair = f"{words[i]['text']} {words[i + 1]['text']}"
            if target_re.fullmatch(re.sub(r"\s+", " ", pair).strip()):
                # Find the line containing the phrase
                for li, ln in enumerate(raw_lines):
                    if re.search(r"investment\s+overview", ln, re.IGNORECASE):
                        heading_idx = li
                        break
                fontnames = (words[i].get("fontname", "") or "").lower() + " " + (words[i+1].get("fontname", "") or "").lower()
                if any(b in fontnames for b in ["bold", "bd", "black"]):
                    bold_detected = True
                break

        # Fallback: line-wise search
        if heading_idx is None:
            for li, ln in enumerate(raw_lines):
                if re.search(r"investment\s+overview", ln, re.IGNORECASE):
                    heading_idx = li
                    break

        if heading_idx is None:
            st.warning(
                f"'INVESTMENT OVERVIEW' heading not found on page {best_candidate['page']} "
                f"for proposed fund '{fund_key}' ({ticker})."
            )
            results[fund_key] = fund_result
            continue

        # Context snippet for diagnostics
        start_ctx = max(0, heading_idx - context_lines)
        end_ctx = min(len(raw_lines), heading_idx + context_lines + 1)
        fund_result["Overview Context"] = "\n".join(raw_lines[start_ctx:end_ctx])
        fund_result["Overview Bold Detected"] = bold_detected

        # Extract paragraph below heading
        collected = []
        for ln in raw_lines[heading_idx + 1:]:
            if not ln.strip():
                if collected:
                    break
                else:
                    continue
            if is_new_section_heading(ln):
                break
            collected.append(ln.strip())
            if len(collected) >= 60:
                break

        full_text = " ".join(collected)
        sentences = split_into_sentences(full_text)
        paragraph = " ".join(sentences[:3]).strip() if sentences else ""

        # Fallback hierarchy
        if not paragraph:
            if full_text:
                paragraph = full_text
            else:
                next_nonempty = next((l for l in raw_lines[heading_idx + 1:] if l.strip()), "")
                paragraph = next_nonempty.strip()

        fund_result["Overview Paragraph"] = paragraph
        fund_result["Found Investment Overview"] = bool(paragraph)
        results[fund_key] = fund_result

    st.session_state["step16_5_proposed_overview_lookup"] = results
    return results


#───Build Powerpoint─────────────────────────────────────────────────────────────────
def step17_export_to_ppt():
    import streamlit as st
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from io import BytesIO
    from copy import deepcopy
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    import pandas as pd



# ───── Stand Alone Helpers ────────────────────────────────────────────────────────────
    
    def fill_bullet_points(slide, placeholder="[Bullet Point 1]", bullets=None):
        """
        Finds the first shape on `slide` whose text_frame contains `placeholder`,
        clears it, and writes each item in `bullets` as a Cambria-11pt bullet.
        Returns True if replacement occurred.
        """
        from pptx.util import Pt
    
        if bullets is None:
            bullets = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # look for the placeholder in any paragraph
            if not any(placeholder in (p.text or "") for p in tf.paragraphs):
                continue
    
            # clear existing paragraphs by resetting .text
            tf.text = ""
    
            # write each string as its own bullet
            for text in bullets:
                p = tf.add_paragraph()
                p.text      = text
                p.level     = 0
                p.font.name = "Cambria"
                p.font.size = Pt(11)
            return True
    
        return False

    def truncate_to_n_sentences(text, n=3):
        sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        if len(sentences) <= n:
            return " ".join(sentences).strip()
        return " ".join(sentences[:n]).strip()

    def lookup_overview_paragraph(label, lookup_dict):
        """
        Given a proposed label like "Fund Name (TICK)" find best matching key in lookup_dict
        (which was populated in step16_5_proposed_overview_lookup) using fuzzy matching,
        then return the Overview Paragraph.
        """
        import re
        from rapidfuzz import fuzz
    
        # canonicalize: strip ticker parenthesis and punctuation
        base_name = re.sub(r"\s*\(.*\)$", "", label).strip()
        def normalize(s):
            return re.sub(r"[^A-Za-z0-9 ]+", "", s or "").strip().lower()
    
        target = normalize(base_name)
        best_key = None
        best_score = -1
        for key in lookup_dict.keys():
            score = fuzz.token_sort_ratio(target, normalize(key))
            if score > best_score:
                best_score = score
                best_key = key
    
        if best_key and best_score >= 60:  # threshold you can tweak
            return lookup_dict.get(best_key, {}).get("Overview Paragraph", "")
        # fallback: try exact base_name
        return lookup_dict.get(base_name, {}).get("Overview Paragraph", "")


    
    selected = st.session_state.get("selected_fund")
    if not selected:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    # Get confirmed proposed funds (name + ticker)
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposed = []
    if not confirmed_proposed_df.empty:
        for _, row in confirmed_proposed_df.iterrows():
            name = row.get("Fund Scorecard Name", "")
            ticker = row.get("Ticker", "")
            label = f"{name} ({ticker})" if ticker else name
            proposed.append(label)
    proposed = proposed[:2]  # template supports up to two

    template_path = "assets/writeup&rec_templates.pptx"
    try:
        prs = Presentation(template_path)
    except Exception as e:
        st.error(f"Could not load PowerPoint template: {e}")
        return

    def fill_table_with_styles(table, df_table, bold_row_idx=None, first_col_white=True):
        for i in range(min(len(df_table), len(table.rows) - 1)):
            for j in range(min(len(df_table.columns), len(table.columns))):
                val = df_table.iloc[i, j]
                cell = table.cell(i + 1, j)
                cell.text = str(val) if val is not None else ""
                cell.vertical_alignment = MSO_VERTICAL_ANCHOR.MIDDLE
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        if j == 0:
                            run.font.color.rgb = RGBColor(255, 255, 255) if first_col_white else RGBColor(0, 0, 0)
                        else:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.bold = (bold_row_idx is not None and i == bold_row_idx)

    def fill_text_placeholder_preserving_format(slide, placeholder_text, replacement_text):
        replaced = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                if placeholder_text in full_text:
                    new_text = full_text.replace(placeholder_text, replacement_text)
                    runs = paragraph.runs
                    if len(runs) == 1:
                        runs[0].text = new_text
                    else:
                        for run in runs:
                            run.text = ""
                        avg_len = len(new_text) // len(runs)
                        idx = 0
                        for run in runs[:-1]:
                            run.text = new_text[idx:idx+avg_len]
                            idx += avg_len
                        runs[-1].text = new_text[idx:]
                    replaced = True
        return replaced

    def fill_bullet_points(slide, placeholder="[Bullet Point 1]", bullets=None):
        if bullets is None:
            bullets = st.session_state.get("bullet_points", [])
        if not bullets:
            bullets = ["Performance exceeded benchmark.", "No watch status.", "No action required."]
        replaced = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if any(placeholder in p.text for p in tf.paragraphs):
                tf.clear()
                for b in bullets:
                    p_new = tf.add_paragraph()
                    clean_text = b.replace("**", "")
                    p_new.text = clean_text
                    p_new.level = 0
                    p_new.font.name = "Cambria"
                    p_new.font.size = Pt(11)
                    p_new.font.color.rgb = RGBColor(0, 0, 0)
                    p_new.font.bold = False  # <-- no longer forcing bold
                replaced = True
                break
        return replaced
    # ───── 1) Load template ────────────────────────────────────────────────────────────
    prs = Presentation("assets/writeup&rec_templates.pptx")
    selected = st.session_state.get("selected_fund", "")
    confirmed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposal_names = confirmed_df["Fund Scorecard Name"].dropna().unique().tolist()
    
    # ───── 2) Pull in session data ──────────────────────────────────────────────────────
    selected = st.session_state.get("selected_fund", "")
    ips_icon_table = st.session_state.get("ips_icon_table")
    bullets = st.session_state.get("bullet_points", [])
    ear_df = st.session_state.get("ear_table1_data")  # DataFrame for Expense & Return Table 1
    facts = st.session_state.get("fund_factsheets_data", [])
        # Pull in the DataFrame you saved in Step 15
    df2 = st.session_state.get("ear_table2_data", pd.DataFrame())
    date_label = st.session_state.get("report_date", "")

    # ───── 3) Validate IPS data ──────────────────────────────────────────────────────────
    if ips_icon_table is None or ips_icon_table.empty:
        st.error("IPS screening table not found. Run earlier steps first.")
        return
    row = ips_icon_table[ips_icon_table["Fund Name"] == selected]
    if row.empty:
        st.error("No IPS screening result found for selected fund.")
        return
    row_dict = row.iloc[0].to_dict()

    # ───── 4) FIRST SLIDE LOGIC (Fund Name, IPS table, bullets) ────────────────────────
    slide = prs.slides[0]

    # — Replace "Fund Name" placeholder ─────────────────────────────────────────────
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "Fund Name" in shape.text_frame.text:
            p = shape.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.text = selected
            run.font.name = "Cambria"
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.underline = True
            break

    # — Fill IPS table ───────────────────────────────────────────────────────────────
    def style_run(run, text):
        run.text = text
        run.font.name = "Cambria"
        run.font.size = Pt(12)
        if text == "✔":
            run.font.color.rgb = RGBColor(0, 128, 0)
        elif text == "✗":
            run.font.color.rgb = RGBColor(255, 0, 0)

    # gather values for the IPS table row
    facts_rec = next((f for f in facts if f.get("Matched Fund Name") == selected), {})
    report_date = st.session_state.get("report_date", "")
    vals = [
        facts_rec.get("Category", ""),
        report_date,
        "$",
    ] + [str(row_dict.get(f"IPS Investment Criteria {i}", "")) for i in range(1, 12)] + [row_dict.get("IPS Watch Status", "")]

    table_shape = next((sh for sh in slide.shapes if sh.has_table), None)
    if not table_shape:
        st.error("No table found on the first slide.")
        return
    table = table_shape.table

    # Category → row 1, col 0
    cell = table.cell(1, 0)
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        style_run(para.runs[0], vals[0])
    else:
        style_run(para.add_run(), vals[0])

    # Other values → bottom row
    bottom = len(table.rows) - 1
    for idx, text in enumerate(vals[1:], start=1):
        cell = table.cell(bottom, idx)
        para = cell.text_frame.paragraphs[0]
        if para.runs:
            style_run(para.runs[0], text)
        else:
            style_run(para.add_run(), text)

    # — Fill bullets textbox ──────────────────────────────────────────────────────────
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "[Bullet Point 1]" in shape.text_frame.text:
            tf = shape.text_frame
            tf.text = ""
            for i, bp in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bp
                p.level = 0
                p.font.name = "Cambria"
                p.font.size = Pt(11)
            break

    # ───── Proposal Slides: Replace Headings, Fill Bullets, Delete Extras ─────────────
    import re
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    # build proposal_names & overview_map
    confirmed_df   = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposal_names = confirmed_df["Fund Scorecard Name"].dropna().tolist()
    if not proposal_names:
        ear_df = st.session_state.get("ear_table1_data", pd.DataFrame())
        if "Investment Manager" in ear_df.columns:
            proposal_names = [
                nm.split(" (")[0]
                for nm in ear_df["Investment Manager"].iloc[1:]
            ]
    overview_map = st.session_state.get("step16_5_proposed_overview_lookup", {})
    
    # Replace [Replacement 1]–[Replacement 5], fill bullets, delete extras
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.util import Pt
    import re
    to_delete = []
    
    for i in range(1, 6):
        # a) find the slide with "[Replacement i]"
        slide_idx = None
        slide_obj = None
        for idx, sl in enumerate(prs.slides):
            for shp in sl.shapes:
                if not shp.has_text_frame:
                    continue
                if f"[Replacement {i}]" in (shp.text_frame.text or ""):
                    slide_idx = idx
                    slide_obj = sl
                    break
            if slide_obj:
                break
        if slide_obj is None:
            continue
    
        # b) if there's a fund, replace and fill; else mark for deletion
        if i <= len(proposal_names):
            fund_name = proposal_names[i - 1]
    
            # replace the heading placeholder
            for shp in slide_obj.shapes:
                if not shp.has_text_frame:
                    continue
                txt = shp.text_frame.text or ""
                if f"[Replacement {i}]" in txt:
                    shp.text_frame.text = txt.replace(f"[Replacement {i}]", fund_name)
                    break
    
            # fill the BODY placeholder with overview bullets
            overview = overview_map.get(fund_name, {}).get("Overview Paragraph", "")
            if overview:
                sentences = [
                    s.strip()
                    for s in re.split(r'(?<=[.!?])\s+', overview)
                    if s.strip()
                ][:3]
    
                # target the BODY placeholder directly
                body_ph = next(
                    (ph for ph in slide_obj.placeholders
                        if ph.placeholder_format.type == PP_PLACEHOLDER.BODY),
                    None
                )
                if body_ph:
                    tf = body_ph.text_frame
                    tf.clear()  # clear out the old placeholder text
                    for line in sentences:
                        p = tf.add_paragraph()
                        p.text      = line
                        p.level     = 0
                        p.font.name = "Cambria"
                        p.font.size = Pt(11)
        else:
            to_delete.append(slide_idx)
    
    # c) delete extra slides (reverse order to keep indices valid)
    sldIdLst = prs.slides._sldIdLst
    for idx in sorted(to_delete, reverse=True):
        sldIdLst.remove(sldIdLst[idx])




    # ───── 6) EXPENSE & RETURN SLIDE: Table 1 ────────────────────────────────────────
    # Locate the slide by its placeholder
    slide_expense_and_return = None
    for sl in prs.slides:
        for shape in sl.shapes:
            if shape.has_text_frame and "[Category] – Expense & Return" in shape.text_frame.text:
                slide_expense_and_return = sl
                break
        if slide_expense_and_return:
            break

    if not slide_expense_and_return:
        st.warning("Couldn't find the Expense and Return slide.")
    else:
        # Grab all tables, pick the first (top-left)
        tables = [sh for sh in slide_expense_and_return.shapes if sh.has_table]
        if not tables:
            st.warning("No tables found on Expense and Return slide.")
        else:
            tbl1 = tables[0].table
            tbl_xml = tbl1._tbl

            # Determine how many body rows to add
            existing = len(tbl1.rows) - 1  # header excluded
            needed = len(ear_df) - existing
            if needed > 0:
                base_tr = tbl_xml.tr_lst[1]
                for _ in range(needed):
                    tbl_xml.append(deepcopy(base_tr))

            
            # Fill each DataFrame row into the table
            for r_idx, df_row in enumerate(ear_df.itertuples(index=False), start=1):
                for c_idx, val in enumerate(df_row):
                    cell = tbl1.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]
            
                    # get or create the run
                    if para.runs:
                        run = para.runs[0]
                        run.text = val
                    else:
                        run = para.add_run()
                        run.text = val
            
                    # if this is the Net Expense Ratio column (column index 1):
                    if c_idx == 1:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                    # else: leave the template’s default styling for column 0
    
    # ───── 9c) Expense and Return Slide: Table 2 – Returns ─────────────────────────────

    # Locate the Expense and Return slide
    slide_expense_and_return = None
    for sl in prs.slides:
        for shape in sl.shapes:
            if shape.has_text_frame and "[Category] – Expense & Return" in shape.text_frame.text:
                slide_expense_and_return = sl
                break
        if slide_expense_and_return:
            break

    if not slide_expense_and_return:
        st.warning("Couldn't find the Expense and Return slide.")
    else:
        # Grab the second table (top‐right)
        tables = [sh for sh in slide_expense_and_return.shapes if sh.has_table]
        if len(tables) < 2:
            st.warning("Couldn't find Table 2 on Expense and Return slide.")
        else:
            table2 = tables[1].table
            tbl2_xml = table2._tbl

            # 1) Replace the header "Q_, 20__" with the actual quarter
            hdr_cell = table2.cell(0, 1)
            hdr_para = hdr_cell.text_frame.paragraphs[0]
            if hdr_para.runs:
                hdr_para.runs[0].text = date_label
            else:
                run = hdr_para.add_run()
                run.text = date_label
            # (font/size/color preserved from placeholder)

            # 2) Clone extra body rows if needed
            existing = len(table2.rows) - 1  # exclude header
            needed = len(df2) - existing
            if needed > 0:
                base_tr = tbl2_xml.tr_lst[1]
                for _ in range(needed):
                    tbl2_xml.append(deepcopy(base_tr))

            # 3) Fill each row: selected fund, proposals, then benchmark
            # Determine how many DF rows and which index is the benchmark
            df2 = st.session_state.get("ear_table2_data", pd.DataFrame())
            total_rows = len(df2)
        
            # 3) Fill each row: selected fund, proposals, then benchmark
            for r_idx, row in enumerate(df2.itertuples(index=False), start=1):
                is_benchmark = (r_idx == total_rows)
                for c_idx, val in enumerate(row):
                    cell = table2.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]
        
                    if para.runs:
                        run = para.runs[0]
                        run.text = val
                    else:
                        run = para.add_run()
                        run.text = val
        
                    if c_idx == 0:
                        # Investment Manager: preserve placeholder styling
                        continue
                    else:
                        # Other columns: Cambria 12pt black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        # Only bold the benchmark row
                        run.font.bold = is_benchmark

    # ───── 9d) Expense and Return Slide: Table 3 – Calendar Returns ────────────────────
    from copy import deepcopy
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    import pandas as pd

    # Pull in the DataFrame you saved in Step 15
    df3 = st.session_state.get("ear_table3_data", pd.DataFrame())
    if df3.empty:
        st.warning("No calendar returns data found for Table 3.")
    else:
        # Locate the 3rd table (under Tables 1 & 2)
        tables = [sh for sh in slide_expense_and_return.shapes if sh.has_table]
        if len(tables) < 3:
            st.warning("Couldn't find Table 3 on the Expense and Return slide.")
        else:
            table3 = tables[2].table
            tbl3_xml = table3._tbl

            # --- 1) Update year‐header placeholders ("20__") to actual years ---
            years = list(df3.columns[1:])  # skip "Investment Manager"
            for col_idx, year in enumerate(years, start=1):
                hdr_cell = table3.cell(0, col_idx)
                p = hdr_cell.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = year
                else:
                    run = p.add_run()
                    run.text = year
                # font/name/size/color/bold preserved from placeholder

            # --- 2) Ensure enough rows (header + body) ---
            existing = len(table3.rows) - 1  # exclude header row
            needed = len(df3) - existing
            if needed > 0:
                base_tr = tbl3_xml.tr_lst[1]  # first data row as template
                for _ in range(needed):
                    tbl3_xml.append(deepcopy(base_tr))

            # --- 3) Fill each row; bottom row (benchmark) bolds all data cols ---
            total = len(df3)
            for r_idx, row in enumerate(df3.itertuples(index=False), start=1):
                is_benchmark = (r_idx == total)
                for c_idx, val in enumerate(row):
                    cell = table3.cell(r_idx, c_idx)
                    p = cell.text_frame.paragraphs[0]
                    # get or create run
                    if p.runs:
                        run = p.runs[0]
                        run.text = val
                    else:
                        run = p.add_run()
                        run.text = val

                    if c_idx == 0:
                        # Investment Manager col → preserve placeholder style
                        continue
                    else:
                        # Calendar return cols → Cambria 11 pt black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        # Bold entire benchmark row
                        run.font.bold = is_benchmark

    # ───── Risk Adjusted Statistics Slide: Locate & Replace ─────────────────────────
    # (Make sure this is placed after you’ve defined `selected`, `facts`, and `fs_rec`)

    # Ensure fs_rec is defined:
    fs_rec = next((f for f in facts if f.get("Matched Fund Name") == selected), {})

    # 1) Find the slide
    risk_adjusted_stats = None
    for sl in prs.slides:
        for shape in sl.shapes:
            if shape.has_text_frame and "[Category] – Risk Adjusted Statistics" in shape.text_frame.text:
                risk_adjusted_stats = sl
                break
        if risk_adjusted_stats:
            break

    # 2) Replace placeholder
    if risk_adjusted_stats:
        actual_cat = fs_rec.get("Category", "")
        for shape in risk_adjusted_stats.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "[Category]" in run.text:
                        run.text = run.text.replace("[Category]", actual_cat)
    else:
        st.warning("Couldn't find the Risk Adjusted Statistics slide.")


    # ───── Risk Adjusted Statistics Slide: Table 1 – MPT Statistics Summary ─────────
    from copy import deepcopy
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    import pandas as pd

    # Pull in the DataFrame saved in Step 15
    df_raj = st.session_state.get("raj_table1_data", pd.DataFrame())
    if df_raj.empty:
        st.warning("No MPT Statistics data found for Table 1.")
    else:
        # Locate the first table on the Risk Adjusted Statistics slide
        ras_tables = [sh for sh in risk_adjusted_stats.shapes if sh.has_table]
        if not ras_tables:
            st.warning("No tables found on the Risk Adjusted Statistics slide.")
        else:
            table1 = ras_tables[0].table
            tbl_xml = table1._tbl

            # 1) Add extra rows if needed (clone row template at index 1)
            existing = len(table1.rows) - 1  # exclude header
            needed = len(df_raj) - existing
            if needed > 0:
                base_tr = tbl_xml.tr_lst[1]
                for _ in range(needed):
                    tbl_xml.append(deepcopy(base_tr))

            # 2) Fill each row: selected fund first, then proposals
            for r_idx, row in enumerate(df_raj.itertuples(index=False), start=1):
                for c_idx, val in enumerate(row):
                    cell = table1.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]

                    # get or create run
                    if para.runs:
                        run = para.runs[0]
                        run.text = val
                    else:
                        run = para.add_run()
                        run.text = val

                    if c_idx == 0:
                        # Investment Manager column: preserve placeholder styling
                        continue
                    else:
                        # Other columns: Cambria 11 pt black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # ───── Risk Adjusted Statistics Slide: Table 2 – Risk-Adjusted Returns / Peer Ranking % ─────
    from copy import deepcopy
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    import pandas as pd

    # pull in the DataFrame saved in Step 15
    df_raj2 = st.session_state.get("raj_table2_data", pd.DataFrame())

    if df_raj2.empty:
        st.warning("No Risk-Adjusted Returns / Peer Ranking data found for Table 2.")
    else:
        # assume `risk_adjusted_stats` slide was located earlier
        ras_tables = [sh for sh in risk_adjusted_stats.shapes if sh.has_table]
        if len(ras_tables) < 2:
            st.warning("Couldn't find Table 2 on the Risk Adjusted Statistics slide.")
        else:
            table2 = ras_tables[1].table
            tbl2_xml = table2._tbl

            # 1) add extra rows if needed (clone the first data row at idx 1)
            existing = len(table2.rows) - 1
            needed = len(df_raj2) - existing
            if needed > 0:
                base_tr = tbl2_xml.tr_lst[1]
                for _ in range(needed):
                    tbl2_xml.append(deepcopy(base_tr))

            # 2) fill each row: selected fund first, then proposals
            for r_idx, row in enumerate(df_raj2.itertuples(index=False), start=1):
                for c_idx, val in enumerate(row):
                    cell = table2.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]

                    # get or create run
                    if para.runs:
                        run = para.runs[0]
                        run.text = val
                    else:
                        run = para.add_run()
                        run.text = val

                    if c_idx == 0:
                        # Investment Manager column: preserve placeholder styling
                        continue
                    else:
                        # Other columns: Cambria 11pt black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)


    # ───── Qualitative Factors Slide: Locate & Replace Heading ─────────────────────────
    slide_qualitative_factors = None
    for sl in prs.slides:
        for shape in sl.shapes:
            if shape.has_text_frame and "[Category]– Qualitative Factors" in shape.text_frame.text:
                slide_qualitative_factors = sl
                break
        if slide_qualitative_factors:
            break

    if not slide_qualitative_factors:
        st.warning("Couldn't find the Qualitative Factors slide.")
    else:
        # Replace the [Category] token in the heading
        actual_cat = fs_rec.get("Category", "")
        for shape in slide_qualitative_factors.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "[Category]" in run.text:
                        run.text = run.text.replace("[Category]", actual_cat)

    # ───── Qualitative Factors Slide: Table 1 – Manager Tenure ─────────────────────────
    from copy import deepcopy
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    import pandas as pd

    # Pull in the DataFrame you saved in Step 15
    df_q1 = st.session_state.get("qualfact_table1_data", pd.DataFrame())
    if df_q1.empty:
        st.warning("No Manager Tenure data found for Table 1.")
    else:
        # 1) Locate the correct table by matching its headers
        table1 = None
        for shape in slide_qualitative_factors.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            h0 = tbl.cell(0, 0).text_frame.text.strip()
            h1 = tbl.cell(0, 1).text_frame.text.strip()
            if h0 == "Investment Manager" and h1 == "Manager Tenure":
                table1 = tbl
                tbl_xml = tbl._tbl
                break

        if table1 is None:
            st.warning("Couldn't find the Manager Tenure table.")
        else:
            # 2) Add rows if needed (keep header row intact)
            existing = len(table1.rows) - 1
            needed = len(df_q1) - existing
            if needed > 0:
                base_tr = tbl_xml.tr_lst[1]
                for _ in range(needed):
                    tbl_xml.append(deepcopy(base_tr))

            # 3) Fill each row: selected fund first, then proposals
            for r_idx, row in enumerate(df_q1.itertuples(index=False), start=1):
                for c_idx, val in enumerate(row):
                    cell = table1.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]
                    # get or create run
                    if para.runs:
                        run = para.runs[0]
                        run.text = val
                    else:
                        run = para.add_run()
                        run.text = val

                    if c_idx == 0:
                        # Investment Manager: preserve placeholder style
                        continue
                    else:
                        # Manager Tenure: Cambria 11pt black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # ───── Qualitative Factors Slide: Table 2 – Assets Under Management ────────────
    from copy import deepcopy
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    import pandas as pd

    df_q2 = st.session_state.get("qualfact_table2_data", pd.DataFrame())
    if df_q2.empty:
        st.warning("No Assets data found for Table 2.")
    else:
        # 1) Locate the correct table by header row (must have at least 3 columns)
        table2 = None
        for shape in slide_qualitative_factors.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            if len(tbl.columns) < 3:
                continue
            h0 = tbl.cell(0, 0).text_frame.text.strip()
            h1 = tbl.cell(0, 1).text_frame.text.strip()
            h2 = tbl.cell(0, 2).text_frame.text.strip()
            if (h0, h1, h2) == (
                "Investment Manager",
                "Assets Under Management",
                "Average Market Capitalization",
            ):
                table2 = tbl
                tbl_xml = tbl._tbl
                break

        if table2 is None:
            st.warning("Couldn't find the Assets Under Management table.")
        else:
            # 2) Add rows if there are more proposals than existing body rows
            existing = len(table2.rows) - 1  # exclude header
            needed = len(df_q2) - existing
            if needed > 0:
                # pick a base row index safely (first data row if present, else header)
                base_idx = 1 if len(table2.rows) > 1 else 0
                base_tr = tbl_xml.tr_lst[base_idx]
                for _ in range(needed):
                    tbl_xml.append(deepcopy(base_tr))

            # 3) Fill each row: selected fund first, then proposals
            for r_idx, record in enumerate(df_q2.itertuples(index=False), start=1):
                for c_idx, val in enumerate(record):
                    # guard against columns beyond shape
                    if c_idx >= len(table2.columns):
                        continue
                    cell = table2.cell(r_idx, c_idx)
                    para = cell.text_frame.paragraphs[0]

                    # replace or add run
                    if para.runs:
                        run = para.runs[0]
                        run.text = str(val)
                    else:
                        run = para.add_run()
                        run.text = str(val)

                    if c_idx == 0:
                        # Investment Manager: preserve the template’s styling
                        continue
                    else:
                        # Assets & Avg Market Cap: Cambria, 11 pt, black
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0, 0, 0)

    
    # ───── 6) Save & Download ───────────────────────────────────────────────────────────
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    st.success("PowerPoint Generated")
    st.download_button(
        label="Download Writeup PowerPoint",
        data=out,
        file_name=f"{selected} Writeup.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
# –– Cards ––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
def render_step16_and_16_5_cards(pdf):
    import streamlit as st
    import html
    import re

    def markdown_bold_to_html(text: str) -> str:
        escaped = html.escape(text)
        return re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', escaped)

    selected_fund = st.session_state.get("selected_fund", "—")
    # refresh bullets & overview
    step16_bullet_points(pdf)
    proposed_overview = step16_5_locate_proposed_factsheets_with_overview(pdf, context_lines=3, min_score=60)
    bullet_points = st.session_state.get("bullet_points", [])

    # Build selected fund card
    bullets_html = "".join(
        f"<div style='margin-bottom:6px; line-height:1.2; font-size:0.75rem;'>{markdown_bold_to_html(bp)}</div>"
        for bp in bullet_points
    ) or "<div style='font-size:0.75rem;'>No bullet points available.</div>"

    ips_status = ""
    ips_icon_table = st.session_state.get("ips_icon_table")
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        if not row.empty:
            ips_status = row.iloc[0].get("IPS Watch Status", "")

    status_display = {
        "NW": "No Watch",
        "IW": "Informal Watch",
        "FW": "Formal Watch"
    }.get(ips_status, "")

    status_badge = ""
    if status_display:
        if ips_status == "NW":
            badge_style = "background:#d6f5df; color:#217a3e;"
        elif ips_status == "IW":
            badge_style = "background:#fff3cd; color:#B87333;"
        elif ips_status == "FW":
            badge_style = "background:#f8d7da; color:#c30000;"
        else:
            badge_style = ""
        status_badge = (
            f"<span style='margin-left:8px; font-size:0.55rem; padding:4px 10px; border-radius:12px; "
            f"font-weight:600; display:inline-block; vertical-align:middle; {badge_style}'>"
            f"{html.escape(status_display)}</span>"
        )

    # Shared card style variables for consistency and smaller font
    CARD_BG = "linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%)"
    CARD_BORDER = "1.2px solid #b5d0eb"
    FONT_FAMILY = "system-ui,-apple-system,BlinkMacSystemFont,Segoe UI,Roboto,sans-serif"
    HEADING_COLOR = "#1f3f72"
    TEXT_COLOR = "#244369"
    LINE_HEIGHT = "1.3"

    selected_card = f"""
    <div style="
        background: {CARD_BG};
        color: {TEXT_COLOR};
        border-radius: 1.5rem;
        box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
        padding: 1.2rem 1.4rem;
        border: {CARD_BORDER};
        font-family: {FONT_FAMILY};
        line-height: {LINE_HEIGHT};
        max-width:100%;
        font-size: 0.75rem;
    ">
        <div style="display:flex; align-items:center; gap:6px; margin-bottom:6px;">
            <div style="font-size:0.85rem; font-weight:700; color:{HEADING_COLOR};">
                {html.escape(selected_fund)}
            </div>
            {status_badge}
        </div>
        {bullets_html}
    </div>
    """

    # Build proposed fund overview card(s)
    proposed_cards = ""
    if not proposed_overview:
        proposed_cards = f"""
        <div style="
            background: {CARD_BG};
            color: {TEXT_COLOR};
            border-radius: 1.5rem;
            box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
            padding: 1.2rem 1.4rem;
            border: {CARD_BORDER};
            font-family: {FONT_FAMILY};
            line-height: {LINE_HEIGHT};
            max-width:100%;
            font-size:0.75rem;
        ">
            <div style="font-size:0.85rem; font-weight:700; margin-bottom:6px; color:{HEADING_COLOR};">
                Proposed Fund Investment Overviews
            </div>
            <div style="font-size:0.75rem;">No proposed funds or overview data available.</div>
        </div>
        """
    else:
        for fund, info in proposed_overview.items():
            ticker = info.get("Ticker", "")
            name_label = f"{fund} ({ticker})" if ticker else fund
            paragraph = info.get("Overview Paragraph", "").strip()
            if not paragraph:
                snippet = "_No overview paragraph extracted._"
            else:
                snippet = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html.escape(paragraph))

            proposed_cards += f"""
            <div style="
                background: {CARD_BG};
                color: {TEXT_COLOR};
                border-radius: 1.5rem;
                box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
                padding: 1rem 1.3rem;
                border: {CARD_BORDER};
                font-family: {FONT_FAMILY};
                line-height: {LINE_HEIGHT};
                margin-bottom:1rem;
                max-width:100%;
                font-size:0.75rem;
            ">
                <div style="font-weight:700; font-size:0.85rem; margin-bottom:4px; color:{HEADING_COLOR};">{html.escape(name_label)}</div>
                <div style="font-size:0.7rem; line-height:1.25;">{snippet}</div>
            </div>
            """

    col1, col2 = st.columns(2, gap="large")
    with col1:
        st.markdown(selected_card, unsafe_allow_html=True)
    with col2:
        st.markdown(proposed_cards, unsafe_allow_html=True)


    # Spacer
    st.markdown("<div style='height:16px;'></div>", unsafe_allow_html=True)

# –– Main App –––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
def run():
    st.title("Writeup & Rec Test")
    uploaded = st.file_uploader("Upload MPI PDF to Generate Writup & Rec PPTX", type="pdf")
    if not uploaded:
        return

    with pdfplumber.open(uploaded) as pdf:
        # --- Initial metadata extraction ---
        first = pdf.pages[0].extract_text() or ""
        process_page1(first)
        show_report_summary()

        # --- TOC ---
        with st.expander("Table of Contents", expanded=False):
            toc_text = "".join((pdf.pages[i].extract_text() or "") for i in range(min(3, len(pdf.pages))))
            process_toc(toc_text)

        # --- All Fund Details ---
        with st.expander("All Fund Details", expanded=True):
            # IPS / Scorecard
            with st.expander("IPS Investment Screening", expanded=True):
                sp = st.session_state.get("scorecard_page")
                tot = st.session_state.get("total_options")
                pp = st.session_state.get("performance_page")
                factsheets_page = st.session_state.get("factsheets_page")
                if sp and tot is not None and pp:
                    step3_5_6_scorecard_and_ips(pdf, sp, pp, factsheets_page, tot)
                else:
                    st.error("Missing scorecard, performance page, or total options")

            # Factsheets
            with st.expander("Fund Factsheets", expanded=True):
                names = [b.get("Fund Name") for b in st.session_state.get("fund_blocks", [])]
                step6_process_factsheets(pdf, names)

            with st.expander("Fund Facts (sub-headings)", expanded=False):
                step12_process_fund_facts(pdf)

            with st.expander("Returns", expanded=False):
                step7_extract_returns(pdf)
                step8_calendar_returns(pdf)

            with st.expander("MPT Statistics Summary", expanded=False):
                step9_risk_analysis_3yr(pdf)
                step10_risk_analysis_5yr(pdf)
                step11_create_summary()

            with st.expander("Risk-Adjusted Returns", expanded=False):
                step13_process_risk_adjusted_returns(pdf)
                step14_extract_peer_risk_adjusted_return_rank(pdf)

        # --- Derive bullet context fields once (safe defaults) ---
        report_date = st.session_state.get("report_date", "")
        m = re.match(r"(\d)(?:st|nd|rd|th)\s+QTR,\s*(\d{4})", report_date or "")
        quarter = m.group(1) if m else ""
        year = m.group(2) if m else ""

        for itm in st.session_state.get("fund_performance_data", []):
            try:
                qtd = float(itm.get("QTD") or 0)
            except:
                qtd = 0.0
            try:
                bench_qtd = float(itm.get("Bench QTD") or 0)
            except:
                bench_qtd = 0.0
            itm["Perf Direction"] = "overperformed" if qtd >= bench_qtd else "underperformed"
            itm["Quarter"] = quarter
            itm["Year"] = year
            diff_bps = round((qtd - bench_qtd) * 100, 1)
            itm["QTD_bps_diff"] = str(diff_bps)
            fund_pct = f"{qtd:.2f}%"
            bench_pct = f"{bench_qtd:.2f}%"
            itm["QTD_pct_diff"] = f"{(qtd - bench_qtd):.2f}%"
            itm["QTD_vs"] = f"{fund_pct} vs. {bench_pct}"

        if "bullet_point_templates" not in st.session_state:
            st.session_state["bullet_point_templates"] = [
                "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], "
                "[Year] by [QTD_bps_diff] bps ([QTD_vs])."
            ]

        # --- Cards (proposed / watch / fail) ---
        extract_proposed_scorecard_blocks(pdf)
        fail_card_html, fail_css = get_ips_fail_card_html()
        proposed_card_html, proposed_css = get_proposed_fund_card_html()
        watch_summary_card_html, watch_summary_css = get_watch_summary_card_html()

        col1, col2 = st.columns(2, gap="large")
        with col1:
            if proposed_card_html:
                st.markdown(proposed_card_html, unsafe_allow_html=True)
            if watch_summary_card_html:
                st.markdown(watch_summary_card_html, unsafe_allow_html=True)
        with col2:
            if fail_card_html:
                st.markdown(fail_card_html, unsafe_allow_html=True)
        st.markdown(f"{fail_css}\n{proposed_css}\n{watch_summary_css}", unsafe_allow_html=True)

        # --- Single Fund Writeup ---
        with st.expander("Single Fund Write Up", expanded=False):
            step15_display_selected_fund()

        # --- Replaced: show Step 16 & 16.5 as side-by-side cards ---
        render_step16_and_16_5_cards(pdf)

        # --- Export to PowerPoint (always visible, clean UI) ---
        step17_export_to_ppt()



if __name__ == "__main__":
    run()
