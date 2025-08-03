import re
import streamlit as st
import pdfplumber
from calendar import month_name
import pandas as pd
from rapidfuzz import fuzz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from io import BytesIO
import yfinance as yf

#───Performance Table──────────────────────────────────────────────────────────────────

def extract_performance_table(pdf, performance_page, fund_names, end_page=None):
    import re
    from rapidfuzz import fuzz

    # Decide where to stop in the PDF
    end = end_page if end_page is not None else len(pdf.pages) + 1

    # 1. Get all lines from the section
    lines = []
    for pnum in range(performance_page - 1, end - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 2. Prepare regex
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    # 3. For each fund, try to pull numbers
    perf_data = []
    for name in fund_names:
        item = {"Fund Scorecard Name": name}
        tk = ""  # You’ll fill in ticker later from st.session_state["tickers"] or similar
        # a) Exact-ticker match: you’ll want to match this if you have tickers already
        idx = next(
            (i for i, ln in enumerate(lines)
             if name in ln),
            None
        )
        # b) Fuzzy-name fallback if not found
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                continue  # Can't find a match

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        clean += [None] * (8 - len(clean))  # pad

        # d) Map to columns
        item["QTD"] = clean[0]
        item["1Yr"] = clean[2]
        item["3Yr"] = clean[3]
        item["5Yr"] = clean[4]
        item["10Yr"] = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 3Yr, 5Yr from next lines
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]
        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None

        perf_data.append(item)
    return perf_data

#───Utility──────────────────────────────────────────────────────────────────

def extract_report_date(text):
    """
    Extracts and formats the report date from a block of text.
    Returns a string like '2nd QTR, 2024' for quarter-end dates,
    or 'As of March 12, 2024' for other dates.
    """
    dates = re.findall(r'(\d{1,2})/(\d{1,2})/(20\d{2})', text or "")
    for month, day, year in dates:
        m, d = int(month), int(day)
        # Quarter-end mapping
        quarter_map = {(3,31): "1st", (6,30): "2nd", (9,30): "3rd", (12,31): "4th"}
        if (m, d) in quarter_map:
            return f"{quarter_map[(m, d)]} QTR, {year}"
        # Fallback: month-day-year as human readable
        return f"As of {month_name[m]} {d}, {year}"
    return None

#───Page 1──────────────────────────────────────────────────────────────────

def process_page1(text):
    """
    Extracts 'report_date', 'total_options', 'prepared_for', and 'prepared_by' from text.
    Populates st.session_state with those keys.
    """
    # Extract report date
    report_date = extract_report_date(text)
    if report_date:
        st.session_state['report_date'] = report_date
    else:
        st.session_state['report_date'] = None

    # Extract total options
    m = re.search(r"Total Options:\s*(\d+)", text or "")
    st.session_state['total_options'] = int(m.group(1)) if m else None

    # Extract prepared for
    m = re.search(r"Prepared For:\s*\n(.*)", text or "")
    st.session_state['prepared_for'] = m.group(1).strip() if m else None

    # Extract prepared by, default if blank or says "mpi stylus"
    m = re.search(r"Prepared By:\s*(.*)", text or "")
    pb = m.group(1).strip() if m else ""
    if not pb or "mpi stylus" in pb.lower():
        pb = "Procyon Partners, LLC"
    st.session_state['prepared_by'] = pb

#───Info Card──────────────────────────────────────────────────────────────────

def show_report_summary():
    report_date    = st.session_state.get('report_date', 'N/A')
    total_options  = st.session_state.get('total_options', 'N/A')
    prepared_for   = st.session_state.get('prepared_for', 'N/A')
    prepared_by    = st.session_state.get('prepared_by', 'N/A')

    st.markdown(f"""
        <div style="
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.5rem;
            box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
            padding: 1.2rem 2.3rem 1.2rem 2.3rem;
            margin-bottom: 2rem;
            font-size: 1.08rem;
            border: 1.2px solid #b5d0eb;">
            <div style="color:#244369;">
                <b>Report Date:</b> {report_date} <br>
                <b>Total Options:</b> {total_options} <br>
                <b>Prepared For:</b> {prepared_for} <br>
                <b>Prepared By:</b> {prepared_by}
            </div>
        </div>
    """, unsafe_allow_html=True)

#───Step 2: Table of Contents Extraction──────────────────────────────────────────────────────────

def process_toc(text):
    perf = re.search(r"Fund Performance[^\d]*(\d{1,3})", text or "")
    cy   = re.search(r"Fund Performance: Calendar Year\s+(\d{1,3})", text or "")
    r3yr = re.search(r"Risk Analysis: MPT Statistics \(3Yr\)\s+(\d{1,3})", text or "")
    r5yr = re.search(r"Risk Analysis: MPT Statistics \(5Yr\)\s+(\d{1,3})", text or "")

    sc            = re.search(r"Fund Scorecard\s+(\d{1,3})", text or "")
    sc_prop       = re.search(r"Fund Scorecard:\s*Proposed Funds\s+(\d{1,3})", text or "")

    fs            = re.search(r"Fund Factsheets\s+(\d{1,3})", text or "")
    fs_prop       = re.search(r"Fund Factsheets:\s*Proposed Funds\s+(\d{1,3})", text or "")

    perf_page     = int(perf.group(1)) if perf else None
    cy_page       = int(cy.group(1)) if cy else None
    r3yr_page     = int(r3yr.group(1)) if r3yr else None
    r5yr_page     = int(r5yr.group(1)) if r5yr else None
    sc_page       = int(sc.group(1)) if sc else None
    sc_prop_page  = int(sc_prop.group(1)) if sc_prop else None
    fs_page       = int(fs.group(1)) if fs else None
    fs_prop_page  = int(fs_prop.group(1)) if fs_prop else None

    # Store in session state for downstream use
    st.session_state['performance_page'] = perf_page
    st.session_state['calendar_year_page'] = cy_page
    st.session_state['r3yr_page'] = r3yr_page
    st.session_state['r5yr_page'] = r5yr_page
    st.session_state['scorecard_page'] = sc_page
    st.session_state['scorecard_proposed_page'] = sc_prop_page
    st.session_state['factsheets_page'] = fs_page
    st.session_state['factsheets_proposed_page'] = fs_prop_page

#───IPS Invesment Screening──────────────────────────────────────────────────────────────────
import streamlit as st
import re
import pdfplumber
import pandas as pd
import yfinance as yf

def infer_fund_type_guess(ticker):
    """Infer 'Active' or 'Passive' based on Yahoo Finance info (name and summary)."""
    try:
        if not ticker:
            return ""
        info = yf.Ticker(ticker).info
        name = (info.get("longName") or info.get("shortName") or "").lower()
        summary = (info.get("longBusinessSummary") or "").lower()
        if "index" in name or "index" in summary:
            return "Passive"
        if "track" in summary and "index" in summary:
            return "Passive"
        if "actively managed" in summary or "actively-managed" in summary:
            return "Active"
        if "outperform" in summary or "manager selects" in summary:
            return "Active"
        return ""
    except Exception:
        return ""

def extract_scorecard_blocks(pdf, scorecard_page):
    metric_labels = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)"
    ]
    pages, fund_blocks, fund_name, metrics = [], [], None, []
    for p in pdf.pages[scorecard_page-1:]:
        pages.append(p.extract_text() or "")
    lines = "\n".join(pages).splitlines()
    for line in lines:
        if not any(metric in line for metric in metric_labels) and line.strip():
            if fund_name and metrics:
                fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
            fund_name = re.sub(r"Fund (Meets Watchlist Criteria|has been placed on watchlist for not meeting .* out of 14 criteria)", "", line.strip()).strip()
            metrics = []
        for metric in metric_labels:
            if metric in line:
                m = re.match(r"^(.*?)\s+(Pass|Review|Fail)\s*(.*)", line.strip())
                if m:
                    metric_name, status, info = m.groups()
                    metrics.append({"Metric": metric_name, "Status": status, "Info": info.strip()})
    if fund_name and metrics:
        fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
    return fund_blocks

def extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page=None):
    import re
    from rapidfuzz import fuzz

    def normalize_name(name):
        # strip the watchlist suffix and punctuation, lowercase
        cleaned = re.sub(
            r"has been placed on watchlist for not meeting .*? criteria",
            "",
            name,
            flags=re.IGNORECASE
        )
        cleaned = re.sub(r"[^A-Za-z0-9 ]+", "", cleaned)  # remove punctuation
        return cleaned.strip().lower()

    end_page = factsheets_page - 1 if factsheets_page else len(pdf.pages)
    all_lines = []
    for p in pdf.pages[performance_page - 1:end_page]:
        txt = p.extract_text() or ""
        all_lines.extend([ln.strip() for ln in txt.splitlines() if ln.strip()])

    # Step 1: Collect candidate (raw_name, ticker) pairs from lines with uppercase tickers length 2-5
    candidate_pairs = []  # list of tuples (normalized_raw_name, ticker, raw_name_original)
    ticker_rx = re.compile(r"\b([A-Z]{2,5})\b")
    for ln in all_lines:
        matches = ticker_rx.findall(ln)
        if not matches:
            continue
        for ticker in set(matches):
            parts = ln.rsplit(ticker, 1)
            if len(parts) >= 1:
                raw_name = parts[0].strip()
                if not raw_name:
                    continue
                norm_raw = normalize_name(raw_name)
                if not norm_raw:
                    continue
                candidate_pairs.append((norm_raw, ticker, raw_name))

    # Step 2: For each fund_name, find best candidate fuzzy match (one-to-one)
    assigned = {}
    used_tickers = set()

    # Precompute normalized expected names
    norm_expected_list = [
        (name, normalize_name(name))
        for name in fund_names
    ]

    for fund_name, norm_expected in norm_expected_list:
        best = (None, None, 0)  # (ticker, raw_name_original, score)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            score = fuzz.token_sort_ratio(norm_expected, norm_raw)
            if score > best[2]:
                best = (ticker, raw_name, score)
        if best[2] >= 70:  # threshold, adjust if needed
            assigned[fund_name] = best[0]
            used_tickers.add(best[0])
        else:
            assigned[fund_name] = ""

    # Step 3: Fallback for unmatched names: looser heuristic and line-based fallback
    for name in fund_names:
        if assigned.get(name):
            continue
        norm_expected = normalize_name(name)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            if norm_expected in norm_raw or norm_raw in norm_expected:
                assigned[name] = ticker
                used_tickers.add(ticker)
                break
        if not assigned.get(name):
            for ln in all_lines:
                if name.lower() in ln.lower():
                    m = re.search(r"\b([A-Z]{2,5})\b", ln)
                    if m:
                        assigned[name] = m.group(1)
                        break

    # Final cleanup: ensure non-None strings
    return {k: (v if v else "") for k, v in assigned.items()}

def scorecard_to_ips(fund_blocks, fund_types, tickers):
    # Maps for converting scorecard to IPS criteria (from your logic)
    metrics_order = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)",
    ]
    active_map  = [0,1,3,6,10,2,4,7,11,5,None]
    passive_map = [0,8,3,6,12,9,4,7,13,5,None]
    ips_labels = [f"IPS Investment Criteria {i+1}" for i in range(11)]
    ips_results, raw_results = [], []
    for fund in fund_blocks:
        fund_name = fund["Fund Name"]
        fund_type = fund_types.get(fund_name, "Passive" if "index" in fund_name.lower() else "Active")
        metrics = fund["Metrics"]
        scorecard_status = [next((m["Status"] for m in metrics if m["Metric"] == label), None) for label in metrics_order]
        idx_map = passive_map if fund_type == "Passive" else active_map
        ips_status = [scorecard_status[m_idx] if m_idx is not None else "Pass" for m_idx in idx_map]
        review_fail = sum(1 for status in ips_status if status in ["Review","Fail"])
        watch_status = "FW" if review_fail >= 6 else "IW" if review_fail >= 5 else "NW"
        def iconify(status): return "✔" if status == "Pass" else "✗" if status in ("Review", "Fail") else ""
        row = {
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: iconify(ips_status[i]) for i in range(11)},
            "IPS Watch Status": watch_status,
        }
        ips_results.append(row)
        raw_results.append({
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: ips_status[i] for i in range(11)},
            "IPS Watch Status": watch_status,
        })
    return pd.DataFrame(ips_results), pd.DataFrame(raw_results)

def watch_status_color(val):
    if val == "FW":
        return "background-color:#f8d7da; color:#c30000; font-weight:600;"
    if val == "IW":
        return "background-color:#fff3cd; color:#B87333; font-weight:600;"
    if val == "NW":
        return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
    return ""

def step3_5_6_scorecard_and_ips(pdf, scorecard_page, performance_page, factsheets_page, total_options):
    import pandas as pd
    import streamlit as st

    # --- 1. Extract scorecard blocks ---
    fund_blocks = extract_scorecard_blocks(pdf, scorecard_page)
    fund_names = [fund["Fund Name"] for fund in fund_blocks]
    if not fund_blocks:
        st.error("Could not extract fund scorecard blocks. Check the PDF and page number.")
        return

    # --- 2. Extract tickers ---
    tickers = extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page)

    # --- Prepare inferred fund type guesses/defaults ---
    inferred_guesses = []
    for name in fund_names:
        guess = ""
        if tickers.get(name):
            guess = infer_fund_type_guess(tickers.get(name, "")) or ""
        inferred_guesses.append("Passive" if guess.lower() == "passive" else ("Passive" if "index" in name.lower() else "Active"))

    # Build base df for editor
    df_types_base = pd.DataFrame({
        "Fund Name":       fund_names,
        "Ticker":          [tickers.get(n, "") for n in fund_names],
        "Inferred Type":   inferred_guesses,
        "Fund Type":       inferred_guesses,  # starts same as inferred
    })

    # --- Toggleable editor display ---
    if "show_edit_fund_type" not in st.session_state:
        st.session_state["show_edit_fund_type"] = False
    toggle_label = "Hide Fund Type Editor" if st.session_state["show_edit_fund_type"] else "Edit Fund Type"
    if st.button(toggle_label, key="toggle_edit_fund_type"):
        st.session_state["show_edit_fund_type"] = not st.session_state["show_edit_fund_type"]

    # Decide fund_types mapping
    if st.session_state["show_edit_fund_type"]:
        st.markdown("### Fund Type Overrides")
        st.caption("Edit any fund type here; once you modify at least one, your edits take precedence over inferred types.")
        edited_types = st.data_editor(
            df_types_base,
            column_config={
                "Fund Type": st.column_config.SelectboxColumn("Fund Type", options=["Active", "Passive"]),
            },
            disabled=["Fund Name", "Ticker", "Inferred Type"],
            hide_index=True,
            key="data_editor_fundtype_ips",
            use_container_width=True,
        )
        # Determine whether any manual edit occurred (diff between Fund Type and Inferred Type)
        manual_override = any(
            row["Fund Type"] != row["Inferred Type"]
            for _, row in edited_types.iterrows()
        )
        if manual_override:
            fund_types = {row["Fund Name"]: row["Fund Type"] for _, row in edited_types.iterrows()}
        else:
            # none edited: use inferred
            fund_types = {row["Fund Name"]: row["Inferred Type"] for _, row in edited_types.iterrows()}
    else:
        # editor hidden: use inferred guesses
        fund_types = {name: inferred_guesses[i] for i, name in enumerate(fund_names)}

    # --- 4. IPS conversion ---
    df_icon, df_raw = scorecard_to_ips(fund_blocks, fund_types, tickers)

    # --- IPS Results ---
    st.subheader("IPS Screening Results")
    st.markdown(
        '<div style="display:flex; gap:1rem; margin-bottom:0.5rem;">'
        '<div style="padding:4px 10px; background:#d6f5df; border-radius:4px; font-weight:600;">NW: No Watch</div>'
        '<div style="padding:4px 10px; background:#fff3cd; border-radius:4px; font-weight:600;">IW: Informal Watch</div>'
        '<div style="padding:4px 10px; background:#f8d7da; border-radius:4px; font-weight:600;">FW: Formal Watch</div>'
        '</div>', unsafe_allow_html=True
    )

    if df_icon.empty:
        st.info("No IPS screening data available.")
    else:
        # Compact numbered criteria
        display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
        display_df = df_icon.rename(columns=display_columns)

        def iconify(s):
            if s == "Pass":
                return "✔"
            if s in ("Review", "Fail"):
                return "✗"
            return ""

        compact_df = display_df.copy()
        for i in range(1, 12):
            orig = f"IPS Investment Criteria {i}"
            if orig in compact_df.columns:
                compact_df[str(i)] = compact_df[orig].apply(iconify)
                compact_df.drop(columns=[orig], inplace=True)
        cols_order = ["Fund Name", "Ticker", "Fund Type"] + [str(i) for i in range(1, 12)] + ["IPS Watch Status"]
        compact_df = compact_df[[c for c in cols_order if c in compact_df.columns]]

        def watch_style(val):
            if val == "NW":
                return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
            if val == "IW":
                return "background-color:#fff3cd; color:#B87333; font-weight:600;"
            if val == "FW":
                return "background-color:#f8d7da; color:#c30000; font-weight:600;"
            return ""

        styled = compact_df.style.applymap(watch_style, subset=["IPS Watch Status"])
        st.dataframe(styled, use_container_width=True, hide_index=True)

        # --- Summary badges (bottom) ---
        def summarize_watch(df):
            counts = df["IPS Watch Status"].value_counts().to_dict()
            return {
                "No Watch": counts.get("NW", 0),
                "Informal Watch": counts.get("IW", 0),
                "Formal Watch": counts.get("FW", 0),
            }

        summary = summarize_watch(df_icon)
        st.markdown("---")
        st.markdown("### Watch Summary")
        b1, b2, b3 = st.columns(3, gap="small")
        with b1:
            st.metric("No Watch", summary["No Watch"])
        with b2:
            st.metric("Informal Watch", summary["Informal Watch"])
        with b3:
            st.metric("Formal Watch", summary["Formal Watch"])

    # --- Persist state downstream ---
    st.session_state["fund_blocks"] = fund_blocks
    st.session_state["fund_types"] = fund_types
    st.session_state["fund_tickers"] = tickers
    st.session_state["ips_icon_table"] = df_icon
    st.session_state["ips_raw_table"] = df_raw

    # --- Performance extraction ---
    perf_data = extract_performance_table(
        pdf,
        performance_page,
        fund_names,
        factsheets_page
    )
    for itm in perf_data:
        itm["Ticker"] = tickers.get(itm["Fund Scorecard Name"], "")

    st.session_state["fund_performance_data"] = perf_data
    st.session_state["tickers"] = tickers  # legacy compatibility


#───Step 6:Factsheets Pages──────────────────────────────────────────────────────────────────

def step6_process_factsheets(pdf, fund_names, suppress_output=True):
    # If you ever want UI, set suppress_output=False when calling

    factsheet_start = st.session_state.get("factsheets_page")
    total_declared = st.session_state.get("total_options")
    performance_data = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in st.session_state.get("tickers", {}).items()
    ]

    if not factsheet_start:
        if not suppress_output:
            st.error("❌ 'Fund Factsheets' page number not found in TOC.")
        return

    matched_factsheets = []
    for i in range(factsheet_start - 1, len(pdf.pages)):
        page = pdf.pages[i]
        words = page.extract_words(use_text_flow=True)
        header_words = [w['text'] for w in words if w['top'] < 100]
        first_line = " ".join(header_words).strip()

        if not first_line or "Benchmark:" not in first_line or "Expense Ratio:" not in first_line:
            continue

        ticker_match = re.search(r"\b([A-Z]{5})\b", first_line)
        ticker = ticker_match.group(1) if ticker_match else ""
        fund_name_raw = first_line.split(ticker)[0].strip() if ticker else first_line

        best_score = 0
        matched_name = matched_ticker = ""
        for item in performance_data:
            ref = f"{item['Fund Scorecard Name']} {item['Ticker']}".strip()
            score = fuzz.token_sort_ratio(f"{fund_name_raw} {ticker}".lower(), ref.lower())
            if score > best_score:
                best_score, matched_name, matched_ticker = score, item['Fund Scorecard Name'], item['Ticker']

        def extract_field(label, text, stop=None):
            try:
                start = text.index(label) + len(label)
                rest = text[start:]
                if stop and stop in rest:
                    return rest[:rest.index(stop)].strip()
                return rest.split()[0]
            except Exception:
                return ""

        benchmark = extract_field("Benchmark:", first_line, "Category:")
        category  = extract_field("Category:", first_line, "Net Assets:")
        net_assets= extract_field("Net Assets:", first_line, "Manager Name:")
        manager   = extract_field("Manager Name:", first_line, "Avg. Market Cap:")
        avg_cap   = extract_field("Avg. Market Cap:", first_line, "Expense Ratio:")
        expense   = extract_field("Expense Ratio:", first_line)

        matched_factsheets.append({
            "Page #": i + 1,
            "Parsed Fund Name": fund_name_raw,
            "Parsed Ticker": ticker,
            "Matched Fund Name": matched_name,
            "Matched Ticker": matched_ticker,
            "Benchmark": benchmark,
            "Category": category,
            "Net Assets": net_assets,
            "Manager Name": manager,
            "Avg. Market Cap": avg_cap,
            "Expense Ratio": expense,
            "Match Score": best_score,
            "Matched": "✅" if best_score > 20 else "❌"
        })

    df_facts = pd.DataFrame(matched_factsheets)
    st.session_state['fund_factsheets_data'] = matched_factsheets

    # Hide UI output unless suppress_output is False
    if not suppress_output:
        display_df = df_facts[[
            "Matched Fund Name", "Matched Ticker", "Benchmark", "Category",
            "Net Assets", "Manager Name", "Avg. Market Cap", "Expense Ratio", "Matched"
        ]].rename(columns={"Matched Fund Name": "Fund Name", "Matched Ticker": "Ticker"})

        st.dataframe(display_df, use_container_width=True)

        matched_count = sum(1 for r in matched_factsheets if r["Matched"] == "✅")
        st.write(f"Matched {matched_count} of {len(matched_factsheets)} factsheet pages.")
        if matched_count == total_declared:
            st.success(f"All {matched_count} funds matched the declared Total Options from Page 1.")
        else:
            st.error(f"Mismatch: Page 1 declared {total_declared}, but only matched {matched_count}.")

#───Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense Ratio & Bench QTD──────────────────────────────────────────────────────────────────

def step7_extract_returns(pdf):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    # 1) Where to scan
    perf_page = st.session_state.get("performance_page")
    end_page  = st.session_state.get("calendar_year_page") or (len(pdf.pages) + 1)
    perf_data = st.session_state.get("fund_performance_data", [])
    if perf_page is None or not perf_data:
        st.error("❌ Run Step 5 first to populate performance data.")
        return

    # 2) Prep output slots
    fields = [
        "QTD", "1Yr", "3Yr", "5Yr", "10Yr", "Net Expense Ratio",
        "Bench QTD", "Bench 1Yr", "Bench 3Yr", "Bench 5Yr", "Bench 10Yr"
    ]
    for itm in perf_data:
        for f in fields:
            itm.setdefault(f, None)

    # 3) Gather every nonblank line in the Performance section
    lines = []
    for pnum in range(perf_page - 1, end_page - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 4) Regex to pull decimal tokens (with optional % and parentheses)
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    matched = 0
    for item in perf_data:
        name = item["Fund Scorecard Name"]
        tk   = item["Ticker"].upper().strip()

        # a) Exact-ticker match
        idx = next(
            (i for i, ln in enumerate(lines)
             if re.search(rf"\b{re.escape(tk)}\b", ln)),
            None
        )
        # b) Fuzzy-name fallback
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                st.warning(f"⚠️ {name} ({tk}): no match found.")
                continue

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        if len(clean) < 8:
            clean += [None] * (8 - len(clean))

        # d) Map fund returns & net expense
        item["QTD"]               = clean[0]
        item["1Yr"]               = clean[2]
        item["3Yr"]               = clean[3]
        item["5Yr"]               = clean[4]
        item["10Yr"]              = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 1Yr, 3Yr, 5Yr, and 10Yr from the very next line(s)
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 5 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]

        item["Bench QTD"]  = bench_clean[0] if len(bench_clean) > 0 else None
        item["Bench 1Yr"] = bench_clean[1] if len(bench_clean) > 1 else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None
        item["Bench 10Yr"] = bench_clean[5] if len(bench_clean) > 5 else None

        matched += 1

    # 5) Save & display
    st.session_state["fund_performance_data"] = perf_data
    df = pd.DataFrame(perf_data)
    # Hide the per-fund warnings and overall success message, but still alert if something is off
    
    expected_count = len(perf_data)
    if matched < expected_count:
        st.error(f"❌ Only matched {matched} of {expected_count} funds with return data. Check your PDF or extraction logic.")
    # (Do NOT display per-fund warnings or success)

    st.dataframe(
        df[["Fund Scorecard Name", "Ticker"] + fields],
        use_container_width=True
    )


#───Step 8 Calendar Year Returns (funds + benchmarks──────────────────────────────────────────────────────────────────

def step8_calendar_returns(pdf):
    import re, streamlit as st, pandas as pd

    # 1) Figure out section bounds
    cy_page  = st.session_state.get("calendar_year_page")
    end_page = st.session_state.get("r3yr_page", len(pdf.pages) + 1)
    if cy_page is None:
        st.error("❌ 'Fund Performance: Calendar Year' not found in TOC.")
        return

    # 2) Pull every line from that section
    all_lines = []
    for p in pdf.pages[cy_page-1 : end_page-1]:
        all_lines.extend((p.extract_text() or "").splitlines())

    # 3) Identify header & years
    header = next((ln for ln in all_lines if "Ticker" in ln and re.search(r"20\d{2}", ln)), None)
    if not header:
        st.error("❌ Couldn’t find header row with 'Ticker' + year.")
        return
    years = re.findall(r"\b20\d{2}\b", header)
    num_rx = re.compile(r"-?\d+\.\d+%?")

    # — A) Funds themselves —
    fund_map     = st.session_state.get("tickers", {})
    fund_records = []
    for name, tk in fund_map.items():
        ticker = (tk or "").upper()
        idx    = next((i for i, ln in enumerate(all_lines) if ticker in ln.split()), None)
        raw    = num_rx.findall(all_lines[idx-1]) if idx not in (None, 0) else []
        vals   = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec    = {"Name": name, "Ticker": ticker}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        fund_records.append(rec)

    df_fund = pd.DataFrame(fund_records)
    if not df_fund.empty:
        st.markdown("**Fund Calendar‑Year Returns**")
        st.dataframe(df_fund[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["step8_returns"] = fund_records

    # — B) Benchmarks matched back to each fund’s ticker —
    facts         = st.session_state.get("fund_factsheets_data", [])
    bench_records = []
    for f in facts:
        bench_name = f.get("Benchmark", "").strip()
        fund_tkr   = f.get("Matched Ticker", "")
        if not bench_name:
            continue

        # find the first line containing the benchmark name
        idx = next((i for i, ln in enumerate(all_lines) if bench_name in ln), None)
        if idx is None:
            continue
        raw  = num_rx.findall(all_lines[idx])
        vals = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec  = {"Name": bench_name, "Ticker": fund_tkr}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        bench_records.append(rec)

    df_bench = pd.DataFrame(bench_records)
    if not df_bench.empty:
        st.markdown("**Benchmark Calendar‑Year Returns**")
        st.dataframe(df_bench[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["benchmark_calendar_year_returns"] = bench_records
    else:
        st.warning("No benchmark returns extracted.")

#───Step 9: 3‑Yr Risk Analysis – Match & Extract MPT Stats (hidden matching)──────────────────────────────────────────────────────────────────

def step9_risk_analysis_3yr(pdf):
    import re, streamlit as st, pandas as pd
    from rapidfuzz import fuzz

    # 1) Get your fund→ticker map
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (3Yr)” page
    start_page = st.session_state.get("r3yr_page")
    if not start_page:
        st.error("❌ ‘Risk Analysis: MPT Statistics (3Yr)’ page not found; run Step 2 first.")
        return

    # 3) Scan forward until you’ve seen each ticker (no display)
    locs = {}
    for pnum in range(start_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for fname, tk in fund_map.items():
                if fname in locs: 
                    continue
                if tk.upper() in tokens:
                    locs[fname] = {"page": pnum, "line": li}
        if len(locs) == len(fund_map):
            break

    # 4) Extract the first four numeric MPT stats from that same line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, info in locs.items():
        page = pdf.pages[info["page"]-1]
        lines = (page.extract_text() or "").splitlines()
        line = lines[info["line"]]
        nums = num_rx.findall(line)
        nums += [None] * (4 - len(nums))
        alpha, beta, up, down = nums[:4]
        results.append({
            "Fund Name":               name,
            "Ticker":                  fund_map[name].upper(),
            "3 Year Alpha":            alpha,
            "3 Year Beta":             beta,
            "3 Year Upside Capture":   up,
            "3 Year Downside Capture": down
        })

    # 5) Display final table only
    st.session_state["step9_mpt_stats"] = results

#───Step 10: Risk Analysis (5Yr) – Match & Extract MPT Statistics──────────────────────────────────────────────────────────────────

def step10_risk_analysis_5yr(pdf):
    import re, streamlit as st, pandas as pd

    # 1) Your fund→ticker map from Step 5
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (5Yr)” section
    section_page = next(
        (i for i, pg in enumerate(pdf.pages, start=1)
         if "Risk Analysis: MPT Statistics (5Yr)" in (pg.extract_text() or "")),
        None
    )
    if section_page is None:
        st.error("❌ Could not find ‘Risk Analysis: MPT Statistics (5Yr)’ section.")
        return

    # 3) Under‑the‑hood: scan pages until each ticker is located
    locs = {}
    total = len(fund_map)
    for pnum in range(section_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for name, tk in fund_map.items():
                if name in locs:
                    continue
                if tk.upper() in tokens:
                    locs[name] = {"page": pnum, "line": li}
        if len(locs) == total:
            break

    # 4) Wrap‑aware extraction of the first four floats after each ticker line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, tk in fund_map.items():
        info = locs.get(name)
        vals = [None] * 4
        if info:
            page = pdf.pages[info["page"] - 1]
            text_lines = (page.extract_text() or "").splitlines()
            idx = info["line"]
            nums = []
            # look on the line of the ticker and up to the next 2 lines
            for j in range(idx, min(idx + 3, len(text_lines))):
                nums += num_rx.findall(text_lines[j])
                if len(nums) >= 4:
                    break
            nums += [None] * (4 - len(nums))
            vals = nums[:4]
        else:
            st.warning(f"⚠️ {name} ({tk.upper()}): not found after page {section_page}.")

        alpha5, beta5, up5, down5 = vals
        results.append({
            "Fund Name":               name,
            "Ticker":                  tk.upper(),
            "5 Year Alpha":            alpha5,
            "5 Year Beta":             beta5,
            "5 Year Upside Capture":   up5,
            "5 Year Downside Capture": down5,
        })

    # 5) Save & display only the consolidated table
    st.session_state["step10_mpt_stats"] = results

#───Step 11: Combined MPT Statistics Summary──────────────────────────────────────────────────────────────────

def step11_create_summary(pdf=None):
    import pandas as pd
    import streamlit as st

    # 1) Load your 3‑Yr and 5‑Yr stats from session state
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    if not mpt3 or not mpt5:
        st.error("❌ Missing MPT stats. Run Steps 9 & 10 first.")
        return

    # 2) Build DataFrames
    df3 = pd.DataFrame(mpt3)  # contains "3 Year Alpha", "3 Year Beta", etc.
    df5 = pd.DataFrame(mpt5)  # contains "5 Year Alpha", "5 Year Beta", etc.

    # 3) Merge on Fund Name & Ticker
    df = pd.merge(
        df3,
        df5,
        on=["Fund Name", "Ticker"],
        how="outer",
        suffixes=("_3yr", "_5yr")
    )

    # 4) Build the Investment Manager column
    df.insert(0, "Investment Manager", df["Fund Name"] + " (" + df["Ticker"] + ")")

    # 5) Select & order the columns
    df = df[[
        "Investment Manager",
        "3 Year Alpha",
        "5 Year Alpha",
        "3 Year Beta",
        "5 Year Beta",
        "3 Year Upside Capture",
        "3 Year Downside Capture",
        "5 Year Upside Capture",
        "5 Year Downside Capture"
    ]]

    # 6) Display
    st.session_state["step11_summary"] = df.to_dict("records")
    st.dataframe(df)

#───Step 12: Extract “FUND FACTS” & Its Table Details in One Go──────────────────────────────────────────────────────────────────

def step12_process_fund_facts(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # the exact labels and the order you want them in the table
    labels = [
        "Manager Tenure Yrs.",
        "Expense Ratio",
        "Expense Ratio Rank",
        "Total Number of Holdings",
        "Turnover Ratio"
    ]

    records = []
    # scan each factsheet page
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = pdf.pages[pnum-1].extract_text().splitlines()

        for idx, line in enumerate(lines):
            if line.lstrip().upper().startswith("FUND FACTS"):
                # grab the next 8 lines (should contain your 5 labels)
                snippet = lines[idx+1 : idx+1+8]
                rec = {"Fund Name": fund_name, "Ticker": ticker}
                for lab in labels:
                    val = None
                    for ln in snippet:
                        norm = " ".join(ln.strip().split())
                        if norm.startswith(lab):
                            rest = norm[len(lab):].strip(" :\t")
                            m = re.match(r"(-?\d+\.\d+)", rest)
                            val = m.group(1) if m else (rest.split()[0] if rest else None)
                            break
                    rec[lab] = val
                records.append(rec)
                break  # move on to the next page once Fund Facts is processed

    if not records:
        st.warning("No Fund Facts tables found.")
        return

    # save & show
    st.session_state["step12_fund_facts_table"] = records

#───Step 13: Extract Risk‑Adjusted Returns Metrics──────────────────────────────────────────────────────────────────

def step13_process_risk_adjusted_returns(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # which metrics to pull
    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    num_rx  = re.compile(r"-?\d+\.\d+")

    records = []
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()

        # find the heading
        for idx, line in enumerate(lines):
            norm = " ".join(line.strip().split()).upper()
            if norm.startswith("RISK-ADJUSTED RETURNS"):
                snippet = lines[idx+1 : idx+1+6]  # grab next few lines
                rec = {"Fund Name": fund_name, "Ticker": ticker}

                for metric in metrics:
                    # find the snippet line for this metric
                    text_line = next(
                        ( " ".join(ln.strip().split())
                          for ln in snippet
                          if ln.strip().upper().startswith(metric.upper()) ),
                        None
                    ) or ""
                    # extract up to 4 numbers
                    nums = num_rx.findall(text_line)
                    nums += [None] * (4 - len(nums))

                    # assign into rec
                    rec[f"{metric} 1Yr"]  = nums[0]
                    rec[f"{metric} 3Yr"]  = nums[1]
                    rec[f"{metric} 5Yr"]  = nums[2]
                    rec[f"{metric} 10Yr"] = nums[3]

                records.append(rec)
                break  # done with this page

    if not records:
        st.warning("No 'RISK‑ADJUSTED RETURNS' tables found.")
        return

    # save & show
    st.session_state["step13_risk_adjusted_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#───Step 14: Peer Risk-Adjusted Return Rank──────────────────────────────────────────────────────────────────

def step14_extract_peer_risk_adjusted_return_rank(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.write("Peer Rank")

    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    page_map = {
        f["Page #"]:(f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    records = []

    for pnum, (fund, ticker) in page_map.items():
        page = pdf.pages[pnum-1]
        text = page.extract_text() or ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        # 1) locate the Risk-Adjusted Returns header
        try:
            risk_idx = next(i for i, ln in enumerate(lines)
                            if "RISK-ADJUSTED RETURNS" in ln.upper())
        except StopIteration:
            st.warning(f"⚠️ {fund} ({ticker}): Risk-Adjusted header not found.")
            continue

        # 2) find all “1 Yr 3 Yrs 5 Yrs 10 Yrs” lines after that
        header_idxs = [i for i, ln in enumerate(lines)
                       if re.match(r"1\s*Yr", ln)]
        peer_header_idxs = [i for i in header_idxs if i > risk_idx]

        if not peer_header_idxs:
            st.warning(f"⚠️ {fund} ({ticker}): peer header not found.")
            continue

        # take the *second* header occurrence (first is Risk-Adjusted, next is Peer)
        peer_hdr = peer_header_idxs[0] if len(peer_header_idxs)==1 else peer_header_idxs[1]

        rec = {"Fund Name": fund, "Ticker": ticker}

        # 3) read the three lines immediately below that header
        for offset, metric in enumerate(metrics, start=1):
            if peer_hdr + offset < len(lines):
                parts = lines[peer_hdr + offset].split()
                # parts[0:2] = metric name words, parts[2:6] = the four integer ranks
                vals = parts[2:6] if len(parts) >= 6 else []
            else:
                vals = []

            # fill into record (pad with None if missing)
            for idx, period in enumerate(["1Yr","3Yr","5Yr","10Yr"]):
                rec[f"{metric} {period}"] = vals[idx] if idx < len(vals) else None

            if len(vals) < 4:
                st.warning(f"⚠️ {fund} ({ticker}): only {len(vals)} peer values found for '{metric}'.")

        records.append(rec)

    if not records:
        st.warning("❌ No Peer Risk-Adjusted Return Rank data extracted.")
        return

    df = pd.DataFrame(records)
    st.session_state["step14_peer_rank_table"] = records
    st.dataframe(df, use_container_width=True)

#───Step 14.5: IPS Fail Table──────────────────────────────────────────────────────────────────

def step14_5_ips_fail_table():
    import streamlit as st
    import pandas as pd

    df = st.session_state.get("ips_icon_table")
    if df is None or df.empty:
        return

    fail_df = df[df["IPS Watch Status"].isin(["FW", "IW"])][["Fund Name", "IPS Watch Status"]]
    if fail_df.empty:
        return

    table_html = fail_df.rename(columns={
        "Fund Name": "Fund",
        "IPS Watch Status": "Watch Status"
    }).to_html(index=False, border=0, justify="center", classes="ips-fail-table")

    st.markdown(f"""
    <div style='
        background: linear-gradient(120deg, #e6f0fb 85%, #c8e0f6 100%);
        color: #23395d;
        border-radius: 1.3rem;
        box-shadow: 0 2px 14px rgba(44,85,130,0.08), 0 1px 4px rgba(36,67,105,0.07);
        padding: 1.6rem 2.0rem 1.6rem 2.0rem;
        max-width: 650px;
        margin: 1.4rem auto 1.2rem auto;
        border: 1.5px solid #b5d0eb;'>
        <div style='font-weight:700; color:#23395d; font-size:1.15rem; margin-bottom:0.5rem; letter-spacing:-0.5px;'>
            Funds on Watch
        </div>
        <div style='font-size:1rem; margin-bottom:1rem; color:#23395d;'>
            The following funds failed five or more IPS criteria and are currently on watch.
        </div>
        {table_html}
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <style>
    .ips-fail-table {
        width: 100%;
        border-collapse: collapse;
        margin-top: 0.7em;
        font-family: 'Inter', 'Segoe UI', Arial, sans-serif;
    }
    .ips-fail-table th, .ips-fail-table td {
        border: none;
        padding: 0.48em 1.1em;
        text-align: left;
        font-size: 1.07em;
    }
    .ips-fail-table th {
        background: #244369;
        color: #fff;
        font-weight: 700;
        letter-spacing: 0.01em;
    }
    .ips-fail-table td {
        color: #244369;
    }
    .ips-fail-table tr:nth-child(even) {background: #e6f0fb;}
    .ips-fail-table tr:nth-child(odd)  {background: #f8fafc;}
    </style>
    """, unsafe_allow_html=True)

#───Step 14.7: Proposal──────────────────────────────────────────────────────────────────

def extract_proposed_scorecard_blocks(pdf):
    import re
    import streamlit as st
    import pandas as pd
    from rapidfuzz import fuzz

    """
    Step 14.7: On only the 'Fund Scorecard: Proposed Funds' page, fuzzy-match the
    already-extracted fund names/tickers from the performance/scorecard and persist/display
    only those confirmed as proposed. The card view shows only the name and ticker.
    """
    prop_page = st.session_state.get("scorecard_proposed_page")
    if not prop_page:
        st.error("❌ 'Fund Scorecard: Proposed Funds' page number not found in TOC.")
        return pd.DataFrame()

    # 1. Extract lines from just the Proposed Funds scorecard page
    page = pdf.pages[prop_page - 1]
    lines = [ln.strip() for ln in (page.extract_text() or "").splitlines() if ln.strip()]
    if not lines:
        st.warning("No text found on the Proposed Funds scorecard page.")
        return pd.DataFrame()

    # 2. Build candidate list from already-extracted funds (performance/scorecard)
    perf_data = st.session_state.get("fund_performance_data", [])
    if not perf_data:
        st.warning("No performance/scorecard fund names available to match against.")
        return pd.DataFrame()

    candidate_funds = []
    for item in perf_data:
        name = item.get("Fund Scorecard Name", "").strip()
        ticker = item.get("Ticker", "").strip().upper()
        if name:
            candidate_funds.append({"Fund Scorecard Name": name, "Ticker": ticker})

    # 3. Fuzzy-match each candidate against the lines on the proposed page
    results = []
    for fund in candidate_funds:
        name = fund["Fund Scorecard Name"]
        ticker = fund["Ticker"]
        best_score = 0
        best_line = ""
        for line in lines:
            score_name = fuzz.token_sort_ratio(name.lower(), line.lower())
            score_ticker = fuzz.token_sort_ratio(ticker.lower(), line.lower()) if ticker else 0
            score = max(score_name, score_ticker)
            if score > best_score:
                best_score = score
                best_line = line
        found = best_score >= 70  # threshold; adjust if needed
        results.append({
            "Fund Scorecard Name": name,
            "Ticker": ticker,
            "Found on Proposed": "✅" if found else "❌",
            "Match Score": best_score,
            "Matched Line": best_line if found else ""
        })

    df = pd.DataFrame(results)

    # 4. Keep only confirmed proposed funds and persist independently of selection
    df_confirmed = df[df["Found on Proposed"] == "✅"].copy()
    st.session_state["proposed_funds_confirmed_df"] = df_confirmed

    # 5. Display styled summary card with only Fund and Ticker shown
    if df_confirmed.empty:
        st.markdown(f"""
        <div style="
            background: linear-gradient(120deg, #fff8f0 85%, #ffe9d8 100%);
            color: #8a5a2b;
            border-radius: 1.2rem;
            padding: 1rem 1.5rem;
            border: 1px solid #f0d4b5;
            margin-bottom: 1rem;
            font-size:1rem;
        ">
            No confirmed proposed funds were found on the Proposed Funds scorecard page.
        </div>
        """, unsafe_allow_html=True)
    else:
        # Simplified display: only name and ticker
        display_df = df_confirmed[["Fund Scorecard Name", "Ticker"]].rename(columns={
            "Fund Scorecard Name": "Fund",
        })

        table_html = display_df.to_html(index=False, border=0, justify="center", classes="proposed-fund-table")

        st.markdown(f"""
        <div style='
            background: linear-gradient(120deg, #e6f0fb 85%, #c8e0f6 100%);
            color: #23395d;
            border-radius: 1.3rem;
            box-shadow: 0 2px 14px rgba(44,85,130,0.08), 0 1px 4px rgba(36,67,105,0.07);
            padding: 1.6rem 2.0rem;
            max-width: 100%;
            margin: 1.4rem auto 1.2rem auto;
            border: 1.5px solid #b5d0eb;'>
            <div style='font-weight:700; color:#23395d; font-size:1.15rem; margin-bottom:0.5rem; letter-spacing:-0.5px;'>
                Confirmed Proposed Funds
            </div>
            <div style='font-size:1rem; margin-bottom:1rem; color:#23395d;'>
                The following funds were identified on the Proposed Funds scorecard page.
            </div>
            {table_html}
        </div>
        <style>
        .proposed-fund-table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 0.7em;
            font-family: 'Inter', 'Segoe UI', Arial, sans-serif;
        }}
        .proposed-fund-table th, .proposed-fund-table td {{
            border: none;
            padding: 0.48em 1.1em;
            text-align: left;
            font-size: 1em;
        }}
        .proposed-fund-table th {{
            background: #244369;
            color: #fff;
            font-weight: 700;
            letter-spacing: 0.01em;
        }}
        .proposed-fund-table td {{
            color: #23395d;
        }}
        .proposed-fund-table tr:nth-child(even) {{background: #e6f0fb;}}
        .proposed-fund-table tr:nth-child(odd)  {{background: #f8fafc;}}
        </style>
        """, unsafe_allow_html=True)

    return df_confirmed


#───Step 15: Single Fund──────────────────────────────────────────────────────────────────

def step15_display_selected_fund():
    import pandas as pd
    import streamlit as st
    import re

    facts = st.session_state.get("fund_factsheets_data", [])
    if not facts:
        st.info("Run Steps 1–14 to populate data before viewing fund details.")
        return

    fund_names = [f["Matched Fund Name"] for f in facts]
    selected_fund = st.selectbox("Select a fund to view details:", fund_names)
    st.session_state.selected_fund = selected_fund

    # --- NEW: pull confirmed proposed funds once, independent of selection ---
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposed_fund_names = (
        confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        if not confirmed_proposed_df.empty else []
    )

    st.write(f"Details for: {selected_fund}")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    factsheet_rec = next((row for row in factsheets if row["Matched Fund Name"] == selected_fund), None)
    
    fund_facts_table = st.session_state.get("step12_fund_facts_table", [])
    
    # Filter out metadata rows if present
    fund_facts_table = [row for row in fund_facts_table if row.get("Fund Name") and row.get("Fund Name").lower() != "metadata"]
    
    # Robust matching
    facts_rec = next((row for row in fund_facts_table if row.get("Fund Name") == selected_fund), None)
    if not facts_rec and factsheet_rec:
        factsheet_ticker = factsheet_rec.get("Matched Ticker")
        facts_rec = next((row for row in fund_facts_table if row.get("Ticker") == factsheet_ticker), None)
    if not facts_rec:
        facts_rec = next(
            (row for row in fund_facts_table if selected_fund.lower() in row.get("Fund Name", "").lower()),
            None
        )
    
    left_box = (
        f"""<div style='
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.2rem;
            box-shadow: 0 2px 12px rgba(44,85,130,0.09), 0 1px 4px rgba(36,67,105,0.07);
            padding: 1rem 1.2rem;
            min-width: 220px;
            max-width: 260px;
            margin: 0.3rem 1.2rem 0.3rem 0;
            border: 1.2px solid #b5d0eb;
            font-size: 1rem;
            display: inline-block;
            vertical-align: top;'>
            <div><b>Category:</b> {factsheet_rec.get("Category", "—")}</div>
            <div><b>Benchmark:</b> {factsheet_rec.get("Benchmark", "—")}</div>
            <div><b>Net Assets:</b> {factsheet_rec.get("Net Assets", "—")}</div>
            <div><b>Manager:</b> {factsheet_rec.get("Manager Name", "—")}</div>
            <div><b>Avg. Market Cap:</b> {factsheet_rec.get("Avg. Market Cap", "—")}</div>
        </div>"""
        if factsheet_rec else "<div style='display:inline-block; min-width:220px; color:#666;'>No factsheet info found.</div>"
    )
    
    right_box = (
        f"""<div style='
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.2rem;
            box-shadow: 0 2px 12px rgba(44,85,130,0.09), 0 1px 4px rgba(36,67,105,0.07);
            padding: 1rem 1.2rem;
            min-width: 220px;
            max-width: 260px;
            margin: 0.3rem 0 0.3rem 0;
            border: 1.2px solid #b5d0eb;
            font-size: 1rem;
            display: inline-block;
            vertical-align: top;'>
            <div><b>Manager Tenure:</b> {facts_rec.get("Manager Tenure Yrs.", "—")}</div>
            <div><b>Expense Ratio:</b> {facts_rec.get("Expense Ratio", "—")}</div>
            <div><b>Expense Ratio Rank:</b> {facts_rec.get("Expense Ratio Rank", "—")}</div>
            <div><b>Total Number of Holdings:</b> {facts_rec.get("Total Number of Holdings", "—")}</div>
            <div><b>Turnover Ratio:</b> {facts_rec.get("Turnover Ratio", "—")}</div>
        </div>"""
        if facts_rec else "<div style='display:inline-block; min-width:220px; color:#666;'>No Fund Facts available.</div>"
    )
    
    st.markdown(
        f"""
        <div style='
            width:100%;
            display:flex;
            flex-wrap:wrap;
            justify-content:center;
            align-items:flex-start;
            gap:24px;
            margin: 0.6rem 0 2rem 0;
        '>
            {left_box}{right_box}
        </div>
        """,
        unsafe_allow_html=True
    )


    # --- Slide 1 Table: IPS Results ---
    ips_icon_table = st.session_state.get("ips_icon_table")
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)

    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        if not row.empty:
            row_dict = row.iloc[0].to_dict()
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            row_df = pd.DataFrame([{
                "Category": fs_rec.get("Category", "") if fs_rec else "",
                "Time Period": st.session_state.get("report_date", ""),
                "Plan Assets": "$",  # Or replace with actual variable if you store this elsewhere!
                **{display_columns.get(k, k): v for k, v in row_dict.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row_dict.get("IPS Watch Status", "")
            }])

            def color_bool(v):
                return "background-color: green" if v == "✔" else ("background-color: red" if v == "✗" else "")
            
            def style_status(v):
                if v == "NW": return "background-color: green; color: white; font-weight: 600;"
                if v == "IW": return "background-color: orange; color: white; font-weight: 600;"
                if v == "FW": return "background-color: red; color: white; font-weight: 600;"
                return ""
            
            styled = row_df.style.applymap(color_bool, subset=[str(i) for i in range(1, 12)]) \
                                 .applymap(style_status, subset=["IPS Status"])
            
            st.dataframe(styled, use_container_width=True)

        else:
            st.warning("No IPS screening result found for selected fund.")
    else:
        st.warning("IPS screening table not found. Run earlier steps first.")

    # --- Slide 2 Table 1 ---
    st.markdown("**Net Expense Ratio**")
    
    def format_expense(val):
        if val is None or val == "":
            return ""
        s = str(val)
        if s.endswith("%"):
            return s
        return f"{s}%"
    
    def build_expense_row(fund_name):
        perf_data = st.session_state.get("fund_performance_data", [])
        item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        net_exp = format_expense(item.get("Net Expense Ratio", ""))
        return {
            "Investment Manager": inv_mgr,
            "Net Expense Ratio": net_exp
        }
    
    # Selected fund row
    row_selected = build_expense_row(selected_fund)
    
    # Proposed fund(s) — persistent, independent of selection
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_expense_row(pf))
    
    # Assemble final table: selected fund first, then proposed(s)
    all_rows = [row_selected] + proposed_rows
    df_slide2_table1 = pd.DataFrame(all_rows)
    
    # Save & display
    st.session_state["slide2_table1_data"] = df_slide2_table1
    st.dataframe(df_slide2_table1, use_container_width=True)


    # --- Slide 2 Table 2: Returns ---
    st.markdown("**Returns**")
    date_label = st.session_state.get("report_date", "QTD")
    
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s == "" else f"{s}%"
    
    def build_return_row(fund_name):
        perf_data = st.session_state.get("fund_performance_data", [])
        item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
    
        qtd   = append_pct(item.get("QTD", ""))
        one   = append_pct(item.get("1Yr", ""))
        three = append_pct(item.get("3Yr", ""))
        five  = append_pct(item.get("5Yr", ""))
        ten   = append_pct(item.get("10Yr", ""))
    
        return {
            "Investment Manager": inv_mgr,
            date_label:           qtd,
            "1 Year":             one,
            "3 Year":             three,
            "5 Year":             five,
            "10 Year":            ten
        }
    
    # Selected fund row
    row_selected = build_return_row(selected_fund)
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_return_row(pf))
    
    # Benchmark row (same as before)
    fs_rec = next((f for f in st.session_state.get("fund_factsheets_data", []) if f["Matched Fund Name"] == selected_fund), {})
    bench_name = fs_rec.get("Benchmark", "") if fs_rec else ""
    bench_ticker = fs_rec.get("Matched Ticker", "") if fs_rec else ""
    bench_inv_mgr = f"{bench_name} ({bench_ticker})" if bench_name and bench_ticker else (bench_name or "Benchmark")
    
    # Need the selected fund's perf_item to get benchmark returns
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    bench_qtd   = append_pct(perf_item.get("Bench QTD", ""))
    bench_one   = append_pct(perf_item.get("Bench 1Yr", ""))
    bench_3yr   = append_pct(perf_item.get("Bench 3Yr", ""))
    bench_5yr   = append_pct(perf_item.get("Bench 5Yr", ""))
    bench_ten   = append_pct(perf_item.get("Bench 10Yr", ""))
    
    row_benchmark = {
        "Investment Manager": bench_inv_mgr,
        date_label:           bench_qtd,
        "1 Year":             bench_one,
        "3 Year":             bench_3yr,
        "5 Year":             bench_5yr,
        "10 Year":            bench_ten
    }
    
    # Assemble: selected, proposed(s), then benchmark
    all_rows = [row_selected] + proposed_rows + [row_benchmark]
    df_slide2_table2 = pd.DataFrame(all_rows)
    
    # Save & display
    st.session_state["slide2_table2_data"] = df_slide2_table2
    st.dataframe(df_slide2_table2, use_container_width=True)

    # --- Slide 2 Table 3: Calendar Returns ---
    st.markdown("**Calendar Returns**")
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    if not fund_cy or not bench_cy:
        st.error("❌ No calendar year returns data found. Ensure Step 8 has been run correctly.")
        return
    
    def build_cy_row(name, is_selected=True):
        # find fund record
        fund_rec = next((r for r in fund_cy if r.get("Name") == name), None)
        if not fund_rec:
            return None
        ticker = fund_rec.get("Ticker", "")
        row = {"Investment Manager": f"{name} ({ticker})" if ticker else name}
        year_cols = [col for col in fund_rec.keys() if re.match(r"20\d{2}", col)]
        for year in year_cols:
            row[year] = fund_rec.get(year, "")
        return row, year_cols
    
    # Selected
    selected_row_tuple = build_cy_row(selected_fund)
    if not selected_row_tuple:
        st.error(f"❌ Could not find data for selected fund: {selected_fund}")
        return
    row_selected, year_cols = selected_row_tuple
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            tup = build_cy_row(pf)
            if tup:
                row_pf, _ = tup
                proposed_rows.append(row_pf)
    
    # Benchmark for selected fund
    bench_rec = None
    fund_rec = next((r for r in fund_cy if r.get("Name") == selected_fund), None)
    if fund_rec:
        bench_rec = next(
            (r for r in bench_cy if r.get("Name") == selected_fund or r.get("Ticker") == fund_rec.get("Ticker")),
            None
        )
    if not bench_rec:
        st.error(f"❌ Could not find benchmark data for selected fund: {selected_fund}")
        return
    bench_ticker = bench_rec.get("Ticker", "")
    bench_name_display = bench_rec.get("Name", "Benchmark")
    row_benchmark = {"Investment Manager": f"{bench_name_display} ({bench_ticker})" if bench_ticker else bench_name_display}
    for year in year_cols:
        row_benchmark[year] = bench_rec.get(year, "")
    
    # Assemble in order: selected, proposed(s), benchmark
    all_rows = [row_selected] + proposed_rows + [row_benchmark]
    df_slide2_3 = pd.DataFrame(all_rows, columns=["Investment Manager"] + year_cols)
    
    st.session_state["slide2_table3_data"] = df_slide2_3
    st.dataframe(df_slide2_3, use_container_width=True)



    # --- Slide 3 Table 1 ---
    st.markdown("**MPT Statistics Summary**")
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_risk_analysis_5yr", []) or st.session_state.get("step10_mpt_stats", [])
    
    def build_mpt_row(fund_name, stats3_list, stats5_list):
        stats3 = next((r for r in stats3_list if r["Fund Name"] == fund_name), {})
        stats5 = next((r for r in stats5_list if r["Fund Name"] == fund_name), {})
        ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        return {
            "Investment Manager":        inv_mgr,
            "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
            "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
            "3 Year Beta":               stats3.get("3 Year Beta", ""),
            "5 Year Beta":               stats5.get("5 Year Beta", ""),
            "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
            "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
            "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
            "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
        }
    
    # Selected fund row
    row_selected = build_mpt_row(selected_fund, mpt3, mpt5)
    
    # Proposed fund(s)
    proposed_rows = []
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        proposed_names = confirmed_proposed_df["Fund Scorecard Name"].unique().tolist()
        for pf in proposed_names:
            proposed_rows.append(build_mpt_row(pf, mpt3, mpt5))
    
    # Assemble: selected then proposed(s)
    all_rows = [row_selected] + proposed_rows
    df_slide3_1 = pd.DataFrame(all_rows)
    
    st.session_state["slide3_table1_data"] = df_slide3_1
    st.dataframe(df_slide3_1, use_container_width=True)

    # --- Slide 3 Table 2 ---
    st.markdown("**Risk-Adjusted Returns / Peer Ranking %**")
    risk_table = st.session_state.get("step13_risk_adjusted_table", [])
    peer_table = st.session_state.get("step14_peer_rank_table", [])
    
    def build_ratio_row(fund_name):
        risk_rec = next((r for r in risk_table if r.get("Fund Name") == fund_name), {})
        peer_rec = next((r for r in peer_table if r.get("Fund Name") == fund_name), {})
        ticker = risk_rec.get("Ticker") or peer_rec.get("Ticker", "")
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        def frac(metric, period):
            r = risk_rec.get(f"{metric} {period}", "")
            p = peer_rec.get(f"{metric} {period}", "")
            return f"{r} / {p}"
        return {
            "Investment Manager": inv_mgr,
            "3 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "3Yr"),
            "5 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "5Yr"),
            "3 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "3Yr"),
            "5 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "5Yr"),
            "3 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "3Yr"),
            "5 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "5Yr"),
        }
    
    # Selected fund
    rows = [build_ratio_row(selected_fund)]
    
    # Proposed fund(s)
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_ratio_row(pf))
    
    df_slide3_2 = pd.DataFrame(rows)

    # Save for Step 17 to use
    st.session_state["slide3_table2_data"] = df_slide3_2
    st.dataframe(df_slide3_2, use_container_width=True)

    # --- Slide 4 Table 1 ---
    st.markdown("**Manager Tenure**")
    blocks = st.session_state.get("fund_blocks", [])
    
    def build_tenure_row(fund_name):
        block = next((b for b in blocks if b["Fund Name"] == fund_name), {})
        raw_tenure = next((m["Info"] for m in block.get("Metrics", []) if m["Metric"] == "Manager Tenure"), "")
        m = re.search(r"(\d+(\.\d+)?)", raw_tenure)
        tenure = f"{m.group(1)} years" if m else raw_tenure
        # attempt to get ticker from performance_data for formatting
        perf_data = st.session_state.get("fund_performance_data", [])
        perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (perf_item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        return {
            "Investment Manager": inv_mgr,
            "Manager Tenure":     tenure
        }
    
    rows = [build_tenure_row(selected_fund)]
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_tenure_row(pf))
    
    df_slide4 = pd.DataFrame(rows)

    # Save for Step 17 to use
    st.session_state["slide4"] = df_slide4
    st.dataframe(df_slide4, use_container_width=True)

    
    # --- Slide 4 Table 2 ---
    st.markdown("**Assets**")
    facts = st.session_state.get("fund_factsheets_data", [])
    
    def build_asset_row(fund_name):
        # factsheet for the fund
        fs_rec = next((f for f in facts if f["Matched Fund Name"] == fund_name), None)
        # performance for ticker formatting
        perf_data = st.session_state.get("fund_performance_data", [])
        perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == fund_name), {})
        ticker = (perf_item.get("Ticker") or "").upper().strip()
        inv_mgr = f"{fund_name} ({ticker})" if ticker else fund_name
        assets = fs_rec.get("Net Assets", "") if fs_rec else ""
        avg_cap = fs_rec.get("Avg. Market Cap", "") if fs_rec else ""
        return {
            "Investment Manager":            inv_mgr,
            "Assets Under Management":       assets,
            "Average Market Capitalization": avg_cap
        }
    
    rows = [build_asset_row(selected_fund)]
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    if not confirmed_proposed_df.empty:
        for pf in confirmed_proposed_df["Fund Scorecard Name"].unique().tolist():
            rows.append(build_asset_row(pf))


    df_slide4_2 = pd.DataFrame(rows)
    
    # Save for Step 17 to use
    st.session_state["slide4_table2_data"] = df_slide4_2
    st.dataframe(df_slide4_2, use_container_width=True)

#───Bullet Points──────────────────────────────────────────────────────────────────

def step16_bullet_points():
    import streamlit as st

    selected_fund = st.session_state.get("selected_fund")
    if not selected_fund:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    perf_data = st.session_state.get("fund_performance_data", [])
    item = next((x for x in perf_data if x["Fund Scorecard Name"] == selected_fund), None)
    if not item:
        st.error(f"❌ Performance data for '{selected_fund}' not found.")
        return

    bullets = []

    # Bullet 1: Performance vs. Benchmark, using template
    template = st.session_state.get("bullet_point_templates", [""])[0]
    b1 = template
    for fld, val in item.items():
        b1 = b1.replace(f"[{fld}]", str(val))
    bullets.append(b1)
    st.markdown(f"- {b1}")

    # Get IPS status from icon table (guaranteed to match Slide 1 Table)
    ips_icon_table = st.session_state.get("ips_icon_table")
    ips_status = None
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        ips_status = row.iloc[0]["IPS Watch Status"] if not row.empty else None

    # Bullet 2: Watch status and return comparison
    if ips_status == "NW":
        b2 = "- This fund is **not on watch**."
        bullets.append(b2)
        st.markdown(b2)
    else:
        status_label = (
            "Formal Watch" if ips_status == "FW" else
            "Informal Watch" if ips_status == "IW" else
            ips_status or "on watch"
        )

        three   = float(item.get("3Yr") or 0)
        bench3  = float(item.get("Bench 3Yr") or 0)
        five    = float(item.get("5Yr") or 0)
        bench5  = float(item.get("Bench 5Yr") or 0)
        bps3 = round((three - bench3) * 100, 1)
        bps5 = round((five  - bench5) * 100, 1)

        # Peer rank logic (safe handling)
        peer = st.session_state.get("step14_peer_rank_table", [])
        raw3 = next((r.get("Sharpe Ratio 3Yr") or r.get("Sharpe Ratio Rank 3Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        raw5 = next((r.get("Sharpe Ratio 5Yr") or r.get("Sharpe Ratio Rank 5Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        try:
            pos3 = "top" if raw3 and int(raw3) <= 50 else "bottom"
        except:
            pos3 = "bottom"
        try:
            pos5 = "top" if raw5 and int(raw5) <= 50 else "bottom"
        except:
            pos5 = "bottom"

        b2 = (
            f"- The fund is now on **{status_label}**. Its 3‑year return trails the benchmark by "
            f"{bps3} bps ({three:.2f}% vs. {bench3:.2f}%) and its 5‑year return trails by "
            f"{bps5} bps ({five:.2f}% vs. {bench5:.2f}%). Its 3‑Yr Sharpe ranks in the {pos3} half of peers "
            f"and the {pos5} half of its 5‑Yr Sharpe ranks."
        )
        bullets.append(b2)
        st.markdown(b2)

    # Bullet 3: Action for Formal Watch only
    if ips_status == "FW":
        # Get confirmed proposed funds (stays constant regardless of selected fund)
        confirmed = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
        proposals = []
        if not confirmed.empty:
            # Unique ordered list of proposed fund names with tickers
            seen = set()
            for name, ticker in zip(confirmed["Fund Scorecard Name"], confirmed["Ticker"]):
                display = f"{name} ({ticker})" if ticker else name
                if display not in seen:
                    seen.add(display)
                    proposals.append(display)
        replacement = ", ".join(proposals) if proposals else "a proposed fund"
        b3 = f"- **Action:** Consider replacing this fund with {replacement}."
        bullets.append(b3)
        st.markdown(b3)



#───Build Powerpoint──────────────────────────────────────────────────────────────────
def step17_export_to_ppt():
    import streamlit as st
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from io import BytesIO
    import pandas as pd

    selected = st.session_state.get("selected_fund")
    if not selected:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    # Get confirmed proposed funds (name + ticker)
    confirmed_proposed_df = st.session_state.get("proposed_funds_confirmed_df", pd.DataFrame())
    proposed = []
    if not confirmed_proposed_df.empty:
        for _, row in confirmed_proposed_df.iterrows():
            name = row.get("Fund Scorecard Name", "")
            ticker = row.get("Ticker", "")
            label = f"{name} ({ticker})" if ticker else name
            proposed.append(label)
    # Limit to two proposals because template supports up to two
    proposed = proposed[:2]

    template_path = "assets/writeup&rec_template.pptx"
    try:
        prs = Presentation(template_path)
    except Exception as e:
        st.error(f"Could not load PowerPoint template: {e}")
        return

    def fill_table_with_styles(table, df_table, bold_row_idx=None, first_col_white=True):
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        for i in range(min(len(df_table), len(table.rows) - 1)):
            for j in range(min(len(df_table.columns), len(table.columns))):
                val = df_table.iloc[i, j]
                cell = table.cell(i + 1, j)
                cell.text = str(val) if val is not None else ""
                cell.vertical_alignment = MSO_VERTICAL_ANCHOR.MIDDLE
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        if j == 0:
                            run.font.color.rgb = RGBColor(255, 255, 255) if first_col_white else RGBColor(0, 0, 0)
                        else:
                            run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.bold = (bold_row_idx is not None and i == bold_row_idx)

    def fill_text_placeholder_preserving_format(slide, placeholder_text, replacement_text):
        replaced = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                if placeholder_text in full_text:
                    new_text = full_text.replace(placeholder_text, replacement_text)
                    runs = paragraph.runs
                    if len(runs) == 1:
                        runs[0].text = new_text
                    else:
                        for run in runs:
                            run.text = ""
                        avg_len = len(new_text) // len(runs)
                        idx = 0
                        for run in runs[:-1]:
                            run.text = new_text[idx:idx+avg_len]
                            idx += avg_len
                        runs[-1].text = new_text[idx:]
                    replaced = True
        return replaced

    def fill_bullet_points(slide, placeholder="[Bullet Point 1]", bullets=None):
        if bullets is None:
            bullets = st.session_state.get("bullet_points", None)
        if not bullets:
            bullets = ["Performance exceeded benchmark.", "No watch status.", "No action required."]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                if placeholder in p.text:
                    shape.text_frame.clear()
                    for b in bullets:
                        p_new = shape.text_frame.add_paragraph()
                        clean_text = b.replace("**", "")
                        p_new.text = clean_text
                        p_new.level = 0
                        p_new.font.name = "Cambria"
                        p_new.font.size = Pt(11)
                        p_new.font.color.rgb = RGBColor(0, 0, 0)
                        p_new.font.bold = True
                    return True
        return False

    # Gather session data for selected fund
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected), {})
    category = fs_rec.get("Category", "N/A")

    # IPS slide data
    df_slide1 = None
    ips_icon_table = st.session_state.get("ips_icon_table")
    if ips_icon_table is not None and not ips_icon_table.empty:
        filtered = ips_icon_table[ips_icon_table["Fund Name"] == selected]
        if not filtered.empty:
            row = filtered.iloc[0]
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            table_data = {
                **{display_columns.get(k, k): v for k, v in row.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row.get("IPS Watch Status", "")
            }
            table_data["Category"] = category
            table_data["Time Period"] = st.session_state.get("report_date", "")
            table_data["Plan Assets"] = "$"
            headers = ["Category", "Time Period", "Plan Assets"] + [str(i+1) for i in range(11)] + ["IPS Status"]
            df_slide1 = pd.DataFrame([table_data], columns=headers)

    # Slide 2 raw tables (with ordering injection of proposed between selected and benchmark)
    # Build Slide 2 Table 2 manually to guarantee order: selected, proposed(s), benchmark
    perf_item = next((p for p in st.session_state.get("fund_performance_data", []) if p.get("Fund Scorecard Name") == selected), {})
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s == "" else f"{s}%"
    date_label = st.session_state.get("report_date", "QTD")
    qtd   = append_pct(perf_item.get("QTD", ""))
    one   = append_pct(perf_item.get("1Yr", ""))
    three = append_pct(perf_item.get("3Yr", ""))
    five  = append_pct(perf_item.get("5Yr", ""))
    ten   = append_pct(perf_item.get("10Yr", ""))
    inv_mgr_selected = f"{selected} ({perf_item.get('Ticker','')})"
    row_fund = {
        "Investment Manager": inv_mgr_selected,
        date_label:           qtd,
        "1 Year":             one,
        "3 Year":             three,
        "5 Year":             five,
        "10 Year":            ten
    }

    # Proposed fund(s) rows
    proposed_rows = []
    for prop in proposed:
        # prop is "Name (TICKER)"
        # Try to retrieve its performance and expense ratio similarly
        raw_name = prop.split(" (")[0]
        perf_item_prop = next((p for p in st.session_state.get("fund_performance_data", []) if p.get("Fund Scorecard Name") == raw_name), {})
        qtd_p   = append_pct(perf_item_prop.get("QTD", ""))
        one_p   = append_pct(perf_item_prop.get("1Yr", ""))
        three_p = append_pct(perf_item_prop.get("3Yr", ""))
        five_p  = append_pct(perf_item_prop.get("5Yr", ""))
        ten_p   = append_pct(perf_item_prop.get("10Yr", ""))
        inv_mgr_prop = prop  # already has ticker in parentheses
        row_prop = {
            "Investment Manager": inv_mgr_prop,
            date_label:           qtd_p,
            "1 Year":             one_p,
            "3 Year":             three_p,
            "5 Year":             five_p,
            "10 Year":            ten_p
        }
        proposed_rows.append(row_prop)

    # Benchmark row (must come last)
    fs_rec_selected = fs_rec
    bench_name = fs_rec_selected.get("Benchmark", "") if fs_rec_selected else ""
    bench_ticker = fs_rec_selected.get("Matched Ticker", "") if fs_rec_selected else ""
    bench_inv_mgr = f"{bench_name} ({bench_ticker})" if bench_name else "Benchmark"
    bench_qtd   = append_pct(perf_item.get("Bench QTD", ""))
    bench_one   = append_pct(perf_item.get("Bench 1Yr", ""))
    bench_3yr   = append_pct(perf_item.get("Bench 3Yr", ""))
    bench_5yr   = append_pct(perf_item.get("Bench 5Yr", ""))
    bench_ten   = append_pct(perf_item.get("Bench 10Yr", ""))
    row_benchmark = {
        "Investment Manager": bench_inv_mgr,
        date_label:           bench_qtd,
        "1 Year":             bench_one,
        "3 Year":             bench_3yr,
        "5 Year":             bench_5yr,
        "10 Year":            bench_ten
    }

    # Assemble Slide 2 Table 2: selected, proposed(s), benchmark
    rows_for_table2 = [row_fund] + proposed_rows + [row_benchmark]
    df_slide2_table2 = pd.DataFrame(rows_for_table2)

    # Slide 2 Table 3 (Calendar Year): similar ordering logic
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    # Build calendar returns rows
    def build_cy_row(name, source_list):
        rec = next((r for r in source_list if r.get("Name") == name or r.get("Ticker") == name), None)
        if not rec:
            return None
        year_cols = [col for col in rec.keys() if re.match(r"20\d{2}", col)]
        inv_mgr = f"{name} ({rec.get('Ticker','')})"
        row = {"Investment Manager": inv_mgr}
        for year in year_cols:
            row[year] = rec.get(year, "")
        return row, year_cols

    # Selected fund calendar row
    fund_row_cy, cy_years = build_cy_row(selected, fund_cy) or ({}, [])
    bench_row_cy, _ = build_cy_row(bench_inv_mgr, bench_cy) or ({}, [])
    # Proposed funds calendar rows
    proposed_cy_rows = []
    for prop in proposed:
        raw_name = prop.split(" (")[0]
        row_prop_cy, _ = build_cy_row(raw_name, fund_cy) or (None, [])
        if row_prop_cy:
            proposed_cy_rows.append(row_prop_cy)
    # Assemble table3: selected, proposed(s), benchmark
    rows_for_table3 = []
    if fund_row_cy:
        rows_for_table3.append(fund_row_cy)
    rows_for_table3.extend(proposed_cy_rows)
    if bench_row_cy:
        rows_for_table3.append(bench_row_cy)
    df_slide2_table3 = pd.DataFrame(rows_for_table3, columns=["Investment Manager"] + (cy_years if cy_years else []))

    # Slide 2 Table 1: Net Expense Ratio (selected + proposed(s))
    def build_expense_row(name):
        perf_item_local = next((p for p in st.session_state.get("fund_performance_data", []) if p.get("Fund Scorecard Name") == name), {})
        net_exp = perf_item_local.get("Net Expense Ratio", "")
        if net_exp and not str(net_exp).endswith("%"):
            net_exp = f"{net_exp}%"
        inv_mgr_local = f"{name} ({perf_item_local.get('Ticker','')})"
        return {"Investment Manager": inv_mgr_local, "Net Expense Ratio": net_exp}

    rows_table1 = [build_expense_row(selected)]
    for prop in proposed:
        raw_name = prop.split(" (")[0]
        rows_table1.append(build_expense_row(raw_name))
    df_slide2_table1 = pd.DataFrame(rows_table1)

    # Slide 3 and others remain similar to prior logic (you can insert proposed fund context elsewhere if needed)
    df_slide1_table1 = df_slide1
    df_slide3_table1 = st.session_state.get("slide3_table1_data")
    df_slide3_table2 = st.session_state.get("slide3_table2_data")
    df_slide4_table1 = st.session_state.get("slide4")
    df_slide4_table2 = st.session_state.get("slide4_table2_data")

    # --- Fill Slide 1 ---
    slide1 = prs.slides[0]
    if not fill_text_placeholder_preserving_format(slide1, "[Fund Name]", selected):
        st.warning("Could not find the [Fund Name] placeholder on Slide 1.")
    if df_slide1_table1 is not None:
        filled = False
        for shape in slide1.shapes:
            if shape.has_table:
                table = shape.table
                if df_slide1_table1 is not None and len(table.columns) == len(df_slide1_table1.columns):
                    fill_table_with_styles(table, df_slide1_table1, first_col_white=False)
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table on Slide 1 to fill.")
    else:
        st.warning("Slide 1 IPS data not found in session state.")
    bullets = st.session_state.get("bullet_points", None)
    if not fill_bullet_points(slide1, "[Bullet Point 1]", bullets):
        st.warning("Could not find bullet points placeholder on Slide 1.")

    # --- Fill Slide 2 ---
    slide2 = prs.slides[3]
    fill_text_placeholder_preserving_format(slide2, "[Category]", category)
    # Table 1
    if df_slide2_table1 is None:
        st.warning("Slide 2 Table 1 data not found.")
    else:
        for shape in slide2.shapes:
            if shape.has_table and len(shape.table.columns) == len(df_slide2_table1.columns):
                fill_table_with_styles(shape.table, df_slide2_table1)
                break
    # Table 2 (with header adjustment)
    quarter_label = st.session_state.get("report_date", "QTD")
    for shape in slide2.shapes:
        if shape.has_table and len(shape.table.columns) == len(df_slide2_table2.columns):
            table = shape.table
            # second cell header
            table.cell(0, 1).text = quarter_label
            for c in range(len(table.columns)):
                cell = table.cell(0, c)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.bold = True
            fill_table_with_styles(table, df_slide2_table2, bold_row_idx=len(proposed))  # bold benchmark if it's last
            break
    # Table 3 (calendar year)
    for shape in slide2.shapes:
        if shape.has_table and df_slide2_table3 is not None and len(shape.table.columns) == len(df_slide2_table3.columns):
            table = shape.table
            # replace headers
            for c, col in enumerate(df_slide2_table3.columns):
                cell = table.cell(0, c)
                cell.text = str(col)
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.bold = True
            fill_table_with_styles(table, df_slide2_table3, bold_row_idx=len(proposed))  # benchmark bottom
            break

    # --- Replacement placeholders for proposed funds ---
    slide_repl1 = prs.slides[1]
    if proposed:
        fill_text_placeholder_preserving_format(slide_repl1, "[Replacement 1]", proposed[0])
    if len(proposed) > 1:
        slide_repl2 = prs.slides[2]
        fill_text_placeholder_preserving_format(slide_repl2, "[Replacement 2]", proposed[1])
    else:
        # remove the second replacement slide if only one proposed fund
        try:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[2])
        except Exception:
            pass  # safe fallback

    # --- Fill Slide 3 ---
    slide3 = prs.slides[4 if len(proposed) > 1 else 3]  # index shifts if slide 2 removed
    fill_text_placeholder_preserving_format(slide3, "[Category]", category)
    if df_slide3_table1 is None or df_slide3_table2 is None:
        st.warning("Slide 3 table data not found.")
    else:
        tables = [shape.table for shape in slide3.shapes if shape.has_table]
        if len(tables) >= 1 and df_slide3_table1 is not None:
            if len(df_slide3_table1.columns) == len(tables[0].columns):
                fill_table_with_styles(tables[0], df_slide3_table1)
        if len(tables) >= 2 and df_slide3_table2 is not None:
            if len(df_slide3_table2.columns) == len(tables[1].columns):
                fill_table_with_styles(tables[1], df_slide3_table2)

    # --- Fill Slide 4 ---
    slide4_index = 5 if len(proposed) > 1 else 4
    slide4 = prs.slides[slide4_index]
    qualitative_placeholder = f"[Category]– Qualitative Factors"
    qualitative_replacement = f"{category} - Qualitative Factors"
    if not fill_text_placeholder_preserving_format(slide4, qualitative_placeholder, qualitative_replacement):
        st.warning(f"Could not find placeholder '{qualitative_placeholder}' on Slide 4.")
    if df_slide4_table1 is not None:
        for shape in slide4.shapes:
            if shape.has_table and len(shape.table.columns) == len(df_slide4_table1.columns):
                fill_table_with_styles(shape.table, df_slide4_table1)
                break
    if df_slide4_table2 is not None:
        for shape in slide4.shapes:
            if shape.has_table and len(shape.table.columns) == len(df_slide4_table2.columns):
                fill_table_with_styles(shape.table, df_slide4_table2)
                break

    # --- Save and offer download ---
    output = BytesIO()
    prs.save(output)
    st.success("Powerpoint Generated")
    st.download_button(
        label="Download Writeup PowerPoint",
        data=output.getvalue(),
        file_name=f"{selected} Writeup.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


#───Main App──────────────────────────────────────────────────────────────────

def run():
    import re
    st.title("Writeup & Rec")
    uploaded = st.file_uploader("Upload MPI PDF to Generate Writup & Rec PPTX", type="pdf")
    if not uploaded:
        return

    with pdfplumber.open(uploaded) as pdf:
        # Step 1
        first = pdf.pages[0].extract_text() or ""
        process_page1(first)
        show_report_summary()

        # Step 2
        with st.expander("Table of Contents", expanded=False):
            toc_text = "".join((pdf.pages[i].extract_text() or "") for i in range(min(3, len(pdf.pages))))
            process_toc(toc_text)

        # --- Combined core details grouped ---
        with st.expander("All Fund Details", expanded=True):
            # 1. IPS Investment Screening
            with st.expander("IPS Investment Screening", expanded=True):
                sp = st.session_state.get('scorecard_page')
                tot = st.session_state.get('total_options')
                pp = st.session_state.get('performance_page')
                factsheets_page = st.session_state.get('factsheets_page')
                if sp and tot is not None and pp:
                    step3_5_6_scorecard_and_ips(pdf, sp, pp, factsheets_page, tot)
                else:
                    st.error("Missing scorecard, performance page, or total options")

            # 2. Fund Factsheets
            with st.expander("Fund Factsheets", expanded=True):
                names = [b['Fund Name'] for b in st.session_state.get('fund_blocks', [])]
                step6_process_factsheets(pdf, names)
                
            # 3. Extract Fund Facts sub-headings (Step 12) so Step 15 has data
            with st.expander("Fund Facts (sub-headings)", expanded=False):
                step12_process_fund_facts(pdf)
                
            # 3. Returns (annualized + calendar)
            with st.expander("Returns", expanded=False):
                step7_extract_returns(pdf)
                step8_calendar_returns(pdf)

            # 4. MPT Statistics Summary (requires risk analyses first)
            with st.expander("MPT Statistics Summary", expanded=False):
                step9_risk_analysis_3yr(pdf)
                step10_risk_analysis_5yr(pdf)
                step11_create_summary()

            # 5. Risk-Adjusted Returns and Peer Rank
            with st.expander("Risk-Adjusted Returns", expanded=False):
                step13_process_risk_adjusted_returns(pdf)
                step14_extract_peer_risk_adjusted_return_rank(pdf)

        # Data prep for bullet points (unchanged)
        report_date = st.session_state.get("report_date", "")
        m = re.match(r"(\d)(?:st|nd|rd|th)\s+QTR,\s*(\d{4})", report_date)
        quarter = m.group(1) if m else ""
        year = m.group(2) if m else ""

        for itm in st.session_state.get("fund_performance_data", []):
            qtd = float(itm.get("QTD") or 0)
            bench_qtd = float(itm.get("Bench QTD") or 0)
            itm["Perf Direction"] = "overperformed" if qtd >= bench_qtd else "underperformed"
            itm["Quarter"] = quarter
            itm["Year"] = year
            diff_bps = round((qtd - bench_qtd) * 100, 1)
            itm["QTD_bps_diff"] = str(diff_bps)
            fund_pct = f"{qtd:.2f}%"
            bench_pct = f"{bench_qtd:.2f}%"
            itm["QTD_pct_diff"] = f"{(qtd - bench_qtd):.2f}%"
            itm["QTD_vs"] = f"{fund_pct} vs. {bench_pct}"

        if "bullet_point_templates" not in st.session_state:
            st.session_state["bullet_point_templates"] = [
                "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], "
                "[Year] by [QTD_bps_diff] bps ([QTD_vs])."
            ]

        # Step 14.5: IPS Fail Table
        step14_5_ips_fail_table()

        extract_proposed_scorecard_blocks(pdf)

        # Step 15: View Single Fund Details
        with st.expander("Single Fund Write Up", expanded=False):
            step15_display_selected_fund()

        # Bullet Points
        with st.expander("Bullet Points", expanded=False):
            step16_bullet_points()

        # PowerPoint
        with st.expander("Export to Powerpoint", expanded=False):
            step17_export_to_ppt()

if __name__ == "__main__":
    run()

