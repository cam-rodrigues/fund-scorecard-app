import streamlit as st
import pandas as pd
import app_pages.write_up_processor as write_up_processor  
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO 
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches

def run():
    st.set_page_config(page_title="IPS Summary Table", layout="wide")
    st.title("Write Up Output")

    uploaded_file = st.file_uploader("Upload MPI PDF", type=["pdf"], key="writeup_upload")

    if not uploaded_file:
        st.warning("Please upload an MPI PDF to continue.")
        return

    if "fund_blocks" not in st.session_state:
        st.session_state["suppress_criteria_display"] = True
        st.session_state["suppress_scorecard_table"] = True
        st.session_state["suppress_matching_confirmation"] = True  
        write_up_processor.process_mpi(uploaded_file)
        st.success("File processed.")
    

    # Load data
    ips_results = st.session_state.get("step8_results", [])
    factsheet_data = st.session_state.get("fund_factsheets_data", [])
    quarter = st.session_state.get("report_quarter", "Unknown")

    if not ips_results or not factsheet_data:
        st.error("Missing processed data.")
        return

    # === Section 1: Full Table with All Funds ===
    rows = []
    for result in ips_results:
        fund_name = result["Fund Name"]
        ips_status = result["Overall IPS Status"]
        metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in result["IPS Metrics"]]
        metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

        fund_fact = next((f for f in factsheet_data if f["Matched Fund Name"] == fund_name), {})
        category = fund_fact.get("Category", "N/A")
        ticker = fund_fact.get("Matched Ticker", "N/A")

        row = {
            "Fund Name": fund_name,
            "Ticker": ticker,
            "Category": category,
            "Time Period": quarter,
            "Plan Assets": "$"
        }
        for i in range(11):
            row[str(i + 1)] = metric_results[i]
        row["IPS Status"] = ips_status

        rows.append(row)

    df_summary = pd.DataFrame(rows)
    st.subheader("IPS Summary Table (All Funds)")
    st.dataframe(df_summary, use_container_width=True)

    # === Section 2: View One Fund in Detail ===
    st.subheader("View Individual Fund Summary")
    fund_names = [res["Fund Name"] for res in ips_results]
    selected_fund = st.selectbox("Select a Fund", fund_names)

    selected_result = next((r for r in ips_results if r["Fund Name"] == selected_fund), None)
    selected_facts = next((f for f in factsheet_data if f["Matched Fund Name"] == selected_fund), {})

    if not selected_result:
        st.error("Selected fund not found.")
        return

    category = selected_facts.get("Category", "N/A")
    ticker = selected_facts.get("Matched Ticker", "N/A")
    ips_status = selected_result["Overall IPS Status"]
    metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in selected_result["IPS Metrics"]]
    metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

    row = {
        "Fund Name": selected_fund,
        "Ticker": ticker,
        "Category": category,
        "Time Period": quarter,
        "Plan Assets": "$"
    }
    for i in range(11):
        row[str(i + 1)] = metric_results[i]
    row["IPS Status"] = ips_status

    df_one = pd.DataFrame([row])
    st.dataframe(df_one, use_container_width=True)

    # === Section 3: Additional Factsheet Info ===
    st.subheader("Fund Factsheet Information")
    if selected_facts:
        st.markdown(f"""
        - **Ticker:** {selected_facts.get('Matched Ticker', 'N/A')}
        - **Benchmark:** {selected_facts.get('Benchmark', 'N/A')}
        - **Category:** {selected_facts.get('Category', 'N/A')}
        - **Net Assets:** {selected_facts.get('Net Assets', 'N/A')}
        - **Manager Name:** {selected_facts.get('Manager Name', 'N/A')}
        - **Avg. Market Cap:** {selected_facts.get('Avg. Market Cap', 'N/A')}
        - **Expense Ratio:** {selected_facts.get('Expense Ratio', 'N/A')}
        """)
    else:
        st.warning("No factsheet data found for the selected fund.")


def run():
    st.set_page_config(page_title="IPS Summary Table", layout="wide")
    st.title("Upload MPI & View IPS Summary Table")

    uploaded_file = st.file_uploader("Upload MPI PDF", type=["pdf"], key="writeup_upload")

    if not uploaded_file:
        st.warning("Please upload an MPI PDF to continue.")
        return

    if "fund_blocks" not in st.session_state:
        st.session_state["suppress_criteria_display"] = True
        st.session_state["suppress_scorecard_table"] = True
        st.session_state["suppress_matching_confirmation"] = True  
        write_up_processor.process_mpi(uploaded_file)
        st.success("File processed.")

    # Load data
    ips_results = st.session_state.get("step8_results", [])
    factsheet_data = st.session_state.get("fund_factsheets_data", [])
    quarter = st.session_state.get("report_quarter", "Unknown")

    if not ips_results or not factsheet_data:
        st.error("Missing processed data.")
        return

    # === Section 1: Full Table with All Funds ===
    rows = []
    for result in ips_results:
        fund_name = result["Fund Name"]
        ips_status = result["Overall IPS Status"]
        metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in result["IPS Metrics"]]
        metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

        fund_fact = next((f for f in factsheet_data if f["Matched Fund Name"] == fund_name), {})
        category = fund_fact.get("Category", "N/A")
        ticker = fund_fact.get("Matched Ticker", "N/A")

        row = {
            "Fund Name": fund_name,
            "Ticker": ticker,
            "Category": category,
            "Time Period": quarter,
            "Plan Assets": "$"
        }
        for i in range(11):
            row[str(i + 1)] = metric_results[i]
        row["IPS Status"] = ips_status

        rows.append(row)

    df_summary = pd.DataFrame(rows)
    st.subheader("IPS Summary Table (All Funds)")
    st.dataframe(df_summary, use_container_width=True)

    # === Section 2: View One Fund in Detail ===
    st.subheader("View Individual Fund Summary")
    fund_names = [res["Fund Name"] for res in ips_results]
    selected_fund = st.selectbox("Select a Fund", fund_names)

    selected_result = next((r for r in ips_results if r["Fund Name"] == selected_fund), None)
    selected_facts = next((f for f in factsheet_data if f["Matched Fund Name"] == selected_fund), {})

    if not selected_result:
        st.error("Selected fund not found.")
        return

    category = selected_facts.get("Category", "N/A")
    ticker = selected_facts.get("Matched Ticker", "N/A")
    ips_status = selected_result["Overall IPS Status"]
    metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in selected_result["IPS Metrics"]]
    metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

    row = {
        "Fund Name": selected_fund,
        "Ticker": ticker,
        "Category": category,
        "Time Period": quarter,
        "Plan Assets": "$"
    }
    for i in range(11):
        row[str(i + 1)] = metric_results[i]
    row["IPS Status"] = ips_status

    df_one = pd.DataFrame([row])
    st.dataframe(df_one, use_container_width=True)

    # === Section 3: Additional Factsheet Info ===
    st.subheader("Fund Factsheet Information")
    if selected_facts:
        st.markdown(f"""
        - **Ticker:** {selected_facts.get('Matched Ticker', 'N/A')}
        - **Benchmark:** {selected_facts.get('Benchmark', 'N/A')}
        - **Category:** {selected_facts.get('Category', 'N/A')}
        - **Net Assets:** {selected_facts.get('Net Assets', 'N/A')}
        - **Manager Name:** {selected_facts.get('Manager Name', 'N/A')}
        - **Avg. Market Cap:** {selected_facts.get('Avg. Market Cap', 'N/A')}
        - **Expense Ratio:** {selected_facts.get('Expense Ratio', 'N/A')}
        """)
    else:
        st.warning("No factsheet data found for the selected fund.")

def generate_watchlist_slide(df, selected_fund):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches

    def set_cell_border(cell, border_color=RGBColor(0, 0, 0)):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        r, g, b = border_color[0], border_color[1], border_color[2]
        hex_color = "%02X%02X%02X" % (r, g, b)
        for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
            ln = OxmlElement(line)
            ln.set("w", "12700")
            solidFill = OxmlElement("a:solidFill")
            srgbClr = OxmlElement("a:srgbClr")
            srgbClr.set("val", hex_color)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)

    def format_quarter(raw):
        import re
        raw = str(raw).strip()
    
        # Patterns like "Q1: 3/31/2025" or "Q1, 2025"
        match = re.search(r"Q([1-4])[,:\s-]*(\d{4})?", raw, re.IGNORECASE)
        if match:
            qtr, year = match.groups()
            suffix = {"1": "1st", "2": "2nd", "3": "3rd", "4": "4th"}[qtr]
            return f"{suffix} QTR {year}" if year else f"{suffix} QTR"
    
        return raw  # fallback if it doesn't match

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Layout 6 is typically a blank slide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add Procyon logo in upper-right corner
    logo_path = "assets/procyon_logo.png"
    logo_width = Inches(1.25)  # Smaller, proportionate size
    
    # Position it ~9.0 inches from left (slide is ~10" wide), minus logo width
    logo_left = Inches(9.0 - 1.25)
    logo_top = Inches(0.2)
    
    slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)



    # Manually add left-aligned title textbox to match subheading
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_box.line.fill.background()  # ✅ Removes the border/outline of the textbox
    
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Investment Watchlist"
    run.font.size = Pt(20)
    run.font.name = "HelveticaNeueLT Std Lt Ext"
    run.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.LEFT


    top = Inches(1.1)
    subheading = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.3))
    tf = subheading.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = selected_fund
    run.font.name = "Cambria"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(0, 0, 0)

    matching_rows = df[df["Fund Name"] == selected_fund]
    rows = len(matching_rows)
    cols = 15
    col_widths = [1.2, 1.2, 1.2] + [0.4] * 11 + [1.2]

    table_top = Inches(1.5)
    table_left = Inches(0.3)
    table_width = Inches(9)
    table_height = Inches(0.25 * (rows + 1))

    table = slide.shapes.add_table(rows + 1, cols, table_left, table_top, table_width, table_height).table

    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    headers = ["Category", "Time Period", "Plan Assets"] + [str(i) for i in range(1, 12)] + ["IPS Status"]

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        set_cell_border(cell)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        text_frame = cell.text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        p = text_frame.paragraphs[0]
        p.font.name = "Cambria"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    for row_idx, (_, r) in enumerate(matching_rows.iterrows(), start=1):
        row_vals = [
            r.get("Category", ""),
            format_quarter(r.get("Time Period", "")),
            r.get("Plan Assets", ""),
        ] + [r.get(str(i), "") for i in range(1, 12)] + [r.get("IPS Status", "")]

        for col_idx, val in enumerate(row_vals):
            cell = table.cell(row_idx, col_idx)
            set_cell_border(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            text_frame = cell.text_frame
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0

            p = text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = "Cambria"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

            if val == "Pass" and col_idx != 14:
                p.text = "✔"
                p.font.color.rgb = RGBColor(0, 176, 80)
            elif val == "Review" and col_idx != 14:
                p.text = "✖"
                p.font.color.rgb = RGBColor(192, 0, 0)
            elif col_idx == 14:
                p.text = ""
                val_str = str(val).strip().lower()

                if val_str == "formal warning":
                    badge_text = "FW"
                    badge_color = RGBColor(192, 0, 0)
                    font_color = RGBColor(255, 255, 255)
                elif val_str == "informal warning":
                    badge_text = "IW"
                    badge_color = RGBColor(255, 165, 0)
                    font_color = RGBColor(255, 255, 255)
                elif val_str == "passed ips screen":
                    badge_text = "✔"
                    badge_color = RGBColor(0, 176, 80)
                    font_color = RGBColor(255, 255, 255)
                else:
                    continue

                badge_left = table_left + sum(Inches(w) for w in col_widths[:col_idx]) + Inches(0.3)
                badge_top = table_top + Inches(0.25 * row_idx) + Inches(0.06)

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left=badge_left,
                    top=badge_top,
                    width=Inches(0.5),
                    height=Inches(0.25),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = badge_color
                shape.line.color.rgb = RGBColor(255, 255, 255)

                tf = shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = badge_text
                run.font.bold = True
                run.font.size = Pt(12 if badge_text == "✔" else 11)
                run.font.color.rgb = font_color
            else:
                p.text = str(val)

    return prs


# === Streamlit App ===
def run():
    st.set_page_config(page_title="IPS Summary Table", layout="wide")
    st.title("Upload MPI & View IPS Summary Table")

    uploaded_file = st.file_uploader("Upload MPI PDF", type=["pdf"], key="writeup_upload")
    if not uploaded_file:
        st.warning("Please upload an MPI PDF to continue.")
        return

    if "fund_blocks" not in st.session_state:
        st.session_state["suppress_criteria_display"] = True
        st.session_state["suppress_scorecard_table"] = True
        st.session_state["suppress_matching_confirmation"] = True
        write_up_processor.process_mpi(uploaded_file)
        st.success("File processed.")

    ips_results = st.session_state.get("step8_results", [])
    factsheet_data = st.session_state.get("fund_factsheets_data", [])
    quarter = st.session_state.get("report_quarter", "Unknown")

    if not ips_results or not factsheet_data:
        st.error("Missing processed data.")
        return

    rows = []
    for result in ips_results:
        fund_name = result["Fund Name"]
        ips_status = result["Overall IPS Status"]
        metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in result["IPS Metrics"]]
        metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

        fund_fact = next((f for f in factsheet_data if f["Matched Fund Name"] == fund_name), {})
        category = fund_fact.get("Category", "N/A")
        ticker = fund_fact.get("Matched Ticker", "N/A")

        row = {
            "Fund Name": fund_name,
            "Ticker": ticker,
            "Category": category,
            "Time Period": quarter,
            "Plan Assets": "$"
        }
        for i in range(11):
            row[str(i + 1)] = metric_results[i]
        row["IPS Status"] = ips_status
        rows.append(row)

    df_summary = pd.DataFrame(rows)
    st.subheader("IPS Summary Table (All Funds)")
    st.dataframe(df_summary, use_container_width=True)

    st.subheader("View Individual Fund Summary")
    fund_names = [res["Fund Name"] for res in ips_results]
    selected_fund = st.selectbox("Select a Fund", fund_names)

    selected_result = next((r for r in ips_results if r["Fund Name"] == selected_fund), None)
    selected_facts = next((f for f in factsheet_data if f["Matched Fund Name"] == selected_fund), {})

    if not selected_result:
        st.error("Selected fund not found.")
        return

    category = selected_facts.get("Category", "N/A")
    ticker = selected_facts.get("Matched Ticker", "N/A")
    ips_status = selected_result["Overall IPS Status"]
    metric_results = [m["Status"] if m["Status"] in ("Pass", "Review") else "N/A" for m in selected_result["IPS Metrics"]]
    metric_results = metric_results[:11] + ["N/A"] * (11 - len(metric_results))

    row = {
        "Fund Name": selected_fund,
        "Ticker": ticker,
        "Category": category,
        "Time Period": quarter,
        "Plan Assets": "$"
    }
    for i in range(11):
        row[str(i + 1)] = metric_results[i]
    row["IPS Status"] = ips_status

    df_one = pd.DataFrame([row])
    st.dataframe(df_one, use_container_width=True)

    st.subheader("Fund Factsheet Information")
    if selected_facts:
        st.markdown(f"""
        - **Ticker:** {selected_facts.get('Matched Ticker', 'N/A')}
        - **Benchmark:** {selected_facts.get('Benchmark', 'N/A')}
        - **Category:** {selected_facts.get('Category', 'N/A')}
        - **Net Assets:** {selected_facts.get('Net Assets', 'N/A')}
        - **Manager Name:** {selected_facts.get('Manager Name', 'N/A')}
        - **Avg. Market Cap:** {selected_facts.get('Avg. Market Cap', 'N/A')}
        - **Expense Ratio:** {selected_facts.get('Expense Ratio', 'N/A')}
        """)
    else:
        st.warning("No factsheet data found for the selected fund.")

    # === Export PowerPoint ===
    if "summary_df" not in st.session_state:
        st.session_state["summary_df"] = df_summary

    st.markdown("---")
    st.subheader("Export Selected Fund to PowerPoint")

    if st.button("Export to PowerPoint"):
        if selected_fund and not st.session_state["summary_df"].empty:
            ppt = generate_watchlist_slide(st.session_state["summary_df"], selected_fund)
            output = BytesIO()
            ppt.save(output)

            st.download_button(
                label="Download PowerPoint File",
                data=output.getvalue(),
                file_name=f"{selected_fund}_Watchlist.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        else:
            st.warning("Please select a fund and ensure data is loaded.")
