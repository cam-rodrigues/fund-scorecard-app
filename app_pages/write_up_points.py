import re
import streamlit as st
import pdfplumber
from calendar import month_name
import pandas as pd
from rapidfuzz import fuzz
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO


# === Utility: Extract & Label Report Date ===
def extract_report_date(text):
    # find the first quarter‐end or any mm/dd/yyyy
    dates = re.findall(r'(\d{1,2})/(\d{1,2})/(20\d{2})', text or "")
    for month, day, year in dates:
        m, d = int(month), int(day)
        # quarter‐end mapping
        if (m, d) in [(3,31), (6,30), (9,30), (12,31)]:
            q = { (3,31): "1st", (6,30): "2nd", (9,30): "3rd", (12,31): "4th" }[(m,d)]
            return f"{q} QTR, {year}"
        # fallback: human‐readable
        return f"As of {month_name[m]} {d}, {year}"
    return None

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 1 & 1.5: Page 1 Extraction ===
def process_page1(text):
    report_date = extract_report_date(text)
    if report_date:
        st.session_state['report_date'] = report_date
        st.success(f"Report Date: {report_date}")
    else:
        st.error("Could not detect report date on page 1.")

    m = re.search(r"Total Options:\s*(\d+)", text or "")
    st.session_state['total_options'] = int(m.group(1)) if m else None

    m = re.search(r"Prepared For:\s*\n(.*)", text or "")
    st.session_state['prepared_for'] = m.group(1).strip() if m else None

    m = re.search(r"Prepared By:\s*(.*)", text or "")
    pb = m.group(1).strip() if m else ""
    if not pb or "mpi stylus" in pb.lower():
        pb = "Procyon Partners, LLC"
    st.session_state['prepared_by'] = pb

    st.subheader("Page 1 Metadata")
    st.write(f"- Total Options: {st.session_state['total_options']}")
    st.write(f"- Prepared For: {st.session_state['prepared_for']}")
    st.write(f"- Prepared By: {pb}")


# === Step 2: Table of Contents Extraction ===
def process_toc(text):
    perf = re.search(r"Fund Performance[^\d]*(\d{1,3})", text or "")
    sc   = re.search(r"Fund Scorecard\s+(\d{1,3})", text or "")
    fs   = re.search(r"Fund Factsheets\s+(\d{1,3})", text or "")
    cy   = re.search(r"Fund Performance: Calendar Year\s+(\d{1,3})", text or "")
    r3yr = re.search(r"Risk Analysis: MPT Statistics \(3Yr\)\s+(\d{1,3})", text or "")
    r5yr = re.search(r"Risk Analysis: MPT Statistics \(5Yr\)\s+(\d{1,3})", text or "")

    perf_page = int(perf.group(1)) if perf else None
    sc_page   = int(sc.group(1))   if sc   else None
    fs_page   = int(fs.group(1))   if fs   else None
    cy_page   = int(cy.group(1))   if cy   else None
    r3yr_page = int(r3yr.group(1)) if r3yr else None
    r5yr_page = int(r5yr.group(1)) if r5yr else None

    st.subheader("Table of Contents Pages")
    st.write(f"- Fund Performance Current vs Proposed Comparison : {perf_page}")
    st.write(f"- Fund Performance Calendar Year : {cy_page}")
    st.write(f"- MPT 3Yr Risk Analysis : {r3yr_page}")
    st.write(f"- MPT 5Yr Risk Analysis : {r5yr_page}")
    st.write(f"- Fund Scorecard:   {sc_page}")
    st.write(f"- Fund Factsheets :  {fs_page}")
    


    # Store in session state for future reference
    st.session_state['performance_page'] = perf_page
    st.session_state['scorecard_page']   = sc_page
    st.session_state['factsheets_page']  = fs_page
    st.session_state['calendar_year_page'] = cy_page
    st.session_state['r3yr_page'] = r3yr_page
    st.session_state['r5yr_page'] = r5yr_page

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 3 ===
def step3_process_scorecard(pdf, start_page, declared_total):
    pages = []
    for p in pdf.pages[start_page-1:]:
        txt = p.extract_text() or ""
        if "Fund Scorecard" in txt:
            pages.append(txt)
        else:
            break
    lines = "\n".join(pages).splitlines()

    idx = next((i for i,l in enumerate(lines) if "Criteria Threshold" in l), None)
    if idx is not None:
        lines = lines[idx+1:]

    fund_blocks = []
    name = None
    metrics = []

    for i,line in enumerate(lines):
        m = re.match(r"^(.*?)\s+(Pass|Review)\s+(.+)$", line.strip())
        if not m:
            continue
        metric, _, info = m.groups()

        if metric == "Manager Tenure":
            if name and metrics:
                fund_blocks.append({"Fund Name": name, "Metrics": metrics})
            # find the fund name from the previous non-blank line
            prev = ""
            for j in range(i-1, -1, -1):
                if lines[j].strip():
                    prev = lines[j].strip()
                    break
            name = re.sub(r"Fund (Meets Watchlist Criteria|has been placed.*)", "", prev).strip()
            metrics = []

        if name:
            metrics.append({"Metric": metric, "Info": info})

    if name and metrics:
        fund_blocks.append({"Fund Name": name, "Metrics": metrics})

    st.session_state["fund_blocks"] = fund_blocks

    st.subheader("Step 3.5: Key Details per Metric")
    for b in fund_blocks:
        st.markdown(f"### {b['Fund Name']}")
        for m in b["Metrics"]:
            st.write(f"- **{m['Metric']}**: {m['Info'].strip()}")

    st.subheader("Step 3.6: Investment Option Count")
    count = len(fund_blocks)
    st.write(f"- Declared: **{declared_total}**")
    st.write(f"- Extracted: **{count}**")
    if count == declared_total:
        st.success("✅ Counts match.")
    else:
        st.error(f"❌ Expected {declared_total}, found {count}.")

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 4: IPS Screening ===
def step4_ips_screen():
    IPS = [
        "Manager Tenure",
        "Excess Performance (3Yr)",
        "R-Squared (3Yr)",
        "Peer Return Rank (3Yr)",
        "Sharpe Ratio Rank (3Yr)",
        "Sortino Ratio Rank (3Yr)",
        "Tracking Error Rank (3Yr)",
        "Excess Performance (5Yr)",
        "R-Squared (5Yr)",
        "Peer Return Rank (5Yr)",
        "Sharpe Ratio Rank (5Yr)",
        "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (5Yr)",
        "Expense Ratio Rank"
    ]
    st.subheader("Step 4: IPS Investment Criteria Screening")

    for b in st.session_state["fund_blocks"]:
        name = b["Fund Name"]
        is_passive = "bitcoin" in name.lower()
        statuses, reasons = {}, {}

        # Manager Tenure ≥3
        info = next((m["Info"] for m in b["Metrics"] if m["Metric"]=="Manager Tenure"), "")
        yrs = float(re.search(r"(\d+\.?\d*)", info).group(1)) if re.search(r"(\d+\.?\d*)", info) else 0
        ok = yrs>=3
        statuses["Manager Tenure"] = ok
        reasons["Manager Tenure"] = f"{yrs} yrs {'≥3' if ok else '<3'}"

        # map each IPS metric
        for metric in IPS[1:]:  # skip tenure
            m = next((x for x in b["Metrics"] if x["Metric"].startswith(metric.split()[0])), None)
            info = m["Info"] if m else ""
            if "Excess Performance" in metric:
                val_m = re.search(r"([-+]?\d*\.\d+)%", info)
                val = float(val_m.group(1)) if val_m else 0
                ok = (val>0) if "3Yr" in metric else (val>0)
                statuses[metric] = ok
                reasons[metric] = f"{val}%"
            elif "R-Squared" in metric:
                pct_m = re.search(r"(\d+\.\d+)%", info)
                pct = float(pct_m.group(1)) if pct_m else 0
                ok = (pct>=95) if is_passive else True
                statuses[metric] = ok
                reasons[metric] = f"{pct}%"
            elif "Peer Return" in metric or "Sharpe Ratio" in metric:
                rank_m = re.search(r"(\d+)", info)
                rank = int(rank_m.group(1)) if rank_m else 999
                ok = rank<=50
                statuses[metric] = ok
                reasons[metric] = f"Rank {rank}"
            elif "Sortino Ratio" in metric or "Tracking Error" in metric:
                rank_m = re.search(r"(\d+)", info)
                rank = int(rank_m.group(1)) if rank_m else 999
                if "Sortino" in metric and not is_passive:
                    ok = rank<=50
                elif "Tracking Error" in metric and is_passive:
                    ok = rank<90
                else:
                    ok = True
                statuses[metric] = ok
                reasons[metric] = f"Rank {rank}"
            elif "Expense Ratio" in metric:
                rank_m = re.search(r"(\d+)", info)
                rank = int(rank_m.group(1)) if rank_m else 999
                ok = rank<=50
                statuses[metric] = ok
                reasons[metric] = f"Rank {rank}"

        # count fails
        fails = sum(not v for v in statuses.values())
        if fails<=4:
            overall="Passed IPS Screen"
        elif fails==5:
            overall="Informal Watch (IW)"
        else:
            overall="Formal Watch (FW)"

        st.markdown(f"### {name} ({'Passive' if is_passive else 'Active'})")
        st.write(f"**Overall:** {overall} ({fails} fails)")
        for m in IPS:
            sym = "✅" if statuses.get(m,False) else "❌"
            st.write(f"- {sym} **{m}**: {reasons.get(m,'—')}")

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 5: Fund Performance Section Extraction (with fallback) ===
def step5_process_performance(pdf, start_page, fund_names):
    # figure out where the section ends
    end_page = st.session_state.get("factsheets_page") or (len(pdf.pages) + 1)

    # gather all lines and the raw text
    all_lines = []
    perf_text = ""
    for p in pdf.pages[start_page-1 : end_page-1]:
        txt = p.extract_text() or ""
        perf_text += txt + "\n"
        all_lines.extend(txt.splitlines())

    # first pass: normalized line → ticker (1–5 uppercase letters)
    mapping = {}
    for ln in all_lines:
        m = re.match(r"(.+?)\s+([A-Z]{1,5})$", ln.strip())
        if not m:
            continue
        raw_name, ticker = m.groups()
        norm = re.sub(r'[^A-Za-z0-9 ]+', '', raw_name).strip().lower()
        mapping[norm] = ticker

    # try matching each fund by normalized prefix
    tickers = {}
    for name in fund_names:
        norm_expected = re.sub(r'[^A-Za-z0-9 ]+', '', name).strip().lower()
        found = next(
            (t for raw, t in mapping.items() if raw.startswith(norm_expected)),
            None
        )
        tickers[name] = found

    # if too few, fallback to ordered scrape of every 1–5 letter code
    total = len(fund_names)
    found_count = sum(1 for t in tickers.values() if t)
    if found_count < total:
        all_tks = re.findall(r'\b([A-Z]{1,5})\b', perf_text)
        seen = []
        for tk in all_tks:
            if tk not in seen:
                seen.append(tk)
        tickers = {
            name: (seen[i] if i < len(seen) else None)
            for i, name in enumerate(fund_names)
        }

    # store & display
    st.session_state["tickers"] = tickers
    st.subheader("Step 5: Extracted Tickers")
    for n, t in tickers.items():
        st.write(f"- {n}: {t or '❌ not found'}")

    # validation
    st.subheader("Step 5.5: Ticker Count Validation")
    found_count = sum(1 for t in tickers.values() if t)
    st.write(f"- Expected tickers: **{total}**")
    st.write(f"- Found tickers:    **{found_count}**")
    if found_count == total:
        st.success("✅ All tickers found.")
    else:
        st.error(f"❌ Missing {total - found_count} ticker(s).")

    st.session_state["fund_performance_data"] = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in tickers.items()
    ]


def extract_field(text: str, label: str, stop_at: str = None) -> str:
    """
    Extracts the substring immediately following `label` up to `stop_at` (if provided),
    else returns the first whitespace-delimited token.
    """
    try:
        start = text.index(label) + len(label)
        rest  = text[start:]
        if stop_at and stop_at in rest:
            return rest[:rest.index(stop_at)].strip()
        return rest.split()[0].strip()
    except ValueError:
        return ""

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 6: Fund Factsheets ===
def step6_process_factsheets(pdf, fund_names):
    st.subheader("Step 6: Fund Factsheets Section")
    factsheet_start = st.session_state.get("factsheets_page")
    total_declared = st.session_state.get("total_options")
    performance_data = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in st.session_state.get("tickers", {}).items()
    ]

    if not factsheet_start:
        st.error("❌ 'Fund Factsheets' page number not found in TOC.")
        return

    matched_factsheets = []
    # Iterate pages from factsheet_start to end
    for i in range(factsheet_start - 1, len(pdf.pages)):
        page = pdf.pages[i]
        words = page.extract_words(use_text_flow=True)
        header_words = [w['text'] for w in words if w['top'] < 100]
        first_line = " ".join(header_words).strip()

        if not first_line or "Benchmark:" not in first_line or "Expense Ratio:" not in first_line:
            continue

        ticker_match = re.search(r"\b([A-Z]{5})\b", first_line)
        ticker = ticker_match.group(1) if ticker_match else ""
        fund_name_raw = first_line.split(ticker)[0].strip() if ticker else first_line

        best_score = 0
        matched_name = matched_ticker = ""
        for item in performance_data:
            ref = f"{item['Fund Scorecard Name']} {item['Ticker']}".strip()
            score = fuzz.token_sort_ratio(f"{fund_name_raw} {ticker}".lower(), ref.lower())
            if score > best_score:
                best_score, matched_name, matched_ticker = score, item['Fund Scorecard Name'], item['Ticker']

        def extract_field(label, text, stop=None):
            try:
                start = text.index(label) + len(label)
                rest = text[start:]
                if stop and stop in rest:
                    return rest[:rest.index(stop)].strip()
                return rest.split()[0]
            except Exception:
                return ""

        benchmark = extract_field("Benchmark:", first_line, "Category:")
        category  = extract_field("Category:", first_line, "Net Assets:")
        net_assets= extract_field("Net Assets:", first_line, "Manager Name:")
        manager   = extract_field("Manager Name:", first_line, "Avg. Market Cap:")
        avg_cap   = extract_field("Avg. Market Cap:", first_line, "Expense Ratio:")
        expense   = extract_field("Expense Ratio:", first_line)

        matched_factsheets.append({
            "Page #": i + 1,
            "Parsed Fund Name": fund_name_raw,
            "Parsed Ticker": ticker,
            "Matched Fund Name": matched_name,
            "Matched Ticker": matched_ticker,
            "Benchmark": benchmark,
            "Category": category,
            "Net Assets": net_assets,
            "Manager Name": manager,
            "Avg. Market Cap": avg_cap,
            "Expense Ratio": expense,
            "Match Score": best_score,
            "Matched": "✅" if best_score > 20 else "❌"
        })

    df_facts = pd.DataFrame(matched_factsheets)
    st.session_state['fund_factsheets_data'] = matched_factsheets

    display_df = df_facts[[
        "Matched Fund Name", "Matched Ticker", "Benchmark", "Category",
        "Net Assets", "Manager Name", "Avg. Market Cap", "Expense Ratio", "Matched"
    ]].rename(columns={"Matched Fund Name": "Fund Name", "Matched Ticker": "Ticker"})

    st.dataframe(display_df, use_container_width=True)

    matched_count = sum(1 for r in matched_factsheets if r["Matched"] == "✅")
    if not st.session_state.get("suppress_matching_confirmation", False):
        st.write(f"Matched {matched_count} of {len(matched_factsheets)} factsheet pages.")
        if matched_count == total_declared:
            st.success(f"All {matched_count} funds matched the declared Total Options from Page 1.")
        else:
            st.error(f"Mismatch: Page 1 declared {total_declared}, but only matched {matched_count}.")

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense Ratio & Bench QTD ===
def step7_extract_returns(pdf):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    st.subheader("Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense & Benchmark QTD")

    # 1) Where to scan
    perf_page = st.session_state.get("performance_page")
    end_page  = st.session_state.get("calendar_year_page") or (len(pdf.pages) + 1)
    perf_data = st.session_state.get("fund_performance_data", [])
    if perf_page is None or not perf_data:
        st.error("❌ Run Step 5 first to populate performance data.")
        return

    # 2) Prep output slots
    fields = ["QTD", "1Yr", "3Yr", "5Yr", "10Yr", "Net Expense Ratio", "Bench QTD", "Bench 3Yr", "Bench 5Yr"]
    for itm in perf_data:
        for f in fields:
            itm.setdefault(f, None)

    # 3) Gather every nonblank line in the Performance section
    lines = []
    for pnum in range(perf_page - 1, end_page - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 4) Regex to pull decimal tokens (with optional % and parentheses)
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    matched = 0
    for item in perf_data:
        name = item["Fund Scorecard Name"]
        tk   = item["Ticker"].upper().strip()

        # a) Exact-ticker match
        idx = next(
            (i for i, ln in enumerate(lines)
             if re.search(rf"\b{re.escape(tk)}\b", ln)),
            None
        )
        # b) Fuzzy-name fallback
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                st.warning(f"⚠️ {name} ({tk}): no match found.")
                continue

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        if len(clean) < 8:
            clean += [None] * (8 - len(clean))

        # d) Map fund returns & net expense
        item["QTD"]               = clean[0]
        item["1Yr"]               = clean[2]
        item["3Yr"]               = clean[3]
        item["5Yr"]               = clean[4]
        item["10Yr"]              = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 3Yr, and 5Yr from the very next line (or one more down)
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]

        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None

        matched += 1

    # 5) Save & display
    st.session_state["fund_performance_data"] = perf_data
    df = pd.DataFrame(perf_data)

    st.success(f"✅ Matched {matched} fund(s) with return data.")
    for itm in perf_data:
        missing = [f for f in fields if not itm.get(f)]
        if missing:
            st.warning(f"⚠️ Incomplete for {itm['Fund Scorecard Name']} ({itm['Ticker']}): missing {', '.join(missing)}")

    st.dataframe(
        df[["Fund Scorecard Name", "Ticker"] + fields],
        use_container_width=True
    )

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 8 Calendar Year Returns (funds + benchmarks) ===
def step8_calendar_returns(pdf):
    import re, streamlit as st, pandas as pd

    st.subheader("Step 8: Calendar Year Returns")

    # 1) Figure out section bounds
    cy_page  = st.session_state.get("calendar_year_page")
    end_page = st.session_state.get("r3yr_page", len(pdf.pages) + 1)
    if cy_page is None:
        st.error("❌ 'Fund Performance: Calendar Year' not found in TOC.")
        return

    # 2) Pull every line from that section
    all_lines = []
    for p in pdf.pages[cy_page-1 : end_page-1]:
        all_lines.extend((p.extract_text() or "").splitlines())

    # 3) Identify header & years
    header = next((ln for ln in all_lines if "Ticker" in ln and re.search(r"20\d{2}", ln)), None)
    if not header:
        st.error("❌ Couldn’t find header row with 'Ticker' + year.")
        return
    years = re.findall(r"\b20\d{2}\b", header)
    num_rx = re.compile(r"-?\d+\.\d+%?")

    # — A) Funds themselves —
    fund_map     = st.session_state.get("tickers", {})
    fund_records = []
    for name, tk in fund_map.items():
        ticker = (tk or "").upper()
        idx    = next((i for i, ln in enumerate(all_lines) if ticker in ln.split()), None)
        raw    = num_rx.findall(all_lines[idx-1]) if idx not in (None, 0) else []
        vals   = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec    = {"Name": name, "Ticker": ticker}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        fund_records.append(rec)

    df_fund = pd.DataFrame(fund_records)
    if not df_fund.empty:
        st.markdown("**Fund Calendar‑Year Returns**")
        st.dataframe(df_fund[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["step8_returns"] = fund_records

    # — B) Benchmarks matched back to each fund’s ticker —
    facts         = st.session_state.get("fund_factsheets_data", [])
    bench_records = []
    for f in facts:
        bench_name = f.get("Benchmark", "").strip()
        fund_tkr   = f.get("Matched Ticker", "")
        if not bench_name:
            continue

        # find the first line containing the benchmark name
        idx = next((i for i, ln in enumerate(all_lines) if bench_name in ln), None)
        if idx is None:
            continue
        raw  = num_rx.findall(all_lines[idx])
        vals = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec  = {"Name": bench_name, "Ticker": fund_tkr}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        bench_records.append(rec)

    df_bench = pd.DataFrame(bench_records)
    if not df_bench.empty:
        st.markdown("**Benchmark Calendar‑Year Returns**")
        st.dataframe(df_bench[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["benchmark_calendar_year_returns"] = bench_records
    else:
        st.warning("No benchmark returns extracted.")

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 9: 3‑Yr Risk Analysis – Match & Extract MPT Stats (hidden matching) ===
def step9_risk_analysis_3yr(pdf):
    import re, streamlit as st, pandas as pd
    from rapidfuzz import fuzz

    st.subheader("Step 9: Risk Analysis (3Yr) – MPT Statistics")

    # 1) Get your fund→ticker map
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (3Yr)” page
    start_page = st.session_state.get("r3yr_page")
    if not start_page:
        st.error("❌ ‘Risk Analysis: MPT Statistics (3Yr)’ page not found; run Step 2 first.")
        return

    # 3) Scan forward until you’ve seen each ticker (no display)
    locs = {}
    for pnum in range(start_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for fname, tk in fund_map.items():
                if fname in locs: 
                    continue
                if tk.upper() in tokens:
                    locs[fname] = {"page": pnum, "line": li}
        if len(locs) == len(fund_map):
            break

    # 4) Extract the first four numeric MPT stats from that same line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, info in locs.items():
        page = pdf.pages[info["page"]-1]
        lines = (page.extract_text() or "").splitlines()
        line = lines[info["line"]]
        nums = num_rx.findall(line)
        nums += [None] * (4 - len(nums))
        alpha, beta, up, down = nums[:4]
        results.append({
            "Fund Name":               name,
            "Ticker":                  fund_map[name].upper(),
            "3 Year Alpha":            alpha,
            "3 Year Beta":             beta,
            "3 Year Upside Capture":   up,
            "3 Year Downside Capture": down
        })

    # 5) Display final table only
    st.session_state["step9_mpt_stats"] = results
    df = pd.DataFrame(results)
    st.dataframe(df, use_container_width=True)


#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 10: Risk Analysis (5Yr) – Match & Extract MPT Statistics ===
def step10_risk_analysis_5yr(pdf):
    import re, streamlit as st, pandas as pd

    st.subheader("Step 10: Risk Analysis (5Yr) – MPT Statistics")

    # 1) Your fund→ticker map from Step 5
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (5Yr)” section
    section_page = next(
        (i for i, pg in enumerate(pdf.pages, start=1)
         if "Risk Analysis: MPT Statistics (5Yr)" in (pg.extract_text() or "")),
        None
    )
    if section_page is None:
        st.error("❌ Could not find ‘Risk Analysis: MPT Statistics (5Yr)’ section.")
        return

    # 3) Under‑the‑hood: scan pages until each ticker is located
    locs = {}
    total = len(fund_map)
    for pnum in range(section_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for name, tk in fund_map.items():
                if name in locs:
                    continue
                if tk.upper() in tokens:
                    locs[name] = {"page": pnum, "line": li}
        if len(locs) == total:
            break

    # 4) Wrap‑aware extraction of the first four floats after each ticker line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, tk in fund_map.items():
        info = locs.get(name)
        vals = [None] * 4
        if info:
            page = pdf.pages[info["page"] - 1]
            text_lines = (page.extract_text() or "").splitlines()
            idx = info["line"]
            nums = []
            # look on the line of the ticker and up to the next 2 lines
            for j in range(idx, min(idx + 3, len(text_lines))):
                nums += num_rx.findall(text_lines[j])
                if len(nums) >= 4:
                    break
            nums += [None] * (4 - len(nums))
            vals = nums[:4]
        else:
            st.warning(f"⚠️ {name} ({tk.upper()}): not found after page {section_page}.")

        alpha5, beta5, up5, down5 = vals
        results.append({
            "Fund Name":               name,
            "Ticker":                  tk.upper(),
            "5 Year Alpha":            alpha5,
            "5 Year Beta":             beta5,
            "5 Year Upside Capture":   up5,
            "5 Year Downside Capture": down5,
        })

    # 5) Save & display only the consolidated table
    st.session_state["step10_mpt_stats"] = results
    df = pd.DataFrame(results)
    st.dataframe(df, use_container_width=True)


#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 11: Combined MPT Statistics Summary ===
def step11_create_summary(pdf=None):
    import pandas as pd
    import streamlit as st

    st.subheader("Step 11: MPT Statistics Summary")

    # 1) Load your 3‑Yr and 5‑Yr stats from session state
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    if not mpt3 or not mpt5:
        st.error("❌ Missing MPT stats. Run Steps 9 & 10 first.")
        return

    # 2) Build DataFrames
    df3 = pd.DataFrame(mpt3)  # contains "3 Year Alpha", "3 Year Beta", etc.
    df5 = pd.DataFrame(mpt5)  # contains "5 Year Alpha", "5 Year Beta", etc.

    # 3) Merge on Fund Name & Ticker
    df = pd.merge(
        df3,
        df5,
        on=["Fund Name", "Ticker"],
        how="outer",
        suffixes=("_3yr", "_5yr")
    )

    # 4) Build the Investment Manager column
    df.insert(0, "Investment Manager", df["Fund Name"] + " (" + df["Ticker"] + ")")

    # 5) Select & order the columns
    df = df[[
        "Investment Manager",
        "3 Year Alpha",
        "5 Year Alpha",
        "3 Year Beta",
        "5 Year Beta",
        "3 Year Upside Capture",
        "3 Year Downside Capture",
        "5 Year Upside Capture",
        "5 Year Downside Capture"
    ]]

    # 6) Display
    st.session_state["step11_summary"] = df.to_dict("records")
    st.dataframe(df)

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 12: Extract “FUND FACTS” & Its Table Details in One Go ===
def step12_process_fund_facts(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 12: Fund Facts")

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # the exact labels and the order you want them in the table
    labels = [
        "Manager Tenure Yrs.",
        "Expense Ratio",
        "Expense Ratio Rank",
        "Total Number of Holdings",
        "Turnover Ratio"
    ]

    records = []
    # scan each factsheet page
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = pdf.pages[pnum-1].extract_text().splitlines()

        for idx, line in enumerate(lines):
            if line.lstrip().upper().startswith("FUND FACTS"):
                # grab the next 8 lines (should contain your 5 labels)
                snippet = lines[idx+1 : idx+1+8]
                rec = {"Fund Name": fund_name, "Ticker": ticker}
                for lab in labels:
                    val = None
                    for ln in snippet:
                        norm = " ".join(ln.strip().split())
                        if norm.startswith(lab):
                            rest = norm[len(lab):].strip(" :\t")
                            m = re.match(r"(-?\d+\.\d+)", rest)
                            val = m.group(1) if m else (rest.split()[0] if rest else None)
                            break
                    rec[lab] = val
                records.append(rec)
                break  # move on to the next page once Fund Facts is processed

    if not records:
        st.warning("No Fund Facts tables found.")
        return

    # save & show
    st.session_state["step12_fund_facts_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# === Step 13: Extract Risk‑Adjusted Returns Metrics ===
def step13_process_risk_adjusted_returns(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 13: Risk‑Adjusted Returns")

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # which metrics to pull
    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    num_rx  = re.compile(r"-?\d+\.\d+")

    records = []
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()

        # find the heading
        for idx, line in enumerate(lines):
            norm = " ".join(line.strip().split()).upper()
            if norm.startswith("RISK-ADJUSTED RETURNS"):
                snippet = lines[idx+1 : idx+1+6]  # grab next few lines
                rec = {"Fund Name": fund_name, "Ticker": ticker}

                for metric in metrics:
                    # find the snippet line for this metric
                    text_line = next(
                        ( " ".join(ln.strip().split())
                          for ln in snippet
                          if ln.strip().upper().startswith(metric.upper()) ),
                        None
                    ) or ""
                    # extract up to 4 numbers
                    nums = num_rx.findall(text_line)
                    nums += [None] * (4 - len(nums))

                    # assign into rec
                    rec[f"{metric} 1Yr"]  = nums[0]
                    rec[f"{metric} 3Yr"]  = nums[1]
                    rec[f"{metric} 5Yr"]  = nums[2]
                    rec[f"{metric} 10Yr"] = nums[3]

                records.append(rec)
                break  # done with this page

    if not records:
        st.warning("No 'RISK‑ADJUSTED RETURNS' tables found.")
        return

    # save & show
    st.session_state["step13_risk_adjusted_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––

# == Step 14: Peer Risk-Adjusted Return Rank ==
def step14_extract_peer_risk_adjusted_return_rank(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 14: Peer Risk-Adjusted Return Rank")

    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    page_map = {
        f["Page #"]:(f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    records = []

    for pnum, (fund, ticker) in page_map.items():
        page = pdf.pages[pnum-1]
        text = page.extract_text() or ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        # 1) locate the Risk-Adjusted Returns header
        try:
            risk_idx = next(i for i, ln in enumerate(lines)
                            if "RISK-ADJUSTED RETURNS" in ln.upper())
        except StopIteration:
            st.warning(f"⚠️ {fund} ({ticker}): Risk-Adjusted header not found.")
            continue

        # 2) find all “1 Yr 3 Yrs 5 Yrs 10 Yrs” lines after that
        header_idxs = [i for i, ln in enumerate(lines)
                       if re.match(r"1\s*Yr", ln)]
        peer_header_idxs = [i for i in header_idxs if i > risk_idx]

        if not peer_header_idxs:
            st.warning(f"⚠️ {fund} ({ticker}): peer header not found.")
            continue

        # take the *second* header occurrence (first is Risk-Adjusted, next is Peer)
        peer_hdr = peer_header_idxs[0] if len(peer_header_idxs)==1 else peer_header_idxs[1]

        rec = {"Fund Name": fund, "Ticker": ticker}

        # 3) read the three lines immediately below that header
        for offset, metric in enumerate(metrics, start=1):
            if peer_hdr + offset < len(lines):
                parts = lines[peer_hdr + offset].split()
                # parts[0:2] = metric name words, parts[2:6] = the four integer ranks
                vals = parts[2:6] if len(parts) >= 6 else []
            else:
                vals = []

            # fill into record (pad with None if missing)
            for idx, period in enumerate(["1Yr","3Yr","5Yr","10Yr"]):
                rec[f"{metric} {period}"] = vals[idx] if idx < len(vals) else None

            if len(vals) < 4:
                st.warning(f"⚠️ {fund} ({ticker}): only {len(vals)} peer values found for '{metric}'.")

        records.append(rec)

    if not records:
        st.warning("❌ No Peer Risk-Adjusted Return Rank data extracted.")
        return

    df = pd.DataFrame(records)
    st.session_state["step14_peer_rank_table"] = records
    st.dataframe(df, use_container_width=True)

#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
# === Step 15: Single Fund Details ===
def step15_display_selected_fund():
    import pandas as pd
    import streamlit as st
    import re

    st.subheader("Step 15: Single Fund Details")
    facts = st.session_state.get("fund_factsheets_data", [])
    if not facts:
        st.info("Run Steps 1–14 to populate data before viewing fund details.")
        return

    # Select a fund
    fund_names = [f["Matched Fund Name"] for f in facts]
    selected_fund = st.selectbox("Select a fund to view details:", fund_names)
    
    # Save the selected fund in session state
    st.session_state.selected_fund = selected_fund  # Save the selected fund in session_state

    # Now use this selected fund for further details
    st.write(f"Details for: {selected_fund}")

    # Display the fund details as before
    # --- (Existing code for displaying the details of the selected fund) ---

    # === Step 1: Page 1 Metadata ===
    st.markdown("**Step 1: Page 1 Metadata**")
    st.write(f"- Report Date:   {st.session_state.get('report_date','N/A')}")
    st.write(f"- Total Options: {st.session_state.get('total_options','N/A')}")
    st.write(f"- Prepared For:  {st.session_state.get('prepared_for','N/A')}")
    st.write(f"- Prepared By:   {st.session_state.get('prepared_by','N/A')}")

    # === Step 2: Table of Contents Pages ===
    st.markdown("**Step 2: Table of Contents**")
    for key,label in [
        ("performance_page","Fund Performance Current vs Proposed"),
        ("calendar_year_page","Fund Performance Calendar Year"),
        ("r3yr_page","MPT 3Yr Risk Analysis"),
        ("r5yr_page","MPT 5Yr Risk Analysis"),
        ("scorecard_page","Fund Scorecard"),
        ("factsheets_page","Fund Factsheets")
    ]:
        st.write(f"- {label}: {st.session_state.get(key,'N/A')}")

    # === Step 3: Scorecard Metrics ===
    st.markdown("**Step 3: Scorecard Metrics**")
    blocks = st.session_state.get("fund_blocks", [])
    block = next((b for b in blocks if b["Fund Name"] == selected_fund), None)
    if block:
        for m in block["Metrics"]:
            st.write(f"- {m['Metric']}: {m['Info']}")
    else:
        st.write("_No scorecard data found._")

    # === Slide 1 Table ===
    st.markdown("**Slide 1 Table**")

    # 1) Category from factsheet
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), {})
    category = fs_rec.get("Category","")

    # 2) Build first 11 IPS criteria
    IPS = [
      "Manager Tenure","Excess Performance (3Yr)","R‑Squared (3Yr)",
      "Peer Return Rank (3Yr)","Sharpe Ratio Rank (3Yr)","Sortino Ratio Rank (3Yr)",
      "Tracking Error Rank (3Yr)","Excess Performance (5Yr)","R‑Squared (5Yr)",
      "Peer Return Rank (5Yr)","Sharpe Ratio Rank (5Yr)"
    ]

    # 3) Compute pass/fail statuses for this fund
    statuses = {}
    # Manager Tenure ≥3
    info = next((m["Info"] for m in block["Metrics"] if m["Metric"]=="Manager Tenure"),"")
    yrs  = float(re.search(r"(\d+\.?\d*)",info).group(1)) if re.search(r"(\d+\.?\d*)",info) else 0
    statuses["Manager Tenure"] = (yrs >= 3)
    # Other criteria
    for crit in IPS[1:]:
        raw = next((m["Info"] for m in block["Metrics"] if m["Metric"].startswith(crit.split()[0])),"")
        if "Excess Performance" in crit:
            pct = float(re.search(r"([-+]?\d*\.\d+)%",raw).group(1)) if re.search(r"([-+]?\d*\.\d+)%",raw) else 0
            statuses[crit] = (pct > 0)
        elif "R‑Squared" in crit:
            statuses[crit] = True
        else:
            rk = int(re.search(r"(\d+)",raw).group(1)) if re.search(r"(\d+)",raw) else 999
            statuses[crit] = (rk <= 50)

    # 4) Determine overall IPS Status
    fails = sum(not statuses[c] for c in IPS)
    if   fails <= 4:  overall = "Passed IPS Screen"
    elif fails == 5:  overall = "Informal Watch (IW)"
    else:             overall = "Formal Watch (FW)"

    # 5) Build the DataFrame row
    report_date = st.session_state.get("report_date","")
    row = {
      "Category":    category,
      "Time Period": report_date,
      "Plan Assets": "$"
    }
    for idx, crit in enumerate(IPS, start=1):
        row[str(idx)] = statuses[crit]
    row["IPS Status"] = overall

    df_slide1 = pd.DataFrame([row])

    # 6) Style it
    def color_bool(v): return "background-color: green" if v else "background-color: red"
    def style_status(v):
        if v=="Passed IPS Screen":    return "background-color: green; color: white"
        if "Informal Watch" in v:      return "background-color: orange; color: white"
        if "Formal Watch"   in v:      return "background-color: red;   color: white"
        return ""
    styled = df_slide1.style \
        .applymap(color_bool,   subset=[str(i) for i in range(1,len(IPS)+1)]) \
        .applymap(style_status, subset=["IPS Status"])

    st.dataframe(styled, use_container_width=True)

    # === Slide 3 Table 1 ===
    st.markdown("**Slide 3 Table 1**")
    # grab performance data for the selected fund
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    # build Investment Manager label with ticker
    inv_mgr = f"{selected_fund} ({perf_item.get('Ticker','')})"
    # extract Net Expense Ratio and append '%' if not already present
    net_exp = perf_item.get("Net Expense Ratio", "")
    if net_exp and not str(net_exp).endswith("%"):
        net_exp = f"{net_exp}%"
    # assemble and display
    df_slide3 = pd.DataFrame([{
        "Investment Manager": inv_mgr,
        "Net Expense Ratio":  net_exp
    }])
    st.dataframe(df_slide3, use_container_width=True)

    # === Slide 3 Table 2 ===
    st.markdown("**Slide 3 Table 2**")
    # grab the annualized returns for the selected fund
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name")==selected_fund), {})
    # build Investment Manager label with ticker in parentheses
    inv_mgr    = f"{selected_fund} ({perf_item.get('Ticker','')})"
    # use report_date as the QTD column header
    date_label = st.session_state.get("report_date", "QTD")

    # helper to append '%' if missing
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s=="" else f"{s}%"

    # extract and format each return
    qtd   = append_pct(perf_item.get("QTD",""))
    one   = append_pct(perf_item.get("1Yr",""))
    three = append_pct(perf_item.get("3Yr",""))
    five  = append_pct(perf_item.get("5Yr",""))
    ten   = append_pct(perf_item.get("10Yr",""))

    # assemble the row
    row = {
        "Investment Manager": inv_mgr,
        date_label:           qtd,
        "1 Year":             one,
        "3 Year":             three,
        "5 Year":             five,
        "10 Year":            ten
    }
    df_slide3_2 = pd.DataFrame([row])
    st.dataframe(df_slide3_2, use_container_width=True)

    # === Slide 3 Table 3 ===
    st.markdown("**Slide 3 Table 3**")
    
    # 1) Grab the calendar year returns extracted in Step 8 (fund and benchmark data)
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    
    # Check if data exists
    if not fund_cy or not bench_cy:
        st.error("❌ No calendar year returns data found. Ensure Step 8 has been run correctly.")
        return
    
    # 2) Ensure 'Name' exists in the fund and benchmark records (using 'Name' instead of 'Fund Name')
    # Debugging output to check structure
    # st.write(f"Fund data keys: {fund_cy[0].keys() if fund_cy else 'No data'}")
   #  st.write(f"Benchmark data keys: {bench_cy[0].keys() if bench_cy else 'No data'}")
    
    # 3) Find the selected fund’s record and its benchmark record
    fund_rec = next((r for r in fund_cy if r.get("Name") == selected_fund), None)  # Changed "Fund Name" to "Name"
    if not fund_rec:
        st.error(f"❌ Could not find data for selected fund: {selected_fund}")
        return
    
    # 4) Try to match the benchmark data using Name or Ticker
    benchmark_name = selected_fund  # Assume benchmark matches the fund's name, we can refine this logic if needed
    bench_rec = next((r for r in bench_cy if r.get("Name") == benchmark_name or r.get("Ticker") == fund_rec.get("Ticker")), None)
    
    # If benchmark record is not found
    if not bench_rec:
        st.error(f"❌ Could not find benchmark data for selected fund: {selected_fund}")
        return
    
    # 5) Get the years from the calendar year columns (using the first record)
    year_cols = [col for col in fund_rec.keys() if re.match(r"20\d{2}", col)]
    
    # 6) Prepare the rows for the selected fund and benchmark
    rows = []
    
    # 7) Add the selected fund's data
    row_fund = {"Investment Manager": f"{selected_fund} ({fund_rec.get('Ticker','')})"}
    for year in year_cols:
        row_fund[year] = fund_rec.get(year, "")
    rows.append(row_fund)
    
    # 8) Add the benchmark's data, using the benchmark's name (or fallback)
    row_benchmark = {"Investment Manager": f"{bench_rec.get('Name', 'Benchmark')} ({bench_rec.get('Ticker', '')})"}
    for year in year_cols:
        row_benchmark[year] = bench_rec.get(year, "")
    rows.append(row_benchmark)
    
    # 9) Create a DataFrame for the table
    df_slide3_3 = pd.DataFrame(rows, columns=["Investment Manager"] + year_cols)
    
    # 10) Display the table
    st.dataframe(df_slide3_3, use_container_width=True)


    # === Slide 4 Table 1 ===
    st.markdown("**Slide 4 Table 1**")
    # grab 3‑Yr MPT stats
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    stats3 = next((r for r in mpt3 if r["Fund Name"] == selected_fund), {})
    # grab 5‑Yr MPT stats
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    stats5 = next((r for r in mpt5 if r["Fund Name"] == selected_fund), {})
    # build Investment Manager with ticker
    ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
    inv_mgr = f"{selected_fund} ({ticker})"
    # assemble the row
    row = {
        "Investment Manager":        inv_mgr,
        "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
        "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
        "3 Year Beta":               stats3.get("3 Year Beta", ""),
        "5 Year Beta":               stats5.get("5 Year Beta", ""),
        "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
        "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
        "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
        "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
    }
    df_slide4_1 = pd.DataFrame([row])
    st.dataframe(df_slide4_1, use_container_width=True)

    # === Slide 4 Table 2 ===
    st.markdown("**Slide 4 Table 2**")
    # grab risk‑adjusted returns and peer ranks for the selected fund
    risk_table = st.session_state.get("step13_risk_adjusted_table", [])
    peer_table = st.session_state.get("step14_peer_rank_table", [])
    risk_rec = next((r for r in risk_table if r["Fund Name"] == selected_fund), {})
    peer_rec = next((r for r in peer_table if r["Fund Name"] == selected_fund), {})
    
    # build Investment Manager label with ticker
    ticker = risk_rec.get("Ticker") or peer_rec.get("Ticker", "")
    inv_mgr = f"{selected_fund} ({ticker})"
    
    # helper to combine value and peer rank without calculation
    def frac(metric, period):
        r = risk_rec.get(f"{metric} {period}", "")
        p = peer_rec.get(f"{metric} {period}", "")
        return f"{r} / {p}"
    
    # assemble the row
    row = {
        "Investment Manager": inv_mgr,
        "3 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "3Yr"),
        "5 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "5Yr"),
        "3 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "3Yr"),
        "5 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "5Yr"),
        "3 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "3Yr"),
        "5 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "5Yr"),
    }
    
    df_slide4_2 = pd.DataFrame([row])
    st.dataframe(df_slide4_2, use_container_width=True)

    
    # === Slide 5 Table 1 ===
    st.markdown("**Slide 5 Table 1**")
    # grab the scorecard metrics for the selected fund
    blocks      = st.session_state.get("fund_blocks", [])
    block       = next((b for b in blocks if b["Fund Name"] == selected_fund), {})
    raw_tenure  = next((m["Info"] for m in block.get("Metrics", []) if m["Metric"] == "Manager Tenure"), "")
    # extract just the numeric years and append "years"
    import re
    m = re.search(r"(\d+(\.\d+)?)", raw_tenure)
    tenure = f"{m.group(1)} years" if m else raw_tenure

    # build Investment Manager label with ticker
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    inv_mgr   = f"{selected_fund} ({perf_item.get('Ticker','')})"

    # assemble and display
    df_slide5 = pd.DataFrame([{
        "Investment Manager": inv_mgr,
        "Manager Tenure":     tenure
    }])
    st.dataframe(df_slide5, use_container_width=True)

    # === Slide 5 Table 2 ===
    st.markdown("**Slide 5 Table 2**")
    # grab factsheet details for the selected fund
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), {})
    # grab ticker for label
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p["Fund Scorecard Name"] == selected_fund), {})
    # build Investment Manager label
    inv_mgr    = f"{selected_fund} ({perf_item.get('Ticker','')})"
    # extract Net Assets and Avg. Market Cap
    assets     = fs_rec.get("Net Assets", "")
    avg_cap    = fs_rec.get("Avg. Market Cap", "")
    # assemble and display
    df_slide5_2 = pd.DataFrame([{
        "Investment Manager":             inv_mgr,
        "Assets Under Management":        assets,
        "Average Market Capitalization":  avg_cap
    }])
    st.dataframe(df_slide5_2, use_container_width=True)

# –– Powerpoint ––––––––––––––––––––––––––––––––––––

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def load_template():
    # Load the PowerPoint template from the file path in 'assets/template.pptx'
    template_path = "assets/template.pptx"  # Ensure it's a .pptx file
    prs = Presentation(template_path)  # Load the template
    return prs

def set_cell_border(cell, border_color=RGBColor(0, 0, 0)):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    r, g, b = border_color[0], border_color[1], border_color[2]
    hex_color = "%02X%02X%02X" % (r, g, b)
    for line in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
        ln = OxmlElement(line)
        ln.set("w", "12700")
        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", hex_color)
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)

def format_quarter(raw):
    import re
    raw = str(raw).strip()
    
    # Patterns like "Q1: 3/31/2025" or "Q1, 2025"
    match = re.search(r"Q([1-4])[,:\s-]*(\d{4})?", raw, re.IGNORECASE)
    if match:
        qtr, year = match.groups()
        suffix = {"1": "1st", "2": "2nd", "3": "3rd", "4": "4th"}[qtr]
        return f"{suffix} QTR {year}" if year else f"{suffix} QTR"
    
    return raw  # fallback if it doesn't match

def slide_1_table(selected_fund):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Layout 6 is typically a blank slide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Manually add left-aligned title textbox to match logo padding
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))  # left and top = 0.5"
    title_box.line.fill.background()  # remove outline
    
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Investment Watchlist"  # Change this to "Slide 1 Table"
    run.font.size = Pt(20)
    run.font.name = "Helvetica"
    run.font.color.rgb = RGBColor(33, 43, 88)  # #212b58
    p.alignment = PP_ALIGN.LEFT

    top = Inches(1.1)
    subheading = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.3))
    tf = subheading.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = selected_fund
    run.font.name = "Cambria"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(0, 0, 0)

    # Prepare table with fund data
    matching_rows = df[df["Fund Name"] == selected_fund]
    rows = len(matching_rows)
    cols = 15
    col_widths = [1.2, 1.2, 1.2] + [0.4] * 11 + [1]

    table_top = Inches(1.5)
    table_left = Inches(0.3)
    table_width = Inches(9)
    table_height = Inches(0.25 * (rows + 1))

    table = slide.shapes.add_table(rows + 1, cols, table_left, table_top, table_width, table_height).table

    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    headers = ["Category", "Time Period", "Plan Assets"] + [str(i) for i in range(1, 12)] + ["IPS Status"]

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        set_cell_border(cell)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        text_frame = cell.text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        p = text_frame.paragraphs[0]
        p.font.name = "Cambria"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    for row_idx, (_, r) in enumerate(matching_rows.iterrows(), start=1):
        row_vals = [
            r.get("Category", ""),
            format_quarter(r.get("Time Period", "")),
            r.get("Plan Assets", ""),
        ] + [r.get(str(i), "") for i in range(1, 12)] + [r.get("IPS Status", "")]

        for col_idx, val in enumerate(row_vals):
            cell = table.cell(row_idx, col_idx)
            set_cell_border(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

            text_frame = cell.text_frame
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0

            p = text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = "Cambria"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

            if val == "Pass" and col_idx != 14:
                p.text = "✔"
                p.font.color.rgb = RGBColor(0, 176, 80)
            elif val == "Review" and col_idx != 14:
                p.text = "✖"
                p.font.color.rgb = RGBColor(192, 0, 0)
            elif col_idx == 14:
                p.text = ""
                val_str = str(val).strip().lower()

                if val_str == "formal warning":
                    badge_text = "FW"
                    badge_color = RGBColor(192, 0, 0)
                    font_color = RGBColor(255, 255, 255)
                elif val_str == "informal warning":
                    badge_text = "IW"
                    badge_color = RGBColor(255, 165, 0)
                    font_color = RGBColor(255, 255, 255)
                elif val_str == "passed ips screen":
                    badge_text = "✔"
                    badge_color = RGBColor(0, 176, 80)
                    font_color = RGBColor(255, 255, 255)
                else:
                    continue

                badge_left = table_left + sum(Inches(w) for w in col_widths[:col_idx]) + Inches(0.3)
                badge_top = table_top + Inches(0.25 * row_idx) + Inches(0.06)

                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    left=badge_left,
                    top=badge_top,
                    width=Inches(0.5),
                    height=Inches(0.25),
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = badge_color
                shape.line.color.rgb = RGBColor(255, 255, 255)

                tf = shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = badge_text
                run.font.bold = True
                run.font.size = Pt(12 if badge_text == "✔" else 11)
                run.font.color.rgb = font_color
            else:
                p.text = str(val)

    return prs


#––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
# === Main App ===
def run():
    import re
    st.title("Writeup")
    uploaded = st.file_uploader("Upload MPI PDF", type="pdf")
    if not uploaded:
        return

    # ── Initialize templates exactly once ──–––––––––––––––––––––––––––––––––––––––––––––––
    if "bullet_point_templates" not in st.session_state:
        st.session_state["bullet_point_templates"] = [
            "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], [Year] by [QTD_bps_diff] bps ([QTD_pct_diff])."
        ]
    # ––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
    
    with pdfplumber.open(uploaded) as pdf:
        # Step 1
        with st.expander("Step 1: Details", expanded=False):
            first = pdf.pages[0].extract_text() or ""
            process_page1(first)

        # Step 2
        with st.expander("Step 2: Table of Contents", expanded=False):
            toc_text = "".join((pdf.pages[i].extract_text() or "") for i in range(min(3, len(pdf.pages))))
            process_toc(toc_text)

        # Step 3
        with st.expander("Step 3: Scorecard Metrics", expanded=False):
            sp = st.session_state.get('scorecard_page')
            tot = st.session_state.get('total_options')
            if sp and tot is not None:
                step3_process_scorecard(pdf, sp, tot)
            else:
                st.error("Missing scorecard page or total options")

        # Step 4
        with st.expander("Step 4: IPS Screening", expanded=False):
            step4_ips_screen()

        # Step 5
        with st.expander("Step 5: Fund Performance", expanded=False):
            pp = st.session_state.get('performance_page')
            names = [b['Fund Name'] for b in st.session_state.get('fund_blocks', [])]
            if pp and names:
                step5_process_performance(pdf, pp, names)
            else:
                st.error("Missing performance page or fund blocks")

        # Step 6
        with st.expander("Step 6: Fund Factsheets", expanded=True):
            names = [b['Fund Name'] for b in st.session_state.get('fund_blocks', [])]
            step6_process_factsheets(pdf, names)

        # Step 7
        with st.expander("Step 7: Annualized Returns", expanded=False):
            step7_extract_returns(pdf)

        # ── Data Prep for Bullet Points ───────────────────────────────────────────────
        report_date = st.session_state.get("report_date", "")
        m = re.match(r"(\d)(?:st|nd|rd|th)\s+QTR,\s*(\d{4})", report_date)
        quarter = m.group(1) if m else ""
        year    = m.group(2) if m else ""

        for itm in st.session_state["fund_performance_data"]:
            qtd       = float(itm.get("QTD") or 0)
            bench_qtd = float(itm.get("Bench QTD") or 0)
            itm["Perf Direction"] = "overperformed" if qtd >= bench_qtd else "underperformed"
            itm["Quarter"], itm["Year"] = quarter, year
            itm["QTD_bps_diff"] = str(round((qtd - bench_qtd)*100, 1))
            fund_pct  = f"{qtd:.2f}%"
            bench_pct = f"{bench_qtd:.2f}%"
            itm["QTD_pct_diff"] = f"{(qtd - bench_qtd):.2f}%"
            itm["QTD_vs"] = f"{fund_pct} vs. {bench_pct}"

        # Initialize your template exactly once
        if "bullet_point_templates" not in st.session_state:
            st.session_state["bullet_point_templates"] = [
                "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], "
                "[Year] by [QTD_bps_diff] bps ([QTD_vs])."
            ]

        # ───────────────────────────────────────────────────────────────────────────────
        
        # Step 8: Calendar Year Section
        with st.expander("Step 8: Calendar Year Returns", expanded=False):
            step8_calendar_returns(pdf)

        # Step 9: Match Tickers
        with st.expander("Step 9: Risk Analysis (3Yr)", expanded=False):
            step9_risk_analysis_3yr(pdf)

        # Step 10: Match Tickers
        with st.expander("Step 10: Risk Analysis (5Yr)", expanded=False):
            step10_risk_analysis_5yr(pdf)

        # Step 11: MPT Statistics Summary
        with st.expander("Step 11: MPT Statistics Summary", expanded=False):
            step11_create_summary()
            
        # Step 12: Find Factsheet Sub‑Headings
        with st.expander("Step 12: Fund Facts ", expanded=False):
            step12_process_fund_facts(pdf)

        # Step 13: Risk Adjusted Returns
        with st.expander("Step 13: Risk-Adjusted Returns", expanded=False):
            step13_process_risk_adjusted_returns(pdf)

        # Step 14: Peer Risk-Adjusted Return Rank
        with st.expander("Step 14: Peer Risk-Adjusted Return Rank", expanded=False):
            step14_extract_peer_risk_adjusted_return_rank(pdf)

        # Step 15: View Single Fund Details
        with st.expander("Step 15: Single Fund Details", expanded=False):
            step15_display_selected_fund()
                    
        # –– Bullet Points Section –––––––––––––––––––––––––––––––––––––––––––––––––––––––-
        with st.expander("Bullet Points", expanded=False):
            # Get the selected fund from session state
            selected_fund = st.session_state.get('selected_fund', None)
            
            if selected_fund:
                # Retrieve the performance data for the selected fund
                perf_data = st.session_state.get("fund_performance_data", [])
                item = next(x for x in perf_data if x["Fund Scorecard Name"] == selected_fund)
        
                # First bullet point: Performance vs Benchmark
                filled = "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], [Year] by [QTD_bps_diff] bps ([QTD_vs])."
                for field, val in item.items():
                    filled = filled.replace(f"[{field}]", str(val))
                st.markdown(f"- {filled}")
        
                # Second bullet point: IPS Screening Status
                is_passing_ips = item.get("IPS Status") == "Passed IPS Screen"
                if is_passing_ips:
                    filled = "The Fund passed the IPS Screening."
                else:
                    status = "Informal Watch" if "IW" in item.get("IPS Status", "") else "Formal Watch"
                    
                    # Extract relevant data for the returns and risk-adjusted returns
                    three_year_return = float(item.get("3Yr", 0))  # Convert to float
                    bench_three_year = float(item.get("Bench 3Yr", 0))  # Convert to float
                    five_year_return = float(item.get("5Yr", 0))  # Convert to float
                    bench_five_year = float(item.get("Bench 5Yr", 0))  # Convert to float
                    
                    # Calculate the difference in bps
                    bps_three_year = (three_year_return - bench_three_year) * 100
                    bps_five_year = (five_year_return - bench_five_year) * 100
                
                    # Format as percentages with two decimal places
                    three_year_return_str = f"{three_year_return:.2f}%"
                    five_year_return_str = f"{five_year_return:.2f}%"
                    bench_three_year_str = f"{bench_three_year:.2f}%"
                    bench_five_year_str = f"{bench_five_year:.2f}%"
                
                    # Peer Risk-Adjusted Return Rank (Step 14)
                    peer_ranks = st.session_state.get("step14_peer_rank_table", [])
                
                    # Get the rank of the selected fund for 3Yr Sharpe and 5Yr Sharpe with fallback
                    rank_3yr = next(
                        (r.get("Sharpe Ratio Rank 3Yr", "Rank Not Available") for r in peer_ranks if r.get("Fund Name") == selected_fund),
                        "Rank Not Available"
                    )
                    rank_5yr = next(
                        (r.get("Sharpe Ratio Rank 5Yr", "Rank Not Available") for r in peer_ranks if r.get("Fund Name") == selected_fund),
                        "Rank Not Available"
                    )
                
                    # Determine if the fund is in the top or bottom half of the peer group
                    rank_3yr_position = "top" if rank_3yr != "Rank Not Available" and int(rank_3yr) <= 50 else "bottom"
                    rank_5yr_position = "top" if rank_5yr != "Rank Not Available" and int(rank_5yr) <= 50 else "bottom"
                
                    # Fill the bullet point for non-passing funds
                    filled = (
                        f"The fund is now on {status}. Its three-year return currently trails the benchmark by "
                        f"{bps_three_year} bps ({three_year_return_str} vs. {bench_three_year_str}) "
                        f"and its five-year return trails by {bps_five_year} bps ({five_year_return_str} vs. {bench_five_year_str}). "
                        f"In addition, the fund’s three-year absolute and risk-adjusted returns, as measured by Sharpe and Sortino ratios, "
                        f"now rank in the {rank_3yr_position} half of their peer group for 3Yr Sharpe and {rank_5yr_position} half for 5Yr Sharpe."
                    )
                
                # Display the second bullet point
                st.markdown(f"- {filled}")
        
                # Third bullet point: Action for Formal Watch
                if status == "Formal Watch":
                    action = "Action: Consider replacing this fund."
                    st.markdown(f"- {action}")
            else:
                st.error("❌ No fund selected. Please select a fund from Step 15.")
                
        #––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––––
        # === Export PowerPoint ===
        st.markdown("---")
        st.subheader("Export Selected Fund to PowerPoint")
        
        # Ensure that df_summary is initialized
        if "summary_df" not in st.session_state:
            st.session_state["slide_1_table"] = slide_1_table  # Make sure df_summary exists
        
        # Check if the selected fund and summary data are available before exporting
        if st.button("Export to PowerPoint"):
            if selected_fund and not st.session_state["slide_1_table"].empty:
                # Pass both the data and selected fund to generate the slide
                ppt_stream = slide_1_table(st.session_state["slide_1_table"], selected_fund)
                output = BytesIO()
                ppt_stream.save(output)
        
                # Display download button for the PowerPoint file
                st.download_button(
                    label="Download PowerPoint File",
                    data=output.getvalue(),
                    file_name=f"{selected_fund}_Watchlist.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            else:
                st.warning("Please select a fund and ensure data is loaded.")


if __name__ == "__main__":
    run()


