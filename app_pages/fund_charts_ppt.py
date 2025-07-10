import os
import fitz  # PyMuPDF
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import streamlit as st
import base64

# === Together API Setup ===
TOGETHER_API_KEY = st.secrets["together"]["api_key"]
TOGETHER_MODEL = "togethercomputer/llava-1.5-7b-hf"
API_URL = "https://api.together.xyz/v1/chat/completions"

HEADERS = {
    "Authorization": f"Bearer {TOGETHER_API_KEY}",
    "Content-Type": "application/json"
}

# === Convert a page to image bytes ===
def extract_image_from_page(page):
    pix = page.get_pixmap(dpi=200)
    return pix.tobytes("png")

# === Ask Together AI to extract fund name and detect chart ===
def analyze_page(image_bytes):
    img_b64 = base64.b64encode(image_bytes).decode("utf-8")
    prompt = (
        "You are analyzing a fund report page.\n"
        "1. What is the **name of the fund** shown on this page?\n"
        "2. Does this page contain any **charts, tables, or graphs**?\n"
        "Respond in this format:\n"
        '{"fund_name": "...", "contains_chart_or_table": true/false}'
    )

    data = {
        "model": TOGETHER_MODEL,
        "messages": [
            {"role": "system", "content": "You are a helpful financial assistant."},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_b64}"}}
                ]
            }
        ],
        "temperature": 0.1,
        "max_tokens": 200
    }

    try:
        res = requests.post(API_URL, headers=HEADERS, json=data)
        content = res.json()["choices"][0]["message"]["content"]
        result = eval(content.strip())  # safe only for controlled response format
        return result
    except Exception as e:
        st.error(f"API failed: {e}")
        return {"fund_name": "Unknown Fund", "contains_chart_or_table": False}

# === Analyze PDF and group chart pages by fund ===
def extract_and_classify_fund_pages(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    fund_to_images = {}
    for i in range(35, len(doc)):  # Start from page 36
        page = doc[i]
        image_bytes = extract_image_from_page(page)
        result = analyze_page(image_bytes)

        if result["contains_chart_or_table"]:
            fund = result["fund_name"] or "Unknown Fund"
            fund_to_images.setdefault(fund, []).append(image_bytes)

    return fund_to_images

# === Build PowerPoint ===
def build_powerpoint(fund_to_images):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for fund, images in fund_to_images.items():
        for i, image in enumerate(images):
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(BytesIO(image), Inches(0.5), Inches(0.75), width=Inches(8.5))

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0].add_run()
            p.text = f"{fund} (Image {i+1}/{len(images)})"
            p.font.size = Pt(16)
            p.font.bold = True

    return prs

# === Streamlit App ===
def run():
    st.set_page_config(layout="wide")
    st.title("🧠 AI-Powered Chart Extractor (Together AI)")
    st.markdown("Extracts fund names and identifies charts/tables from an MPI PDF using Together AI vision.")

    uploaded_pdf = st.file_uploader("Upload MPI PDF", type=["pdf"])

    if uploaded_pdf:
        with st.spinner("Analyzing PDF pages using Together AI..."):
            fund_charts = extract_and_classify_fund_pages(uploaded_pdf)
            if not fund_charts:
                st.warning("No charts or tables detected.")
                return

            pptx_io = BytesIO()
            build_powerpoint(fund_charts).save(pptx_io)
            pptx_io.seek(0)

            st.success(f"✅ Extracted {sum(len(v) for v in fund_charts.values())} images from {len(fund_charts)} funds.")
            st.download_button("📥 Download PowerPoint", pptx_io, "fund_charts_ai.pptx")

